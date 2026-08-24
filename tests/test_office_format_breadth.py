"""Synthetic coverage for Office family aliases and bounded legacy conversion."""

from __future__ import annotations

import io
import sys
import time
import zipfile
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

from friday.agent_runtime._office_attachments import _OFFICE_SUFFIXES
from friday.archive_formats import archive_dispatch_kind
from friday.documents import DocumentExtractor, office_document_candidate
from friday.documents._office_convert import convert_legacy_office, libreoffice_available
from friday.ingestion._files import _STRUCTURED_OFFICE_MIME_TYPES, _STRUCTURED_OFFICE_SUFFIXES


@pytest.fixture(autouse=True)
def _clean_operator_office_environment(monkeypatch: pytest.MonkeyPatch) -> None:
    for name in (
        "FRIDAY_LIBREOFFICE_PATH",
        "JERICHO_LIBREOFFICE_PATH",
        "FRIDAY_LIBREOFFICE_LIBRARY_PATH",
        "JERICHO_LIBREOFFICE_LIBRARY_PATH",
        "LD_LIBRARY_PATH",
    ):
        monkeypatch.delenv(name, raising=False)


def _docx(marker: str = "WORD FAMILY MARKER") -> bytes:
    document = Document()
    document.add_paragraph(marker)
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx(marker: str = "SHEET FAMILY MARKER") -> bytes:
    workbook = Workbook()
    workbook.active["A1"] = marker
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _pptx(marker: str = "SLIDE FAMILY MARKER") -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = marker
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _odf(marker: str = "ODF FAMILY MARKER") -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w") as archive:
        archive.writestr("content.xml", f"<office><text>{marker}</text></office>")
    return output.getvalue()


def _staroffice_xml(kind: str, marker: str) -> tuple[bytes, str, str]:
    """Build a genuine, minimal OpenOffice.org 1.0 Writer/Calc package."""

    document_class, body, extension = {
        "writer": ("text", f"<text:p>{marker}</text:p>", "sxw"),
        "calc": (
            "spreadsheet",
            (
                '<table:table table:name="Sheet1"><table:table-row>'
                f"<table:table-cell><text:p>{marker}</text:p></table:table-cell>"
                "</table:table-row></table:table>"
            ),
            "sxc",
        ),
    }[kind]
    mime_type = f"application/vnd.sun.xml.{kind}"
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<office:document-content xmlns:office="http://openoffice.org/2000/office" '
        'xmlns:text="http://openoffice.org/2000/text" '
        'xmlns:table="http://openoffice.org/2000/table" '
        f'office:class="{document_class}" office:version="1.0">'
        f"<office:body>{body}</office:body></office:document-content>"
    )
    manifest = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<manifest:manifest xmlns:manifest="http://openoffice.org/2001/manifest">'
        f'<manifest:file-entry manifest:media-type="{mime_type}" manifest:full-path="/"/>'
        '<manifest:file-entry manifest:media-type="text/xml" '
        'manifest:full-path="content.xml"/></manifest:manifest>'
    )
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w") as archive:
        archive.writestr("mimetype", mime_type, compress_type=zipfile.ZIP_STORED)
        archive.writestr("content.xml", content)
        archive.writestr("META-INF/manifest.xml", manifest)
    return output.getvalue(), mime_type, extension


def _pdf(marker: str = "PDF FAMILY MARKER") -> bytes:
    output = io.BytesIO()
    canvas = Canvas(output)
    canvas.drawString(72, 720, marker)
    canvas.save()
    return output.getvalue()


def _with_main_content_type(payload: bytes, canonical: str, replacement: str) -> bytes:
    """Turn a generated OOXML file into a real alias container."""

    source = io.BytesIO(payload)
    output = io.BytesIO()
    with zipfile.ZipFile(source) as archive, zipfile.ZipFile(output, "w") as rewritten:
        for member in archive.infolist():
            content = archive.read(member)
            if member.filename == "[Content_Types].xml":
                old = canonical.encode("utf-8")
                new = replacement.encode("utf-8")
                assert content.count(old) == 1
                content = content.replace(old, new)
            rewritten.writestr(member, content)
    return output.getvalue()


_WORD_CANONICAL = "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
_WORD_TYPES = {
    "docm": "application/vnd.ms-word.document.macroEnabled.main+xml",
    "dotx": "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml",
    "dotm": "application/vnd.ms-word.template.macroEnabledTemplate.main+xml",
}
_SHEET_CANONICAL = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"
_SHEET_TYPES = {
    "xlsm": "application/vnd.ms-excel.sheet.macroEnabled.main+xml",
    "xltx": "application/vnd.openxmlformats-officedocument.spreadsheetml.template.main+xml",
    "xltm": "application/vnd.ms-excel.template.macroEnabled.main+xml",
}
_PRESENTATION_CANONICAL = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
_PRESENTATION_TYPES = {
    "pptm": "application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml",
    "potx": "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
    "potm": "application/vnd.ms-powerpoint.template.macroEnabled.main+xml",
    "ppsx": "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml",
    "ppsm": "application/vnd.ms-powerpoint.slideshow.macroEnabled.main+xml",
}
_FIRST_EXPANSION_TARGETS = {
    "dot": "docx",
    "wpt": "docx",
    "wpd": "docx",
    "pages": "docx",
    "xlt": "xlsx",
    "et": "xlsx",
    "ett": "xlsx",
    "numbers": "xlsx",
    "pot": "pptx",
    "pps": "pptx",
    "dpt": "pptx",
    "dps": "pptx",
    "key": "pptx",
    "pub": "odg",
    "vdx": "odg",
    "vsd": "odg",
    "vsdm": "odg",
    "vsdx": "odg",
    "vstx": "odg",
}
_SECOND_EXPANSION_TARGETS = {
    "abw": "docx",
    "hwp": "docx",
    "lwp": "docx",
    "psw": "docx",
    "sdw": "docx",
    "stw": "docx",
    "sxw": "docx",
    "wri": "docx",
    "zabw": "docx",
    "123": "xlsx",
    "dif": "xlsx",
    "gnm": "xlsx",
    "gnumeric": "xlsx",
    "mp": "xlsx",
    "stc": "xlsx",
    "sxc": "xlsx",
    "wb1": "xlsx",
    "wb2": "xlsx",
    "wdb": "xlsx",
    "wk1": "xlsx",
    "wk3": "xlsx",
    "wk4": "xlsx",
    "wks": "xlsx",
    "wq1": "xlsx",
    "wq2": "xlsx",
    "xlc": "xlsx",
    "xlk": "xlsx",
    "xlm": "xlsx",
    "xlw": "xlsx",
    "sdd": "pptx",
    "sti": "pptx",
    "sxi": "pptx",
    "cdr": "odg",
    "cmx": "odg",
    "fh": "odg",
    "fh1": "odg",
    "fh2": "odg",
    "fh3": "odg",
    "fh4": "odg",
    "fh5": "odg",
    "fh6": "odg",
    "fh7": "odg",
    "fh8": "odg",
    "fh9": "odg",
    "fh10": "odg",
    "fh11": "odg",
    "p65": "odg",
    "pm": "odg",
    "pm6": "odg",
    "pmd": "odg",
    "qxd": "odg",
    "qxt": "odg",
    "sda": "odg",
    "std": "odg",
    "sxd": "odg",
    "wpg": "odg",
    "zmf": "odg",
}
_REGISTERED_EXPANSION_TARGETS = {
    **_FIRST_EXPANSION_TARGETS,
    **_SECOND_EXPANSION_TARGETS,
}
_ALL_CONVERSION_TARGETS = {
    "doc": "docx",
    "wps": "pdf",
    "xls": "xlsx",
    "xlsb": "xlsx",
    "ppt": "pptx",
    **_REGISTERED_EXPANSION_TARGETS,
}


@pytest.mark.parametrize("extension", tuple(_WORD_TYPES))
def test_word_ooxml_family_aliases_use_the_docx_parser(extension: str) -> None:
    payload = _with_main_content_type(
        _docx(),
        _WORD_CANONICAL,
        _WORD_TYPES[extension],
    )
    result = DocumentExtractor(secret_values=()).extract(
        payload,
        f"synthetic.{extension}",
        "application/zip",
    )

    assert result.success is True, result.error
    assert "WORD FAMILY MARKER" in result.text
    assert result.metadata["format"] == "docx"
    assert result.metadata["main_content_type_normalized"] is True
    assert result.office_structure_index is not None


@pytest.mark.parametrize("extension", tuple(_SHEET_TYPES))
def test_excel_ooxml_family_aliases_use_the_xlsx_parser(extension: str) -> None:
    payload = _with_main_content_type(
        _xlsx(),
        _SHEET_CANONICAL,
        _SHEET_TYPES[extension],
    )
    result = DocumentExtractor(secret_values=()).extract(
        payload,
        f"synthetic.{extension}",
        "application/zip",
    )

    assert result.success is True, result.error
    assert "SHEET FAMILY MARKER" in result.text
    assert result.metadata["format"] == "xlsx"
    assert result.office_structure_index is not None


@pytest.mark.parametrize("extension", tuple(_PRESENTATION_TYPES))
def test_powerpoint_ooxml_family_aliases_use_the_pptx_parser(extension: str) -> None:
    payload = _with_main_content_type(
        _pptx(),
        _PRESENTATION_CANONICAL,
        _PRESENTATION_TYPES[extension],
    )
    result = DocumentExtractor(secret_values=()).extract(
        payload,
        f"synthetic.{extension}",
        "application/zip",
    )

    assert result.success is True, result.error
    assert "SLIDE FAMILY MARKER" in result.text
    assert result.metadata["format"] == "pptx"
    if extension != "pptm":
        assert result.metadata["main_content_type_normalized"] is True


@pytest.mark.parametrize("extension", ["sldx", "sldm"])
def test_standalone_slide_containers_are_honestly_unsupported(extension: str) -> None:
    result = DocumentExtractor(secret_values=()).extract(
        _pptx(),
        f"synthetic.{extension}",
        "application/zip",
    )

    assert result.success is False
    assert result.error == "unsupported_document_format"


def test_known_suffix_wins_over_an_adversarial_conflicting_mime(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    executable = _fake_libreoffice(tmp_path / "soffice")
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", str(executable))
    extractor = DocumentExtractor(secret_values=())

    legacy_sheet = extractor.extract(
        _xlsx("LEGACY SHEET"),
        "conflict.xls",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    direct_sheet = extractor.extract(
        _xlsx("DIRECT SHEET"),
        "conflict.xlsx",
        "application/msword",
    )
    odf = extractor.extract(
        _odf("ODF PRECEDENCE"),
        "conflict.odt",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    assert legacy_sheet.success is True, legacy_sheet.error
    assert legacy_sheet.metadata["format"] == "xls"
    assert "LEGACY SHEET" in legacy_sheet.text
    assert direct_sheet.success is True, direct_sheet.error
    assert direct_sheet.metadata["format"] == "xlsx"
    assert "DIRECT SHEET" in direct_sheet.text
    assert odf.success is True, odf.error
    assert odf.metadata["format"] == "odt"
    assert "ODF PRECEDENCE" in odf.text


def test_ooxml_mime_is_used_only_for_an_unknown_suffix() -> None:
    result = DocumentExtractor(secret_values=()).extract(
        _docx("MIME FALLBACK"),
        "neutral.bin",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    assert result.success is True, result.error
    assert result.metadata["format"] == "docx"
    assert "MIME FALLBACK" in result.text


@pytest.mark.parametrize(
    ("extension", "expected_format"),
    [
        ("ott", "odt"),
        ("odm", "odt"),
        ("oth", "odt"),
        ("ots", "ods"),
        ("otp", "odp"),
        ("odg", "odg"),
        ("otg", "odg"),
    ],
)
def test_opendocument_family_aliases_use_the_bounded_xml_parser(
    extension: str,
    expected_format: str,
) -> None:
    result = DocumentExtractor(secret_values=()).extract(
        _odf(),
        f"synthetic.{extension}",
        "application/zip",
    )

    assert result.success is True, result.error
    assert "ODF FAMILY MARKER" in result.text
    assert result.metadata["format"] == expected_format


def _fake_libreoffice(path: Path, *, expected_library_path: str | None = None) -> Path:
    path.write_text(
        f"""#!{sys.executable}
import os
import pathlib
import sys
arguments = sys.argv[1:]
required = ['--headless', '--nologo', '--nodefault', '--nolockcheck', '--norestore']
if any(item not in arguments for item in required):
    raise SystemExit(2)
if os.environ.get('LD_LIBRARY_PATH') != {expected_library_path!r}:
    raise SystemExit(3)
source = pathlib.Path(arguments[-1])
target = {_ALL_CONVERSION_TARGETS!r}[source.suffix[1:]]
output = pathlib.Path(arguments[arguments.index('--outdir') + 1]) / ('source.' + target)
output.write_bytes(source.read_bytes())
""",
        encoding="utf-8",
    )
    path.chmod(0o755)
    return path


def _fake_libreoffice_action(path: Path, action: str) -> Path:
    path.write_text(
        f"""#!{sys.executable}
import pathlib
import subprocess
import sys
import time
arguments = sys.argv[1:]
source = pathlib.Path(arguments[-1])
target = pathlib.Path(arguments[arguments.index('--outdir') + 1]) / 'source.xlsx'
{action}
""",
        encoding="utf-8",
    )
    path.chmod(0o755)
    return path


@pytest.mark.parametrize(
    ("legacy_extension", "marker", "target_format"),
    [
        ("doc", "WORD FAMILY MARKER", "docx"),
        ("xls", "SHEET FAMILY MARKER", "xlsx"),
        ("xlsb", "SHEET FAMILY MARKER", "xlsx"),
        ("ppt", "SLIDE FAMILY MARKER", "pptx"),
    ],
)
def test_legacy_office_families_use_one_closed_libreoffice_fallback(
    legacy_extension: str,
    marker: str,
    target_format: str,
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    executable = _fake_libreoffice(tmp_path / "soffice")
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", str(executable))
    payload = {
        "doc": _docx,
        "xls": _xlsx,
        "xlsb": _xlsx,
        "ppt": _pptx,
    }[legacy_extension]()

    result = DocumentExtractor(secret_values=()).extract(
        payload,
        f"legacy.{legacy_extension}",
        "application/octet-stream",
    )

    assert result.success is True, result.error
    assert marker in result.text
    assert result.metadata["format"] == legacy_extension
    assert result.metadata["converted_format"] == target_format
    assert result.metadata["parser"] == "libreoffice"


@pytest.mark.parametrize(
    ("source_format", "target_format"),
    tuple(_REGISTERED_EXPANSION_TARGETS.items()),
)
def test_registered_uncommon_office_families_use_only_their_fixed_safe_target(
    source_format: str,
    target_format: str,
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    marker = f"CONFIRMED {source_format.upper()} MARKER"
    payload = {
        "docx": _docx,
        "xlsx": _xlsx,
        "pptx": _pptx,
        "odg": _odf,
    }[target_format](marker)
    executable = _fake_libreoffice(tmp_path / "soffice")
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", str(executable))

    result = DocumentExtractor(secret_values=()).extract(
        payload,
        f"confirmed.{source_format}",
        "application/zip",
    )

    assert result.success is True, result.error
    assert marker in result.text
    assert result.metadata["format"] == source_format
    assert result.metadata["converted_format"] == target_format
    assert result.metadata["parser"] == "libreoffice"
    if target_format in {"docx", "xlsx"}:
        assert result.office_structure_index is not None
    else:
        assert result.office_structure_index is None


@pytest.mark.parametrize(
    ("mime_type", "source_format", "target_format"),
    [
        ("application/vnd.wordperfect", "wpd", "docx"),
        ("application/x-iwork-pages-sffpages", "pages", "docx"),
        ("application/x-iwork-numbers-sffnumbers", "numbers", "xlsx"),
        ("application/x-iwork-keynote-sffkey", "key", "pptx"),
        ("application/x-mspublisher", "pub", "odg"),
        ("application/vnd.visio", "vsd", "odg"),
        ("application/vnd.lotus-wordpro", "lwp", "docx"),
        ("application/vnd.sun.xml.writer", "sxw", "docx"),
        ("application/vnd.sun.xml.writer.template", "stw", "docx"),
        ("application/vnd.sun.xml.writer.web", "stw", "docx"),
        ("application/x-hwp", "hwp", "docx"),
        ("application/x-mswrite", "wri", "docx"),
        ("application/x-pocket-word", "psw", "docx"),
        ("application/vnd.sun.xml.calc", "sxc", "xlsx"),
        ("application/vnd.sun.xml.calc.template", "stc", "xlsx"),
        ("application/vnd.sun.xml.impress", "sxi", "pptx"),
        ("application/vnd.sun.xml.impress.template", "sti", "pptx"),
        ("application/vnd.sun.xml.draw", "sxd", "odg"),
        ("application/vnd.sun.xml.draw.template", "std", "odg"),
        ("application/x-pagemaker", "pmd", "odg"),
        ("image/x-cmx", "cmx", "odg"),
        ("image/x-freehand", "fh", "odg"),
        ("image/x-wpg", "wpg", "odg"),
    ],
)
def test_confirmed_mime_fallback_uses_one_canonical_source_family(
    mime_type: str,
    source_format: str,
    target_format: str,
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    marker = f"MIME {source_format.upper()} MARKER"
    payload = {
        "docx": _docx,
        "xlsx": _xlsx,
        "pptx": _pptx,
        "odg": _odf,
    }[target_format](marker)
    executable = _fake_libreoffice(tmp_path / "soffice")
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", str(executable))

    result = DocumentExtractor(secret_values=()).extract(
        payload,
        "neutral.bin",
        mime_type,
    )

    assert result.success is True, result.error
    assert marker in result.text
    assert result.metadata["format"] == source_format
    assert result.metadata["converted_format"] == target_format


def test_converted_word_and_sheet_families_are_registered_as_structured() -> None:
    structured_suffixes = {
        f".{source_format}"
        for source_format, target_format in _REGISTERED_EXPANSION_TARGETS.items()
        if target_format in {"docx", "xlsx"}
    }

    assert structured_suffixes <= _STRUCTURED_OFFICE_SUFFIXES
    assert structured_suffixes <= _OFFICE_SUFFIXES
    assert {
        "application/vnd.lotus-1-2-3",
        "application/vnd.lotus-wordpro",
        "application/vnd.sun.xml.calc",
        "application/vnd.sun.xml.calc.template",
        "application/vnd.sun.xml.writer",
        "application/vnd.sun.xml.writer.template",
        "application/vnd.sun.xml.writer.web",
        "application/vnd.wordperfect",
        "application/x-abiword",
        "application/x-gnumeric",
        "application/x-hwp",
        "application/x-iwork-pages-sffpages",
        "application/x-iwork-numbers-sffnumbers",
        "application/x-mswrite",
        "application/x-pocket-word",
    } <= _STRUCTURED_OFFICE_MIME_TYPES


@pytest.mark.parametrize("source_format", tuple(_REGISTERED_EXPANSION_TARGETS))
def test_registered_office_container_is_never_misrouted_as_an_archive(source_format: str) -> None:
    assert office_document_candidate(f"document.{source_format}", "application/octet-stream")
    assert archive_dispatch_kind(f"document.{source_format}", "application/zip") is None


@pytest.mark.parametrize(
    "mime_type",
    [
        "application/vnd.lotus-wordpro",
        "application/vnd.sun.xml.writer",
        "application/vnd.sun.xml.writer.template",
        "application/vnd.sun.xml.writer.web",
        "application/x-hwp",
        "application/x-mswrite",
        "application/x-pocket-word",
        "application/vnd.sun.xml.calc",
        "application/vnd.sun.xml.calc.template",
        "application/vnd.sun.xml.impress",
        "application/vnd.sun.xml.impress.template",
        "application/vnd.sun.xml.draw",
        "application/vnd.sun.xml.draw.template",
        "application/x-pagemaker",
        "image/x-cmx",
        "image/x-freehand",
        "image/x-wpg",
    ],
)
def test_registry_mime_only_office_carriers_enter_the_closed_document_matrix(
    mime_type: str,
) -> None:
    assert office_document_candidate("neutral.bin", mime_type) is True
    assert office_document_candidate("declared.txt", mime_type) is False
    assert office_document_candidate("declared.zip", mime_type) is False
    assert archive_dispatch_kind("neutral.bin", mime_type) is None


@pytest.mark.parametrize("mime_type", ["application/rtf", "text/rtf"])
def test_suffixless_rtf_mime_uses_the_native_rtf_parser(mime_type: str) -> None:
    result = DocumentExtractor(secret_values=()).extract(
        b"{\\rtf1\\ansi MIME RTF MARKER}",
        "neutral.bin",
        mime_type,
    )

    assert result.success is True, result.error
    assert "MIME RTF MARKER" in result.text
    assert result.metadata["format"] == "rtf"
    assert office_document_candidate("neutral.bin", mime_type) is True


@pytest.mark.parametrize("kind", ["writer", "calc"])
def test_real_libreoffice_imports_genuine_staroffice_xml_through_mime_only_dispatch(
    kind: str,
) -> None:
    if not libreoffice_available():
        pytest.skip("real LibreOffice backend is not installed")
    marker = f"GENUINE {kind.upper()} LEGACY CANARY 2026"
    payload, mime_type, extension = _staroffice_xml(kind, marker)

    result = DocumentExtractor(secret_values=(), parse_budget_sec=30).extract(
        payload,
        "neutral.bin",
        mime_type,
    )

    assert result.success is True, result.error
    assert marker in result.text
    assert result.metadata["format"] == extension
    assert result.metadata["parser"] == "libreoffice"


def test_converted_odg_output_uses_the_bounded_xml_reader(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    marker = "BOUNDED ODG MARKER"
    payload = _odf(marker + " " + "0123456789abcdef" * 140_000)
    executable = _fake_libreoffice(tmp_path / "soffice")
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", str(executable))

    result = DocumentExtractor(secret_values=()).extract(
        payload,
        "drawing.pub",
        "application/x-mspublisher",
    )

    assert result.success is True, result.error
    assert result.text.startswith(marker)
    assert result.metadata["format"] == "pub"
    assert result.metadata["converted_format"] == "odg"
    assert result.metadata["source_truncated_for_parse"] is True


def test_libreoffice_rootless_loader_path_is_explicit_and_ambient_state_is_stripped(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    library_dir = tmp_path / "private-runtime" / "lib"
    library_dir.mkdir(parents=True)
    executable = _fake_libreoffice(
        tmp_path / "soffice",
        expected_library_path=str(library_dir.resolve()),
    )
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", str(executable))
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_LIBRARY_PATH", str(library_dir))
    monkeypatch.setenv("LD_LIBRARY_PATH", "/ambient/value/must/not/leak")

    converted = convert_legacy_office(_xlsx(), "xls")

    assert libreoffice_available() is True
    assert converted.success is True
    assert converted.target_format == "xlsx"


def test_relative_libreoffice_executable_fails_closed(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", "relative/soffice")
    monkeypatch.setenv("PATH", "")

    converted = convert_legacy_office(b"synthetic", "xls")

    assert libreoffice_available() is False
    assert converted.success is False
    assert converted.error == "libreoffice_unavailable"


def test_ambiguous_works_family_uses_the_common_bounded_pdf_target(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    marker = "WORKS FAMILY MARKER"
    executable = _fake_libreoffice(tmp_path / "soffice")
    monkeypatch.setenv("FRIDAY_LIBREOFFICE_PATH", str(executable))
    extracted = DocumentExtractor(secret_values=()).extract(
        _pdf(marker),
        "ambiguous.wps",
        "application/vnd.ms-works",
    )

    assert extracted.success is True, extracted.error
    assert marker in extracted.text
    assert extracted.metadata["format"] == "wps"
    assert extracted.metadata["converted_format"] == "pdf"
    assert extracted.office_structure_index is None
    assert ".wps" not in _STRUCTURED_OFFICE_SUFFIXES
    assert ".wps" not in _OFFICE_SUFFIXES


def test_libreoffice_timeout_kills_the_whole_converter_process_group(tmp_path: Path) -> None:
    marker = tmp_path / "escaped-child"
    child = f"import pathlib,time;time.sleep(0.4);pathlib.Path({str(marker)!r}).write_text('escaped')"
    executable = _fake_libreoffice_action(
        tmp_path / "soffice",
        f"subprocess.Popen([sys.executable, '-c', {child!r}])\ntime.sleep(5)",
    )

    converted = convert_legacy_office(
        b"synthetic",
        "xls",
        executable=str(executable),
        deadline=time.monotonic() + 0.1,
    )
    time.sleep(0.6)

    assert converted.success is False
    assert converted.error == "libreoffice_deadline_reached"
    assert marker.exists() is False


def test_libreoffice_nonzero_exit_is_not_treated_as_conversion(tmp_path: Path) -> None:
    executable = _fake_libreoffice_action(tmp_path / "soffice", "raise SystemExit(7)")

    converted = convert_legacy_office(b"synthetic", "xls", executable=str(executable))

    assert converted.success is False
    assert converted.error == "libreoffice_conversion_failed"


def test_libreoffice_missing_output_is_reported_exactly(tmp_path: Path) -> None:
    executable = _fake_libreoffice_action(tmp_path / "soffice", "raise SystemExit(0)")

    converted = convert_legacy_office(b"synthetic", "xls", executable=str(executable))

    assert converted.success is False
    assert converted.error == "libreoffice_output_missing"


def test_libreoffice_oversize_output_is_rejected_before_read(tmp_path: Path) -> None:
    executable = _fake_libreoffice_action(
        tmp_path / "soffice",
        "target.write_bytes(b'x' * 257)",
    )

    converted = convert_legacy_office(
        b"synthetic",
        "xls",
        executable=str(executable),
        max_output_bytes=256,
    )

    assert converted.success is False
    assert converted.error == "libreoffice_output_too_large"


def test_libreoffice_symlink_output_is_rejected_without_following_it(tmp_path: Path) -> None:
    executable = _fake_libreoffice_action(
        tmp_path / "soffice",
        "target.symlink_to(source)",
    )

    converted = convert_legacy_office(b"synthetic", "xls", executable=str(executable))

    assert converted.success is False
    assert converted.error == "libreoffice_output_invalid"


def test_libreoffice_multiply_linked_output_is_rejected(tmp_path: Path) -> None:
    executable = _fake_libreoffice_action(
        tmp_path / "soffice",
        "target.hardlink_to(source)",
    )

    converted = convert_legacy_office(b"synthetic", "xls", executable=str(executable))

    assert converted.success is False
    assert converted.error == "libreoffice_output_invalid"
