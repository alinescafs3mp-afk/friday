"""Safe, in-memory document text extraction.

The extractor never writes archive members to disk.  Archive entry counts,
uncompressed byte totals, individual reads, and recursion depth are bounded so
uploaded files cannot expand without limit.  Optional Office/RTF/RAR support is
activated when the corresponding dependency is installed.
"""

from __future__ import annotations

import bz2
import csv
import gzip
import html
import io
import itertools
import json
import logging
import lzma
import math
import mimetypes
import os
import re
import selectors
import shutil
import stat
import subprocess
import tarfile
import tempfile
import time
import zipfile
import zlib
from collections.abc import Sequence
from contextlib import ExitStack, closing, suppress
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path, PurePosixPath
from typing import Any, Protocol

from bs4 import BeautifulSoup

from friday.archive_formats import archive_dispatch_kind
from friday.archive_passwords import archive_password_candidates
from friday.documents._ocr import LocalOcrResult, extract_local_ocr, local_ocr_available
from friday.documents._office_structure import (
    build_docx_text_and_structure,
    build_xlsx_text_and_structure,
    validate_office_structure_index,
)
from friday.work_budgets import size_scaled_budget_sec, stage_deadline

LOGGER = logging.getLogger(__name__)
#: Чем заменяется собственный credential, найденный в тексте документа.
#: Замена ВИДНА: молчаливая подмена читалась бы как свойство документа.
SECRET_PLACEHOLDER = "[секрет удалён]"
# API tokens minted by Friday are ``jrc_`` plus the unpadded URL-safe base64
# returned by ``secrets.token_urlsafe(32)`` (43 characters today).  Match a
# small lower-bound range rather than one exact length so a harmless change in
# the entropy size cannot reopen the disclosure.  Short words such as
# ``jrc_example`` remain ordinary document text.
_FRIDAY_API_TOKEN = re.compile(r"(?<![A-Za-z0-9_-])jrc_[A-Za-z0-9_-]{40,}")


def _own_secret_values() -> tuple[str, ...]:
    """Credential этого экземпляра — те, что нельзя показывать модели.

    `secret_hygiene` знает их по имени переменной окружения и держит порог длины,
    ниже которого совпадение было бы случайным. Он же — offline-доктор, который
    ищет их в файлах на диске; здесь тот же список используется на входе, где
    файл ещё только становится текстом.
    """
    from friday.secret_hygiene import named_secrets

    return tuple(dict.fromkeys(named_secrets().values()))


def _redact_own_secrets(text: str, secrets: Sequence[str]) -> tuple[str, int]:
    """Убрать из текста credential, сосчитав каждую замену.

    Exact values cover this instance's environment.  The structural Friday
    token shape also covers a token uploaded by another user or restored from
    an archive: the extractor cannot know that value in advance, but it is no
    safer to index or quote it.
    """
    if not text:
        return text, 0
    removed = 0
    for secret in secrets:
        if not secret or secret not in text:
            continue
        removed += text.count(secret)
        text = text.replace(secret, SECRET_PLACEHOLDER)
    text, shaped_removed = _FRIDAY_API_TOKEN.subn(SECRET_PLACEHOLDER, text)
    removed += shaped_removed
    return text, removed


_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".xml",
    ".log",
    ".py",
    ".ps1",
    ".sh",
    ".ini",
    ".cfg",
    ".conf",
    ".toml",
    ".sql",
    ".js",
    ".ts",
    ".css",
    ".rst",
}
_HTML_EXTENSIONS = {".html", ".htm"}
#: Семья OpenDocument: текст лежит в `content.xml` и у документов, и у их
#: шаблонов.  Расширение выбирает формат результата, но не отдельный parser.
_OPENDOCUMENT_EXTENSIONS = {
    ".odt": "odt",
    ".ott": "odt",
    ".odm": "odt",
    ".oth": "odt",
    ".ods": "ods",
    ".ots": "ods",
    ".odp": "odp",
    ".otp": "odp",
    ".odg": "odg",
    ".otg": "odg",
}
_OPENDOCUMENT_MIME_TYPES = {
    "application/vnd.oasis.opendocument.text",
    "application/vnd.oasis.opendocument.text-template",
    "application/vnd.oasis.opendocument.text-master",
    "application/vnd.oasis.opendocument.text-web",
    "application/vnd.oasis.opendocument.spreadsheet",
    "application/vnd.oasis.opendocument.spreadsheet-template",
    "application/vnd.oasis.opendocument.presentation",
    "application/vnd.oasis.opendocument.presentation-template",
    "application/vnd.oasis.opendocument.graphics",
    "application/vnd.oasis.opendocument.graphics-template",
}
_WORDPROCESSING_EXTENSIONS = frozenset({".docx", ".docm", ".dotx", ".dotm"})
_SPREADSHEET_EXTENSIONS = frozenset({".xlsx", ".xlsm", ".xltx", ".xltm"})
_PRESENTATION_EXTENSIONS = frozenset({".pptx", ".pptm", ".potx", ".potm", ".ppsx", ".ppsm"})
_OOXML_EXTENSIONS = {
    **{extension: "docx" for extension in _WORDPROCESSING_EXTENSIONS},
    **{extension: "xlsx" for extension in _SPREADSHEET_EXTENSIONS},
    **{extension: "pptx" for extension in _PRESENTATION_EXTENSIONS},
}
_OOXML_MIME_FORMATS = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.template": "docx",
    "application/vnd.ms-word.document.macroenabled.12": "docx",
    "application/vnd.ms-word.template.macroenabled.12": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.template": "xlsx",
    "application/vnd.ms-excel.sheet.macroenabled.12": "xlsx",
    "application/vnd.ms-excel.template.macroenabled.12": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.presentationml.template": "pptx",
    "application/vnd.openxmlformats-officedocument.presentationml.slideshow": "pptx",
    "application/vnd.ms-powerpoint.presentation.macroenabled.12": "pptx",
    "application/vnd.ms-powerpoint.template.macroenabled.12": "pptx",
    "application/vnd.ms-powerpoint.slideshow.macroenabled.12": "pptx",
}
_OOXML_MIME_TYPES = frozenset(_OOXML_MIME_FORMATS)
_EMAIL_METADATA_EXTENSIONS = {".eml": "eml", ".mht": "mhtml", ".mhtml": "mhtml"}
_MSG_MIME_TYPES = frozenset({"application/vnd.ms-outlook", "application/x-msg"})
_CONVERTED_OFFICE_FORMATS = {
    ".abw": "abw",
    ".doc": "doc",
    ".dot": "dot",
    ".hwp": "hwp",
    ".lwp": "lwp",
    ".psw": "psw",
    ".sdw": "sdw",
    ".stw": "stw",
    ".sxw": "sxw",
    ".wri": "wri",
    ".wps": "wps",
    ".wpt": "wpt",
    ".wpd": "wpd",
    ".zabw": "zabw",
    ".pages": "pages",
    ".123": "123",
    ".dif": "dif",
    ".xls": "xls",
    ".xlsb": "xlsb",
    ".xlt": "xlt",
    ".xlc": "xlc",
    ".xlk": "xlk",
    ".xlm": "xlm",
    ".xlw": "xlw",
    ".et": "et",
    ".ett": "ett",
    ".gnm": "gnm",
    ".gnumeric": "gnumeric",
    ".mp": "mp",
    ".numbers": "numbers",
    ".stc": "stc",
    ".sxc": "sxc",
    ".wb1": "wb1",
    ".wb2": "wb2",
    ".wdb": "wdb",
    ".wk1": "wk1",
    ".wk3": "wk3",
    ".wk4": "wk4",
    ".wks": "wks",
    ".wq1": "wq1",
    ".wq2": "wq2",
    ".ppt": "ppt",
    ".pot": "pot",
    ".pps": "pps",
    ".dpt": "dpt",
    ".dps": "dps",
    ".key": "key",
    ".sdd": "sdd",
    ".sti": "sti",
    ".sxi": "sxi",
    ".pub": "pub",
    ".cdr": "cdr",
    ".cmx": "cmx",
    ".fh": "fh",
    ".fh1": "fh1",
    ".fh2": "fh2",
    ".fh3": "fh3",
    ".fh4": "fh4",
    ".fh5": "fh5",
    ".fh6": "fh6",
    ".fh7": "fh7",
    ".fh8": "fh8",
    ".fh9": "fh9",
    ".fh10": "fh10",
    ".fh11": "fh11",
    ".p65": "p65",
    ".pm": "pm",
    ".pm6": "pm6",
    ".pmd": "pmd",
    ".qxd": "qxd",
    ".qxt": "qxt",
    ".sda": "sda",
    ".std": "std",
    ".sxd": "sxd",
    ".vdx": "vdx",
    ".vsd": "vsd",
    ".vsdm": "vsdm",
    ".vsdx": "vsdx",
    ".vstx": "vstx",
    ".wpg": "wpg",
    ".zmf": "zmf",
}
_CONVERTED_OFFICE_MIME_FORMATS = {
    "application/msword": "doc",
    "application/vnd.ms-works": "wps",
    "application/vnd.wordperfect": "wpd",
    "application/vnd.lotus-wordpro": "lwp",
    "application/vnd.sun.xml.writer": "sxw",
    "application/vnd.sun.xml.writer.template": "stw",
    "application/vnd.sun.xml.writer.web": "stw",
    "application/x-abiword": "abw",
    "application/x-hwp": "hwp",
    "application/x-iwork-pages-sffpages": "pages",
    "application/x-mswrite": "wri",
    "application/x-pocket-word": "psw",
    "application/vnd.ms-excel": "xls",
    "application/x-msexcel": "xls",
    "application/vnd.ms-excel.sheet.binary.macroenabled.12": "xlsb",
    "application/vnd.sun.xml.calc": "sxc",
    "application/vnd.sun.xml.calc.template": "stc",
    "application/x-iwork-numbers-sffnumbers": "numbers",
    "application/x-gnumeric": "gnumeric",
    "application/vnd.lotus-1-2-3": "123",
    "application/vnd.ms-powerpoint": "ppt",
    "application/mspowerpoint": "ppt",
    "application/x-mspowerpoint": "ppt",
    "application/vnd.sun.xml.impress": "sxi",
    "application/vnd.sun.xml.impress.template": "sti",
    "application/x-iwork-keynote-sffkey": "key",
    "application/x-mspublisher": "pub",
    "application/vnd.corel-draw": "cdr",
    "application/vnd.sun.xml.draw": "sxd",
    "application/vnd.sun.xml.draw.template": "std",
    "application/vnd.visio": "vsd",
    "application/x-pagemaker": "pmd",
    "image/x-cmx": "cmx",
    "image/x-freehand": "fh",
    "image/x-wpg": "wpg",
}
_ODF_META_MEMBER = "meta.xml"
_MAX_ODF_METADATA_BYTES = 256 * 1024
_MAX_ODF_SIGNATURE_BYTES = 512 * 1024
_MAX_ODF_SIGNATURE_MEMBERS = 8
_MAX_TECHNICAL_METADATA_XML_BYTES = 512 * 1024
_MAX_TECHNICAL_METADATA_RECORDS = 64
_MAX_PDF_SIGNATURE_FIELDS = 16
_MAX_PDF_FORM_FIELDS_SCANNED = 256
_MAX_EMAIL_METADATA_BYTES = 2 * 1024 * 1024
_MAX_EPUB_CONTAINER_BYTES = 64 * 1024
_MAX_EPUB_PACKAGE_BYTES = 512 * 1024
_ODF_OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
_ODF_META_NS = "urn:oasis:names:tc:opendocument:xmlns:meta:1.0"
_ODF_DC_NS = "http://purl.org/dc/elements/1.1/"
_ODF_XLINK_NS = "http://www.w3.org/1999/xlink"
_XMLDSIG_NS = "http://www.w3.org/2000/09/xmldsig#"
_OOXML_CP_NS = "http://schemas.openxmlformats.org/package/2006/metadata/core-properties"
_OOXML_DC_NS = "http://purl.org/dc/elements/1.1/"
_OOXML_DCTERMS_NS = "http://purl.org/dc/terms/"
_OOXML_EP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
_OOXML_VT_NS = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
_OOXML_CUSTOM_NS = "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
_ODF_STRING_TAGS = {
    f"{{{_ODF_DC_NS}}}title": ("title", 500),
    f"{{{_ODF_DC_NS}}}subject": ("subject", 500),
    f"{{{_ODF_DC_NS}}}creator": ("creator", 300),
    f"{{{_ODF_META_NS}}}initial-creator": ("initial_creator", 300),
    f"{{{_ODF_DC_NS}}}description": ("description", 4000),
    f"{{{_ODF_DC_NS}}}language": ("language", 64),
    f"{{{_ODF_META_NS}}}generator": ("generator", 300),
    f"{{{_ODF_META_NS}}}printed-by": ("printed_by", 300),
}
_ODF_STATISTIC_ATTRIBUTES = {
    f"{{{_ODF_META_NS}}}cell-count": "cell_count",
    f"{{{_ODF_META_NS}}}draw-count": "draw_count",
    f"{{{_ODF_META_NS}}}frame-count": "frame_count",
    f"{{{_ODF_META_NS}}}page-count": "page_count",
    f"{{{_ODF_META_NS}}}paragraph-count": "paragraph_count",
    f"{{{_ODF_META_NS}}}row-count": "row_count",
    f"{{{_ODF_META_NS}}}sentence-count": "sentence_count",
    f"{{{_ODF_META_NS}}}syllable-count": "syllable_count",
    f"{{{_ODF_META_NS}}}word-count": "word_count",
    f"{{{_ODF_META_NS}}}character-count": "character_count",
    f"{{{_ODF_META_NS}}}non-whitespace-character-count": "non_whitespace_character_count",
    f"{{{_ODF_META_NS}}}table-count": "table_count",
    f"{{{_ODF_META_NS}}}image-count": "image_count",
    f"{{{_ODF_META_NS}}}object-count": "object_count",
    f"{{{_ODF_META_NS}}}ole-object-count": "ole_object_count",
}
_OFFICE_EXTENSIONS = {
    *_OOXML_EXTENSIONS,
    *_OPENDOCUMENT_EXTENSIONS,
    *_CONVERTED_OFFICE_FORMATS,
    ".rtf",
}
_OFFICE_MIME_TYPES = frozenset(
    {
        *_OOXML_MIME_TYPES,
        *_OPENDOCUMENT_MIME_TYPES,
        *_CONVERTED_OFFICE_MIME_FORMATS,
        "application/rtf",
        "text/rtf",
    }
)


def office_document_candidate(filename: str, mime_type: str = "") -> bool:
    """Whether Friday's suffix-first dispatch selects the closed Office matrix."""

    suffix = Path(str(filename or "")).suffix.casefold()
    normalized_mime = str(mime_type or "").split(";", 1)[0].strip().casefold()
    if suffix in _OFFICE_EXTENSIONS:
        return True
    if suffix in _KNOWN_DOCUMENT_EXTENSIONS:
        return False
    if archive_dispatch_kind(filename, normalized_mime) is not None:
        return False
    return normalized_mime in _OFFICE_MIME_TYPES


_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
_KNOWN_DOCUMENT_EXTENSIONS = frozenset(
    {
        *_TEXT_EXTENSIONS,
        *_HTML_EXTENSIONS,
        *_OOXML_EXTENSIONS,
        *_OPENDOCUMENT_EXTENSIONS,
        *_CONVERTED_OFFICE_FORMATS,
        *_EMAIL_METADATA_EXTENSIONS,
        *_IMAGE_EXTENSIONS,
        ".epub",
        ".msg",
        ".pdf",
        ".rtf",
        ".sldm",
        ".sldx",
    }
)
_CONTENT_TYPES_MEMBER = "[Content_Types].xml"
_MAX_OOXML_CONTENT_TYPES_BYTES = 512 * 1024
_WORD_MAIN_PART = "/word/document.xml"
_WORD_CANONICAL_MAIN_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
_WORD_ALIAS_MAIN_TYPES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml",
        "application/vnd.ms-word.document.macroEnabled.main+xml",
        "application/vnd.ms-word.template.macroEnabledTemplate.main+xml",
    }
)
_PRESENTATION_MAIN_PART = "/ppt/presentation.xml"
_PRESENTATION_CANONICAL_MAIN_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)
_PRESENTATION_ALIAS_MAIN_TYPES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml",
        "application/vnd.ms-powerpoint.template.macroEnabled.main+xml",
        "application/vnd.ms-powerpoint.slideshow.macroEnabled.main+xml",
    }
)
_MAX_NESTING_DEPTH = 2
_MAX_ARCHIVE_PREVIEW_FILES = 24
_MAX_STRUCTURED_PARSE_BYTES = 2 * 1024 * 1024
_MAX_TEXT_PARSE_BYTES = 8 * 1024 * 1024
_MAX_OFFICE_EXPANDED_BYTES = 128 * 1024 * 1024
_MAX_OFFICE_MEMBER_BYTES = 64 * 1024 * 1024
_MAX_ZIP_RATIO = 500.0
_MAX_TABULAR_ROWS = 100_000
_MAX_PDF_RENDER_AXIS_PIXELS = 16_384
_MAX_ARCHIVE_DICTIONARY_BYTES = 128 * 1024 * 1024
_MAX_7Z_HEADER_BYTES = 8 * 1024 * 1024
_MAX_7Z_AES_CYCLES_POWER = 20
_RAR_DICTIONARY_SWITCH = f"-mdx{_MAX_ARCHIVE_DICTIONARY_BYTES // (1024 * 1024)}m"
_RAR_MEMBER_TIMEOUT_SEC = 20.0
_ARCHIVE_PARSE_BASE_TIMEOUT_SEC = 60.0
_ARCHIVE_PARSE_TIMEOUT_PER_MIB_SEC = 2.0
_ARCHIVE_PARSE_MAX_TIMEOUT_SEC = 300.0

# Дата САМОГО документа, взятая из провенанса файла, а не угаданная из текста.
#
# Существующий фильтр по датам намеренно означает «документ УПОМИНАЕТ дату»:
# документ называет несколько дат, и какая из них его собственная, текст не
# говорит. Но у файла есть собственная дата, которую никто не выдумывал, — её
# записал редактор при сохранении (docProps/core.xml у docx/xlsx, /CreationDate
# у pdf). На архиве владельца это единственный способ отличить документ 2019
# года от документа 2025-го: у всех 1537 объектов дата загрузки одна и та же —
# день импорта.
#
# Берётся ТОЛЬКО из формата. Файловое mtime сюда не годится: копирование на
# флешку переписывает его у всех файлов разом, и «дата документа» стала бы
# датой копирования — то самое угадывание, от которого фильтр и уходил.
_OFFICE_CORE_PROPERTIES = "docProps/core.xml"
_CORE_DATE_RE = re.compile(
    r"<(?:dcterms:)?(created|modified)[^>]*>([0-9]{4}-[0-9]{2}-[0-9]{2})", re.IGNORECASE
)
# 1900 — ниже этого у офисных файлов лежат только служебные нули и мусор
# конвертеров; 2100 — потолок против «31.12.9999», который ставят генераторы.
_DOCUMENT_DATE_MIN = "1900-01-01"
_DOCUMENT_DATE_MAX = "2100-01-01"


def _plausible_document_date(value: str) -> str | None:
    """ISO-дата, если она вообще похожа на настоящую дату документа."""
    candidate = (value or "").strip()[:10]
    if len(candidate) != 10 or candidate[4] != "-" or candidate[7] != "-":
        return None
    if not (_DOCUMENT_DATE_MIN <= candidate < _DOCUMENT_DATE_MAX):
        return None
    try:
        datetime.strptime(candidate, "%Y-%m-%d")  # noqa: DTZ007 - календарная дата без времени
    except ValueError:
        return None
    return candidate


def _office_document_date(content: bytes) -> str | None:
    """`dcterms:created` из docProps/core.xml; при отсутствии — `modified`."""
    with suppress(Exception):
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            if _OFFICE_CORE_PROPERTIES not in archive.namelist():
                return None
            with archive.open(_OFFICE_CORE_PROPERTIES) as handle:
                # Ограничение чтения: core.xml — это килобайты, и раздутый член
                # архива не должен превращаться в чтение сотен мегабайт.
                raw = handle.read(64 * 1024).decode("utf-8", errors="replace")
        found: dict[str, str] = {}
        for match in _CORE_DATE_RE.finditer(raw):
            found.setdefault(match.group(1).casefold(), match.group(2))
        for key in ("created", "modified"):
            candidate = _plausible_document_date(found.get(key, ""))
            if candidate:
                return candidate
    return None


# Форматы, у которых внутри zip лежит docProps/core.xml.
_OFFICE_DATE_EXTENSIONS = set(_OOXML_EXTENSIONS)


def _email_iso_date(raw: str) -> str:
    """Дата письма в виде ISO — из заголовка `Date`, который пишет почтовик."""
    if not raw.strip():
        return ""
    with suppress(Exception):
        from email.utils import parsedate_to_datetime

        return parsedate_to_datetime(raw).date().isoformat()
    return ""


def _pdf_document_date_from_bytes(content: bytes) -> str | None:
    """Дата PDF по байтам — работает и для скана без текстового слоя."""
    with suppress(Exception):
        from pypdf import PdfReader

        return _pdf_document_date(PdfReader(io.BytesIO(content), strict=False))
    return None


def _pdf_document_date(reader: Any) -> str | None:
    """`/CreationDate` вида `D:20230412...` из метаданных PDF."""
    with suppress(Exception):
        info = reader.metadata or {}
        for key in ("/CreationDate", "/ModDate"):
            raw = str(info.get(key) or "")
            digits = raw[2:10] if raw.startswith("D:") else raw[:8]
            if len(digits) == 8 and digits.isdigit():
                candidate = _plausible_document_date(f"{digits[:4]}-{digits[4:6]}-{digits[6:8]}")
                if candidate:
                    return candidate
    return None


@dataclass(frozen=True)
class DocumentResult:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    success: bool = True
    error: str = ""
    # Last on purpose: older parser doubles and callers construct this class
    # positionally as ``text, metadata, success, error``.  The structure is an
    # optional, content-free companion to the immutable legacy text contract.
    office_structure_index: dict[str, Any] | None = None

    def to_dict(self, *, preview_chars: int = 2_000) -> dict[str, Any]:
        return {
            "text": self.text[:preview_chars],
            "text_length": len(self.text),
            "metadata": self.metadata,
            "success": self.success,
            "error": self.error,
            "truncated": len(self.text) > preview_chars,
        }


@dataclass(frozen=True)
class VisualAsset:
    """Bounded, normalized image supplied to the local vision model."""

    data: bytes
    mime_type: str
    source: str
    width: int
    height: int

    def to_dict(self) -> dict[str, Any]:
        return {
            "mime_type": self.mime_type,
            "source": self.source,
            "width": self.width,
            "height": self.height,
            "size_bytes": len(self.data),
        }


@dataclass(frozen=True)
class VisualPageRender:
    """Bounded, contiguous PDF-page render prepared for vision OCR."""

    assets: tuple[VisualAsset, ...]
    pages_total: int
    pages_rendered: int
    pages_truncated: bool
    deadline_reached: bool
    page_cap_reached: bool
    error: str = ""


class ArchiveLimitError(ValueError):
    """An archive exceeded configured expansion or entry limits."""


class ArchivePasswordRequired(ValueError):
    """An encrypted archive needs a password which was not supplied."""


class ArchivePasswordInvalid(ValueError):
    """An archive password was supplied but did not decrypt the archive."""


class ArchiveBackendUnavailable(ValueError):
    """A format-specific decompressor is not installed on this host."""


class ArchiveExtractionError(ValueError):
    """An archive backend failed without exposing its untrusted diagnostics."""


def _safe_archive_member_name(value: Any) -> str:
    """Validate an inert member name as if it could become a path later.

    Friday never writes members to disk, but applying one uniform zip-slip
    boundary prevents a future caller from accidentally turning a previously
    accepted name into a filesystem destination.  Backslashes are separators
    too: archives made on Windows routinely use them.
    """

    raw = str(value or "")
    normalized = raw.replace("\\", "/")
    path = PurePosixPath(normalized)
    if (
        not normalized
        or "\x00" in normalized
        or normalized.startswith("/")
        or normalized.startswith("//")
        or re.match(r"^[A-Za-z]:", normalized)
        or ".." in path.parts
    ):
        raise ArchiveLimitError("Unsafe archive member path")
    return normalized


class _Bounded7zWriter:
    """Small Py7zIO-compatible sink with hard per-member and shared limits."""

    __slots__ = ("_buffer", "_factory", "_size")

    def __init__(self, factory: _Bounded7zFactory) -> None:
        self._buffer = io.BytesIO()
        self._factory = factory
        self._size = 0

    def write(self, data: bytes | bytearray) -> int:
        if self._factory.deadline is not None and time.monotonic() >= self._factory.deadline:
            raise ArchiveExtractionError("7z extraction exceeded its deadline")
        start = self._buffer.tell()
        old_size = self._size
        new_size = max(self._size, start + len(data))
        growth = new_size - old_size
        if (
            new_size > self._factory.member_limit
            or self._factory.written + growth > self._factory.total_limit
        ):
            raise ArchiveLimitError("Decompressed 7z member exceeds configured limit")
        written = self._buffer.write(data)
        self._size = max(self._size, start + written)
        self._factory.written += max(0, self._size - old_size)
        return written

    def read(self, size: int | None = None) -> bytes:
        return self._buffer.read(-1 if size is None else size)

    def seek(self, offset: int, whence: int = 0) -> int:
        return self._buffer.seek(offset, whence)

    def flush(self) -> None:
        self._buffer.flush()

    def close(self) -> None:
        # py7zr closes each logical output after decompression.  The factory
        # still owns this in-memory buffer so Friday can consume it afterwards.
        return None

    def size(self) -> int:
        return self._size


class _Bounded7zFactory:
    """Py7zr WriterFactory-compatible owner of bounded in-memory products."""

    __slots__ = ("deadline", "member_limit", "products", "total_limit", "written")

    def __init__(
        self,
        *,
        member_limit: int,
        total_limit: int,
        deadline: float | None = None,
    ) -> None:
        self.member_limit = max(0, int(member_limit))
        self.total_limit = max(0, int(total_limit))
        self.deadline = deadline
        self.written = 0
        self.products: dict[str, _Bounded7zWriter] = {}

    def create(self, filename: str) -> _Bounded7zWriter:
        name = _safe_archive_member_name(filename)
        if name in self.products:
            raise ArchiveLimitError("Duplicate 7z member name")
        product = _Bounded7zWriter(self)
        self.products[name] = product
        return product

    def get(self, filename: str) -> _Bounded7zWriter:
        return self.products[_safe_archive_member_name(filename)]


class _BinaryReadable(Protocol):
    """Minimal binary stream contract shared by archive/decompressor readers."""

    def read(self, size: int = -1) -> bytes: ...

    def close(self) -> None: ...


class _ArchiveBudget:
    """What ONE upload may spend on unpacking, across every nesting level.

    The limits used to be per archive: each nested member started again with a
    full allowance of 24 previews and `max_archive_uncompressed_bytes`, so the
    ceilings multiplied instead of dividing. Measured: a 3.0 MB zip of 24
    `.tar.gz` members (440 x 250 KiB inside each) expanded to ~2.7 GB and held
    the event loop for 107 seconds, returning `success=True` and raising nothing
    — the operator's "250 MB per upload" was in fact 250 MB per level.
    """

    __slots__ = ("expanded_bytes", "previews")

    def __init__(self, *, previews: int, expanded_bytes: int) -> None:
        self.previews = max(0, int(previews))
        self.expanded_bytes = max(0, int(expanded_bytes))

    def take_preview(self) -> bool:
        if self.previews <= 0:
            return False
        self.previews -= 1
        return True

    def spend_bytes(self, count: int) -> None:
        self.expanded_bytes = max(0, self.expanded_bytes - max(0, int(count)))

    @property
    def spent_out(self) -> bool:
        return self.previews <= 0 or self.expanded_bytes <= 0


class DocumentExtractor:
    def __init__(
        self,
        *,
        max_archive_entries: int = 500,
        max_archive_uncompressed_bytes: int = 250 * 1024 * 1024,
        max_text_chars: int = 2_000_000,
        max_input_bytes: int = 50 * 1024 * 1024,
        parse_budget_sec: float | None = None,
        secret_values: Sequence[str] | None = None,
    ) -> None:
        # Собственные credential этого экземпляра. `None` — не «нет секретов», а
        # «взять из окружения»: умолчание закрытое, потому что забыть его можно
        # только в одном направлении, и это направление — утечка. Пустой кортеж
        # выключает проверку явно (так делают пробы самого разбора).
        #
        # Зачем это здесь, а не в вызывающем: `extract` — ЕДИНСТВЕННАЯ дорога, по
        # которой текст файла попадает и в контекст модели, и в архив, и в
        # поисковый индекс. Ворота на одной из трёх дорог не охраняли бы ничего.
        self.secret_values: tuple[str, ...] = (
            _own_secret_values() if secret_values is None else tuple(str(value) for value in secret_values)
        )
        self.max_archive_entries = max(1, int(max_archive_entries))
        self.max_archive_uncompressed_bytes = max(1024, int(max_archive_uncompressed_bytes))
        self.max_text_chars = max(10_000, int(max_text_chars))
        self.max_input_bytes = max(1, int(max_input_bytes))
        # Срок, а не таймаут: он проверяется ВНУТРИ разбора, между страницами.
        # Внешний `asyncio.timeout` вокруг `to_thread` возвращает управление
        # вызывающему, но поток не останавливает — тот продолжает жечь ядро на
        # патологической странице, и пул, общий со всей фоновой работой,
        # вычерпывается по одному потоку на каждую такую ссылку.
        #
        # Хранится БЮДЖЕТ, а срок ставится на каждый разбор и едет по стеку вызовов
        # — как `_ArchiveBudget`. Готовый `self.deadline` в конструкторе работал
        # только там, где экстрактор одноразовый (веб-путь): `IngestionPipeline`
        # строит его ОДИН раз на процесс, и такой срок сгорал бы через 8 секунд
        # после старта — дальше каждый PDF отдавал бы ноль страниц навсегда.
        # Инстанс-поле не годится и по второй причине: `ingest_file` разбирает в
        # `asyncio.to_thread`, то есть два файла делят один экстрактор в двух
        # потоках, и общий срок один затирал бы другому.
        self.parse_budget_sec: float | None = (
            None if parse_budget_sec is None else max(0.1, float(parse_budget_sec))
        )

    def _pdf_parse_deadline(self, parent_deadline: float | None) -> float | None:
        """Give pypdf its own ceiling without shortening sibling parsers."""

        if self.parse_budget_sec is None:
            return parent_deadline
        return stage_deadline(
            self.parse_budget_sec,
            parent_deadline=parent_deadline,
        )

    @staticmethod
    def _archive_parse_deadline(content: bytes, parent_deadline: float | None) -> float:
        """Bound one complete nested archive by admitted compressed size."""

        return stage_deadline(
            size_scaled_budget_sec(
                size_bytes=len(content),
                base_sec=_ARCHIVE_PARSE_BASE_TIMEOUT_SEC,
                seconds_per_mib=_ARCHIVE_PARSE_TIMEOUT_PER_MIB_SEC,
                maximum_sec=_ARCHIVE_PARSE_MAX_TIMEOUT_SEC,
            ),
            parent_deadline=parent_deadline,
        )

    def extract(
        self,
        content: bytes,
        filename: str,
        mime_type: str = "",
        *,
        archive_password: str | None = None,
        _depth: int = 0,
        _budget: _ArchiveBudget | None = None,
        _deadline: float | None = None,
    ) -> DocumentResult:
        if not isinstance(content, bytes):
            return DocumentResult("", success=False, error="Document content must be bytes")
        if _depth > _MAX_NESTING_DEPTH:
            return DocumentResult("", success=False, error="Archive nesting limit exceeded")
        password_candidates = archive_password_candidates(archive_password)
        # One budget per upload, created at the top and carried down the nesting.
        budget = _budget or _ArchiveBudget(
            previews=_MAX_ARCHIVE_PREVIEW_FILES,
            expanded_bytes=self.max_archive_uncompressed_bytes,
        )
        # A caller-owned parent deadline may bound every nested stage. Do not
        # manufacture one generic deadline here: ``parse_budget_sec`` is the
        # pypdf safety budget, while LibreOffice, RAR and visual OCR have their
        # own measured ceilings. Applying the PDF's eight seconds to every
        # format made the nominal 45-second Office and 20-second RAR budgets
        # unreachable in production.
        deadline = _deadline
        if len(content) > self.max_input_bytes:
            return DocumentResult(
                "",
                {"input_bytes": len(content), "max_input_bytes": self.max_input_bytes},
                False,
                "Document input exceeds configured size limit",
            )

        safe_name = Path(str(filename or "document")).name
        lowered = safe_name.casefold()
        ext = self._compound_extension(lowered)
        detected_mime = mime_type or mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
        normalized_mime = detected_mime.split(";", 1)[0].strip().casefold()
        suffix_known = ext in _KNOWN_DOCUMENT_EXTENSIONS
        ooxml_format = _OOXML_EXTENSIONS.get(ext)
        converted_office_format = _CONVERTED_OFFICE_FORMATS.get(ext)
        odf_format = _OPENDOCUMENT_EXTENSIONS.get(ext)
        if not suffix_known:
            ooxml_format = _OOXML_MIME_FORMATS.get(normalized_mime)
            converted_office_format = _CONVERTED_OFFICE_MIME_FORMATS.get(normalized_mime)
            if normalized_mime in _OPENDOCUMENT_MIME_TYPES:
                odf_format = "opendocument"
        archive_kind = archive_dispatch_kind(safe_name, detected_mime)
        try:
            if archive_kind is not None:
                # The same closed MIME/suffix decision is used by the eager
                # ingestion gate.  Keep it first: an explicitly declared ZIP
                # named ``.bin`` (or even with a misleading ordinary suffix)
                # must not pass persistence before its password is validated.
                # Passwords are credentials, so retries are a closed set of
                # representation variants rather than guesses.  Every attempt
                # consumes the same upload-wide budget and deadline.  Only an
                # explicit invalid-password result permits the next variant;
                # limits, missing backends and corrupt archives fail closed.
                # One archive-wide clock travels through every nested member.
                # Unlike the old generic PDF clock, it scales with the admitted
                # compressed workload and still has a finite outer ceiling.
                deadline = self._archive_parse_deadline(content, deadline)
                for candidate in password_candidates or (None,):
                    try:
                        result = self._extract_archive(
                            content,
                            safe_name,
                            archive_kind,
                            _depth,
                            budget,
                            deadline,
                            candidate,
                        )
                    except ArchivePasswordInvalid:
                        continue
                    break
                else:
                    raise ArchivePasswordInvalid
            elif ext in _TEXT_EXTENSIONS:
                result = self._extract_text(content, ext)
            elif ext in _HTML_EXTENSIONS or (
                not suffix_known and normalized_mime in {"text/html", "application/xhtml+xml"}
            ):
                result = self._extract_html(content)
            elif ext == ".pdf" or (not suffix_known and normalized_mime == "application/pdf"):
                result = self._extract_pdf(
                    content,
                    deadline=self._pdf_parse_deadline(deadline),
                )
            elif ooxml_format == "docx":
                result = self._extract_docx(content)
            elif converted_office_format == "doc":
                result = self._extract_doc(content, deadline=deadline)
            elif ext == ".msg" or (not suffix_known and normalized_mime in _MSG_MIME_TYPES):
                result = self._extract_msg(content)
            elif converted_office_format is not None:
                result = self._extract_converted_office(
                    content,
                    converted_office_format,
                    deadline=deadline,
                )
            elif ooxml_format == "xlsx":
                result = self._extract_xlsx(content)
            elif ooxml_format == "pptx":
                result = self._extract_pptx(content)
            elif odf_format is not None:
                # Таблица и презентация OpenDocument держат текст ровно там же,
                # где документ, — в `content.xml`. Принят был только `.odt`, и
                # это не решение, а недосмотр: у семьи форматов один разборщик.
                result = self._extract_xml_zip_text(content, "content.xml", odf_format)
            elif ext == ".epub" or (not suffix_known and normalized_mime == "application/epub+zip"):
                result = self._extract_epub(content)
            elif ext in {".eml", ".mht", ".mhtml"} or (
                not suffix_known and normalized_mime in {"message/rfc822", "multipart/related"}
            ):
                result = self._extract_email(content)
            elif ext == ".rtf" or (not suffix_known and normalized_mime in {"application/rtf", "text/rtf"}):
                result = self._extract_rtf(content)
            elif not suffix_known and normalized_mime.startswith("text/"):
                result = self._extract_text(content, ext or ".txt")
            else:
                result = DocumentResult(
                    "",
                    {"filename": safe_name, "mime_type": detected_mime, "size": len(content)},
                    False,
                    "unsupported_document_format",
                )
        except ArchivePasswordRequired:
            result = DocumentResult(
                "",
                {"filename": safe_name, "format": (archive_kind or ext).lstrip(".")},
                False,
                "archive_password_required",
            )
        except ArchivePasswordInvalid:
            result = DocumentResult(
                "",
                {"filename": safe_name, "format": (archive_kind or ext).lstrip(".")},
                False,
                "archive_password_invalid",
            )
        except ArchiveBackendUnavailable:
            result = DocumentResult(
                "",
                {"filename": safe_name, "format": (archive_kind or ext).lstrip(".")},
                False,
                "archive_backend_unavailable",
            )
        except ArchiveExtractionError:
            result = DocumentResult(
                "",
                {"filename": safe_name, "format": (archive_kind or ext).lstrip(".")},
                False,
                "archive_extract_failed",
            )
        except ArchiveLimitError:
            result = DocumentResult("", {"filename": safe_name}, False, "archive_limit_exceeded")
        except Exception as exc:  # defensive boundary for optional parsers
            LOGGER.info("Document extraction failed (%s)", type(exc).__name__)
            result = DocumentResult(
                "",
                {"filename": safe_name},
                False,
                f"document_extract_failed:{type(exc).__name__}",
            )

        # Technical/container metadata has one body-free parser for both new
        # uploads and later legacy hydration.  Merge its closed projection even
        # when body extraction failed (for example, a scan-only image): stored
        # properties do not depend on OCR/model availability.
        native_metadata = self.extract_document_metadata(content, safe_name, detected_mime)
        if native_metadata:
            result = DocumentResult(
                result.text,
                {**result.metadata, **native_metadata},
                result.success,
                result.error,
                result.office_structure_index,
            )

        # Убрать собственные credential ДО обреза: иначе граница могла бы
        # разрезать секрет пополам и оставить его половину в тексте.
        # The password is request-ephemeral.  It is not normally part of member
        # contents, but a hostile archive could deliberately echo it in a name or
        # file; redact that exact value before text can reach a prompt or index.
        turn_secrets = (*self.secret_values, *password_candidates)
        redacted_text, secrets_removed = _redact_own_secrets(result.text, turn_secrets)
        text = redacted_text[: self.max_text_chars]
        metadata = {
            "filename": safe_name,
            "mime_type": detected_mime,
            "input_bytes": len(content),
            **result.metadata,
        }
        # Parsers historically used three names for the same generic loss of
        # extractable text.  Normalize them here so every persistent/transient
        # adapter carries one truthful bit.  Keep deadline/page ceilings
        # distinct: their measured diagnostics are more useful than a generic
        # “text budget” warning and already make the source incomplete.
        if metadata.get("rows_truncated") or (
            metadata.get("extraction_truncated")
            and not metadata.get("parse_deadline_reached")
            and not metadata.get("pages_truncated")
        ):
            metadata["text_truncated"] = True
        if secrets_removed:
            # Потеря названа. Молчаливая подмена читалась бы как свойство
            # документа, а человеку важно знать, что его ключ лежит в файле,
            # который он только что загрузил в общий архив.
            metadata["secrets_redacted"] = secrets_removed
        if len(redacted_text) > self.max_text_chars:
            metadata["text_truncated"] = True
            metadata["original_text_chars"] = len(redacted_text)
        # Собственная дата документа снимается ЗДЕСЬ, а не внутри разборщиков, и
        # потому переживает неудачу разбора: скан без текстового слоя, битый docx
        # и файл незнакомого генератора всё равно несут дату, которую записал
        # редактор. На корпусе владельца 35 файлов не читаются вовсе — их место в
        # хронологии от этого не исчезает.
        if "document_date" not in metadata:
            own_date = (
                _pdf_document_date_from_bytes(content)
                if ext == ".pdf" or (not suffix_known and normalized_mime == "application/pdf")
                else (_office_document_date(content) if ext in _OFFICE_DATE_EXTENSIONS else None)
            )
            if own_date:
                metadata["document_date"] = own_date
        office_structure_index = (
            validate_office_structure_index(result.office_structure_index, text)
            if result.office_structure_index is not None
            else None
        )
        return DocumentResult(
            text,
            metadata,
            result.success,
            result.error,
            office_structure_index,
        )

    def extract_document_metadata(
        self,
        content: bytes,
        filename: str,
        mime_type: str = "",
    ) -> dict[str, Any]:
        """Return a closed, header-only technical metadata projection.

        This intentionally does not call ``extract``: legacy hydration may read
        authorised stored bytes solely to answer a metadata question, and must
        not parse ``content.xml``, run OCR/vision, mutate Raw, or build model
        context.  Every parser below is format-specific, byte/record bounded and
        treats embedded values as inert strings.  Unknown formats have no safe
        projection; malformed or oversized known metadata is explicitly partial.
        """

        if not isinstance(content, bytes) or len(content) > self.max_input_bytes:
            return {}
        safe_name = Path(str(filename or "document")).name
        ext = self._compound_extension(safe_name.casefold())
        detected_mime = (
            (mime_type or mimetypes.guess_type(safe_name)[0] or "").split(";", 1)[0].strip().casefold()
        )
        suffix_known = ext in _KNOWN_DOCUMENT_EXTENSIONS
        parser: Any
        if ext in _OPENDOCUMENT_EXTENSIONS or (
            not suffix_known and detected_mime in _OPENDOCUMENT_MIME_TYPES
        ):
            format_name = _OPENDOCUMENT_EXTENSIONS.get(ext, "opendocument")
            parser = self._extract_opendocument_metadata
        elif ext in _OOXML_EXTENSIONS or (not suffix_known and detected_mime in _OOXML_MIME_TYPES):
            format_name = _OOXML_EXTENSIONS.get(ext, _OOXML_MIME_FORMATS.get(detected_mime, "ooxml"))
            parser = self._extract_ooxml_metadata
        elif ext == ".pdf" or (not suffix_known and detected_mime == "application/pdf"):
            format_name = "pdf"
            parser = self._extract_pdf_metadata
        elif ext in _EMAIL_METADATA_EXTENSIONS or (
            not suffix_known and detected_mime in {"message/rfc822", "multipart/related"}
        ):
            format_name = _EMAIL_METADATA_EXTENSIONS.get(ext, "eml")
            parser = self._extract_email_metadata
        elif ext == ".epub" or (not suffix_known and detected_mime == "application/epub+zip"):
            format_name = "epub"
            parser = self._extract_epub_metadata
        elif ext in _IMAGE_EXTENSIONS or (
            not suffix_known
            and detected_mime.startswith("image/")
            and detected_mime not in _CONVERTED_OFFICE_MIME_FORMATS
        ):
            format_name = "image"
            parser = self._extract_image_metadata
        else:
            return {}
        try:
            metadata = parser(content)
        except Exception:  # noqa: BLE001 - optional metadata cannot break body extraction
            return {
                "format": format_name,
                "metadata_parse_status": "unreadable",
                "technical_metadata_incomplete": True,
            }
        return {"format": format_name, **metadata}

    @staticmethod
    def _metadata_local_name(value: Any) -> str:
        return str(value or "").rsplit("}", 1)[-1].split(":", 1)[-1]

    @staticmethod
    def _metadata_inert_text(value: Any, limit: int = 1_000) -> tuple[str, bool]:
        raw = re.sub(r"[\x00-\x1f\x7f]", " ", str(value or ""))
        cleaned = " ".join(raw.split())
        bounded = cleaned[: max(0, int(limit))]
        return bounded, len(cleaned) > len(bounded)

    @classmethod
    def _metadata_element_text(cls, element: Any, limit: int = 1_000) -> tuple[str, bool]:
        return cls._metadata_inert_text("".join(element.itertext()), limit)

    @staticmethod
    def _metadata_xml_root(raw: bytes) -> Any:
        """Parse metadata XML with encoding-safe DTD/entity/network denial."""

        from lxml import etree  # type: ignore[import-untyped]

        parser = etree.XMLParser(
            resolve_entities=False,
            no_network=True,
            load_dtd=False,
            recover=False,
            huge_tree=False,
        )
        try:
            root = etree.fromstring(raw, parser=parser)
        except (etree.LxmlError, LookupError, ValueError) as exc:
            raise ValueError("unreadable metadata XML") from exc
        # ``resolve_entities=False`` prevents expansion, but the unresolved
        # entity node would still be attacker-controlled metadata.  Reject the
        # declaration and any entity node outright.  ``docinfo.doctype`` is
        # decoded by libxml2 first, so UTF-16/UTF-32 cannot bypass this gate.
        if str(root.getroottree().docinfo.doctype or "").strip():
            raise ValueError("DTD is not permitted in metadata XML")
        if any(getattr(node, "tag", None) is etree.Entity for node in root.iter()):
            raise ValueError("entities are not permitted in metadata XML")
        return root

    def _read_metadata_xml_member(
        self,
        archive: zipfile.ZipFile,
        member_name: str,
        *,
        max_bytes: int = _MAX_TECHNICAL_METADATA_XML_BYTES,
    ) -> tuple[Any | None, str]:
        try:
            info = archive.getinfo(member_name)
        except KeyError:
            return None, "absent"
        if info.is_dir() or info.file_size > max_bytes:
            return None, "too_large"
        try:
            with archive.open(info) as stream:
                raw, truncated = self._read_stream_preview(stream, max_bytes)
            if truncated:
                return None, "unsafe_or_truncated"
            return self._metadata_xml_root(raw), "parsed"
        except (OSError, ValueError, zipfile.BadZipFile):
            return None, "unreadable"

    def _extract_ooxml_metadata(self, content: bytes) -> dict[str, Any]:
        """Read bounded OOXML core/app/custom properties without document body."""

        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            self._validate_office_zip(archive)
            parts = {
                name: self._read_metadata_xml_member(archive, name)
                for name in ("docProps/core.xml", "docProps/app.xml", "docProps/custom.xml")
            }

        metadata: dict[str, Any] = {}
        records: list[dict[str, str]] = []
        records_total = 0
        incomplete = any(status not in {"absent", "parsed"} for _root, status in parts.values())

        def add_record(source: str, name: str, value_type: str, raw_value: Any) -> str:
            nonlocal records_total, incomplete
            value, clipped = self._metadata_inert_text(raw_value, 1_000)
            safe_name, name_clipped = self._metadata_inert_text(name, 200)
            safe_type, type_clipped = self._metadata_inert_text(value_type or "string", 40)
            records_total += 1
            incomplete = incomplete or clipped or name_clipped or type_clipped
            shown = value or "(пустое значение)"
            if len(records) < _MAX_TECHNICAL_METADATA_RECORDS and safe_name and safe_type:
                records.append(
                    {
                        "source": source,
                        "name": safe_name,
                        "value_type": safe_type,
                        "value": shown,
                    }
                )
            else:
                incomplete = True
            return value

        core_root, _core_status = parts["docProps/core.xml"]
        core_map = {
            f"{{{_OOXML_DC_NS}}}title": "title",
            f"{{{_OOXML_DC_NS}}}subject": "subject",
            f"{{{_OOXML_DC_NS}}}creator": "creator",
            f"{{{_OOXML_DC_NS}}}description": "description",
            f"{{{_OOXML_DC_NS}}}language": "language",
            f"{{{_OOXML_DC_NS}}}identifier": "identifier",
            f"{{{_OOXML_CP_NS}}}lastModifiedBy": "last_modified_by",
            f"{{{_OOXML_CP_NS}}}revision": "revision",
            f"{{{_OOXML_CP_NS}}}category": "category",
            f"{{{_OOXML_CP_NS}}}contentStatus": "content_status",
            f"{{{_OOXML_CP_NS}}}version": "version",
            f"{{{_OOXML_DCTERMS_NS}}}created": "creation_date",
            f"{{{_OOXML_DCTERMS_NS}}}modified": "modified_date",
            f"{{{_OOXML_CP_NS}}}lastPrinted": "print_date",
        }
        if core_root is not None:
            for element in list(core_root):
                value, clipped = self._metadata_element_text(element, 4_000)
                incomplete = incomplete or clipped
                local_name = self._metadata_local_name(element.tag)
                add_record("OOXML core", local_name, "string", value)
                mapped = core_map.get(element.tag)
                if mapped and value and mapped not in metadata:
                    metadata[mapped] = value
                if element.tag == f"{{{_OOXML_CP_NS}}}keywords" and value:
                    keywords = [item.strip() for item in re.split(r"[;,\n]", value) if item.strip()]
                    metadata["keywords"] = keywords[:32]
                    metadata["keywords_total"] = len(keywords)
                    metadata["keywords_shown"] = min(len(keywords), 32)
                    incomplete = incomplete or len(keywords) > 32

        app_root, _app_status = parts["docProps/app.xml"]
        app_strings = {
            "Application": "application",
            "AppVersion": "application_version",
            "Company": "company",
            "Manager": "manager",
            "Template": "template_name",
            "PresentationFormat": "presentation_format",
        }
        app_counts = {
            "TotalTime": "total_editing_time_minutes",
            "Pages": "page_count",
            "Words": "word_count",
            "Characters": "character_count",
            "CharactersWithSpaces": "characters_with_spaces",
            "Lines": "line_count",
            "Paragraphs": "paragraph_count",
            "Slides": "slide_count",
            "Notes": "note_count",
            "HiddenSlides": "hidden_slide_count",
            "MMClips": "multimedia_clip_count",
            "DocSecurity": "document_security",
        }
        app_booleans = {
            "ScaleCrop": "scale_crop",
            "LinksUpToDate": "links_up_to_date",
            "SharedDoc": "shared_document",
            "HyperlinksChanged": "hyperlinks_changed",
        }
        if app_root is not None:
            for element in list(app_root):
                value, clipped = self._metadata_element_text(element, 4_000)
                incomplete = incomplete or clipped
                local_name = self._metadata_local_name(element.tag)
                add_record("OOXML app", local_name, "string", value)
                if element.tag != f"{{{_OOXML_EP_NS}}}{local_name}":
                    continue
                if local_name in app_strings and value:
                    metadata.setdefault(app_strings[local_name], value)
                elif local_name in app_counts and value:
                    if value.isdecimal():
                        parsed = int(value)
                        metadata[app_counts[local_name]] = min(parsed, 2_147_483_647)
                        incomplete = incomplete or parsed > 2_147_483_647
                    else:
                        incomplete = True
                elif local_name in app_booleans and value:
                    normalized = value.casefold()
                    if normalized in {"true", "1", "yes"}:
                        metadata[app_booleans[local_name]] = True
                    elif normalized in {"false", "0", "no"}:
                        metadata[app_booleans[local_name]] = False
                    else:
                        incomplete = True

        custom_root, _custom_status = parts["docProps/custom.xml"]
        if custom_root is not None:
            for prop in list(custom_root):
                if prop.tag != f"{{{_OOXML_CUSTOM_NS}}}property":
                    incomplete = True
                    continue
                name = str(prop.attrib.get("name") or "")
                children = list(prop)
                if len(children) != 1:
                    add_record("OOXML custom", name or "property", "unknown", "")
                    incomplete = True
                    continue
                value, clipped = self._metadata_element_text(children[0], 4_000)
                incomplete = incomplete or clipped
                value_type = self._metadata_local_name(children[0].tag)
                if not str(children[0].tag).startswith(f"{{{_OOXML_VT_NS}}}"):
                    incomplete = True
                add_record("OOXML custom", name or "property", value_type, value)

        statuses = [status for _root, status in parts.values()]
        if records:
            metadata["stored_properties"] = records
        metadata["stored_properties_total"] = records_total
        metadata["stored_properties_shown"] = len(records)
        if records_total > len(records):
            incomplete = True
        if all(status == "absent" for status in statuses):
            metadata["metadata_parse_status"] = "absent"
        elif incomplete:
            metadata["metadata_parse_status"] = "partial"
            metadata["technical_metadata_incomplete"] = True
        else:
            metadata["metadata_parse_status"] = "parsed"
        own_date = _plausible_document_date(
            str(metadata.get("creation_date") or metadata.get("modified_date") or "")
        )
        if own_date:
            metadata["document_date"] = own_date
        return self._redact_metadata_value(metadata)

    def _extract_pdf_metadata(self, content: bytes) -> dict[str, Any]:
        """Read PDF Info/XMP and stored signature-field facts, never validate."""

        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content), strict=False)
        if reader.is_encrypted:
            try:
                if reader.decrypt("") == 0:
                    return {
                        "metadata_parse_status": "encrypted",
                        "technical_metadata_incomplete": True,
                    }
            except Exception:
                return {
                    "metadata_parse_status": "encrypted",
                    "technical_metadata_incomplete": True,
                }

        metadata: dict[str, Any] = {}
        records: list[dict[str, str]] = []
        records_total = 0
        incomplete = False

        def add_record(source: str, name: Any, value_type: str, raw_value: Any) -> str:
            nonlocal records_total, incomplete
            safe_name, name_clipped = self._metadata_inert_text(name, 200)
            value, value_clipped = self._metadata_inert_text(raw_value, 1_000)
            records_total += 1
            incomplete = incomplete or name_clipped or value_clipped
            if len(records) < _MAX_TECHNICAL_METADATA_RECORDS and safe_name:
                records.append(
                    {
                        "source": source,
                        "name": safe_name,
                        "value_type": value_type,
                        "value": value or "(пустое значение)",
                    }
                )
            else:
                incomplete = True
            return value

        info_map = {
            "/Title": "title",
            "/Author": "creator",
            "/Subject": "subject",
            "/Creator": "generator",
            "/Producer": "producer",
            "/CreationDate": "creation_date",
            "/ModDate": "modified_date",
            "/Trapped": "trapped",
        }
        try:
            info: Any = reader.metadata or {}
            for key, raw_value in info.items():
                name = str(key or "").lstrip("/") or "property"
                value = add_record("PDF Info", name, "string", raw_value)
                mapped = info_map.get(str(key))
                if mapped and value and mapped not in metadata:
                    metadata[mapped] = value
                if str(key) == "/Keywords" and value:
                    keywords = [item.strip() for item in re.split(r"[;,\n]", value) if item.strip()]
                    metadata["keywords"] = keywords[:32]
                    metadata["keywords_total"] = len(keywords)
                    metadata["keywords_shown"] = min(len(keywords), 32)
                    incomplete = incomplete or len(keywords) > 32
        except Exception:
            incomplete = True
        header = str(getattr(reader, "pdf_header", "") or "")
        if header.startswith("%PDF-"):
            metadata["pdf_version"] = header[5:21]

        root: Any = None
        try:
            root = reader.trailer["/Root"].get_object()
        except Exception:
            incomplete = True

        # XMP is often an uncompressed XML metadata stream.  Never ask pypdf to
        # inflate an arbitrary filtered stream without an output limit: a tiny
        # compressed stream could otherwise allocate unbounded memory.  Such a
        # packet is reported partial rather than silently skipped.
        if root is not None:
            try:
                metadata_ref = root.get("/Metadata")
                if metadata_ref is not None:
                    stream = metadata_ref.get_object()
                    filters = stream.get("/Filter")
                    raw_stream = bytes(getattr(stream, "_data", b""))
                    xmp_bytes: bytes | None = None
                    if len(raw_stream) > _MAX_TECHNICAL_METADATA_XML_BYTES:
                        incomplete = True
                    elif not filters:
                        xmp_bytes = raw_stream
                    else:
                        filter_names = (
                            [str(item) for item in filters]
                            if isinstance(filters, (list, tuple))
                            else [str(filters)]
                        )
                        decode_params = stream.get("/DecodeParms")
                        if filter_names == ["/FlateDecode"] and not decode_params:
                            try:
                                inflater = zlib.decompressobj()
                                expanded = inflater.decompress(
                                    raw_stream,
                                    _MAX_TECHNICAL_METADATA_XML_BYTES + 1,
                                )
                                if (
                                    len(expanded) > _MAX_TECHNICAL_METADATA_XML_BYTES
                                    or inflater.unconsumed_tail
                                ):
                                    incomplete = True
                                else:
                                    remaining = _MAX_TECHNICAL_METADATA_XML_BYTES + 1 - len(expanded)
                                    expanded += inflater.flush(max(1, remaining))
                                    if len(expanded) > _MAX_TECHNICAL_METADATA_XML_BYTES or not inflater.eof:
                                        incomplete = True
                                    else:
                                        xmp_bytes = expanded
                            except zlib.error:
                                incomplete = True
                        else:
                            incomplete = True
                    if xmp_bytes:
                        try:
                            xmp_root = self._metadata_xml_root(xmp_bytes)
                        except Exception:
                            incomplete = True
                        else:
                            scanned = 0

                            def walk_xmp(element: Any, path: tuple[str, ...] = ()) -> None:
                                nonlocal scanned, incomplete
                                if scanned >= 256:
                                    incomplete = True
                                    return
                                local = self._metadata_local_name(element.tag)
                                current = (*path[-2:], local)
                                children = list(element)
                                if not children:
                                    scanned += 1
                                    value, clipped = self._metadata_element_text(element, 4_000)
                                    incomplete = incomplete or clipped
                                    name = "/".join(current)
                                    add_record("PDF XMP", name, "xml", value)
                                    folded_path = {item.casefold() for item in current}
                                    for cue, target in (
                                        ("title", "title"),
                                        ("creator", "creator"),
                                        ("description", "description"),
                                        ("language", "language"),
                                        ("createdate", "creation_date"),
                                        ("modifydate", "modified_date"),
                                    ):
                                        if cue in folded_path and value and target not in metadata:
                                            metadata[target] = value
                                for attr_name, attr_value in element.attrib.items():
                                    scanned += 1
                                    if scanned > 256:
                                        incomplete = True
                                        break
                                    add_record(
                                        "PDF XMP",
                                        f"{'/'.join(current)}/@{self._metadata_local_name(attr_name)}",
                                        "attribute",
                                        attr_value,
                                    )
                                for child in children:
                                    walk_xmp(child, current)

                            walk_xmp(xmp_root)
            except Exception:
                incomplete = True

        signature_fields: list[dict[str, str]] = []
        signature_total = 0
        if root is not None:
            try:
                acroform_ref = root.get("/AcroForm")
                if acroform_ref is not None:
                    acroform = acroform_ref.get_object()
                    stack: list[tuple[Any, str, Any, int]] = [
                        (field, "", None, 0) for field in list(acroform.get("/Fields") or [])
                    ]
                    seen: set[tuple[int, int] | int] = set()
                    scanned = 0
                    while stack:
                        field_ref, parent_name, inherited_type, depth = stack.pop()
                        if depth > 8 or scanned >= _MAX_PDF_FORM_FIELDS_SCANNED:
                            incomplete = True
                            break
                        field = field_ref.get_object()
                        identity_ref = getattr(field_ref, "idnum", None)
                        generation = getattr(field_ref, "generation", 0)
                        identity: tuple[int, int] | int = (
                            (int(identity_ref), int(generation))
                            if isinstance(identity_ref, int)
                            else id(field)
                        )
                        if identity in seen:
                            continue
                        seen.add(identity)
                        scanned += 1
                        own_name, clipped = self._metadata_inert_text(field.get("/T"), 200)
                        incomplete = incomplete or clipped
                        field_name = ".".join(item for item in (parent_name, own_name) if item)
                        field_type = field.get("/FT") or inherited_type
                        if str(field_type or "") == "/Sig":
                            signature_total += 1
                            signature: Any = field.get("/V")
                            signature = signature.get_object() if signature is not None else None
                            row: dict[str, str] = {
                                "field_name": field_name or f"signature-{signature_total}",
                            }
                            if signature is not None:
                                for source_key, target_key in (
                                    ("/Name", "signer_name"),
                                    ("/M", "signing_time"),
                                    ("/Reason", "reason"),
                                    ("/Location", "location"),
                                    ("/ContactInfo", "contact_info"),
                                    ("/Filter", "filter"),
                                    ("/SubFilter", "subfilter"),
                                ):
                                    value, was_clipped = self._metadata_inert_text(
                                        signature.get(source_key), 500
                                    )
                                    incomplete = incomplete or was_clipped
                                    if value:
                                        row[target_key] = value
                                if signature.get("/ByteRange") is not None:
                                    row["byte_range_present"] = "да"
                                if signature.get("/Contents") is not None:
                                    row["contents_present"] = "да"
                            if len(signature_fields) < _MAX_PDF_SIGNATURE_FIELDS:
                                signature_fields.append(row)
                            else:
                                incomplete = True
                        for child in list(field.get("/Kids") or []):
                            stack.append((child, field_name, field_type, depth + 1))
            except Exception:
                incomplete = True

        if records:
            metadata["stored_properties"] = records
        metadata["stored_properties_total"] = records_total
        metadata["stored_properties_shown"] = len(records)
        if signature_fields:
            metadata["signature_fields"] = signature_fields
        metadata["signature_fields_total"] = signature_total
        metadata["signature_fields_shown"] = len(signature_fields)
        if signature_total:
            metadata["signature_count"] = signature_total
            metadata["signature_validity"] = "not_checked"
        if records_total > len(records) or signature_total > len(signature_fields):
            incomplete = True
        own_date = _pdf_document_date(reader)
        if own_date:
            metadata["document_date"] = own_date
        if incomplete:
            metadata["metadata_parse_status"] = "partial"
            metadata["technical_metadata_incomplete"] = True
        else:
            metadata["metadata_parse_status"] = "parsed" if records_total or signature_total else "absent"
        return self._redact_metadata_value(metadata)

    def _extract_email_metadata(self, content: bytes) -> dict[str, Any]:
        """Project bounded RFC/MHTML headers only; MIME bodies stay unread."""

        from email import policy
        from email.parser import BytesParser

        source = content[:_MAX_EMAIL_METADATA_BYTES]
        incomplete = len(content) > len(source)
        message = BytesParser(policy=policy.default).parsebytes(source, headersonly=True)
        records: list[dict[str, str]] = []
        total = 0
        metadata: dict[str, Any] = {}
        header_map = {
            "from": "email_from",
            "to": "email_to",
            "cc": "email_cc",
            "bcc": "email_bcc",
            "sender": "email_sender",
            "reply-to": "email_reply_to",
            "subject": "email_subject",
            "date": "email_date",
            "message-id": "message_id",
            "in-reply-to": "in_reply_to",
            "references": "references",
            "content-type": "email_content_type",
            "content-language": "content_language",
        }
        for raw_name, raw_value in message.raw_items():
            total += 1
            name, name_clipped = self._metadata_inert_text(raw_name, 120)
            value, value_clipped = self._metadata_inert_text(raw_value, 2_000)
            incomplete = incomplete or name_clipped or value_clipped
            if len(records) < _MAX_TECHNICAL_METADATA_RECORDS and name:
                records.append(
                    {
                        "source": "RFC header",
                        "name": name,
                        "value_type": "header",
                        "value": value or "(пустое значение)",
                    }
                )
            else:
                incomplete = True
            mapped = header_map.get(name.casefold())
            if mapped and value and mapped not in metadata:
                metadata[mapped] = value
        # Common fields are decoded for display; the complete bounded raw
        # header ledger above remains available as stored technical evidence.
        for header_name, target in header_map.items():
            decoded, clipped = self._metadata_inert_text(
                str(message.get(header_name) or ""),
                2_000,
            )
            incomplete = incomplete or clipped
            if decoded:
                metadata[target] = decoded
        if records:
            metadata["stored_properties"] = records
        metadata["stored_properties_total"] = total
        metadata["stored_properties_shown"] = len(records)
        raw_date = str(message.get("Date") or "")
        own_date = _plausible_document_date(_email_iso_date(raw_date))
        if own_date:
            metadata["document_date"] = own_date
        if incomplete or total > len(records):
            metadata["metadata_parse_status"] = "partial"
            metadata["technical_metadata_incomplete"] = True
        else:
            metadata["metadata_parse_status"] = "parsed" if total else "absent"
        return self._redact_metadata_value(metadata)

    def _extract_epub_metadata(self, content: bytes) -> dict[str, Any]:
        """Read the package OPF metadata selected by EPUB container.xml."""

        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            self._validate_office_zip(archive)
            container, container_status = self._read_metadata_xml_member(
                archive,
                "META-INF/container.xml",
                max_bytes=_MAX_EPUB_CONTAINER_BYTES,
            )
            if container is None:
                return {
                    "metadata_parse_status": (
                        "missing_container" if container_status == "absent" else container_status
                    ),
                    "technical_metadata_incomplete": True,
                }
            rootfiles = [
                item
                for item in container.iter()
                if self._metadata_local_name(item.tag) == "rootfile"
                and str(item.attrib.get("full-path") or "").strip()
            ]
            if not rootfiles:
                return {
                    "metadata_parse_status": "missing_package_metadata",
                    "technical_metadata_incomplete": True,
                }
            package_path = str(rootfiles[0].attrib.get("full-path") or "")
            _safe_archive_member_name(package_path)
            package, package_status = self._read_metadata_xml_member(
                archive,
                package_path,
                max_bytes=_MAX_EPUB_PACKAGE_BYTES,
            )

        if package is None:
            return {
                "metadata_parse_status": package_status,
                "technical_metadata_incomplete": True,
            }
        metadata_element = next(
            (item for item in package.iter() if self._metadata_local_name(item.tag) == "metadata"),
            None,
        )
        if metadata_element is None:
            return {
                "metadata_parse_status": "missing_package_metadata",
                "technical_metadata_incomplete": True,
            }

        metadata: dict[str, Any] = {}
        records: list[dict[str, str]] = []
        total = 0
        incomplete = len(rootfiles) > 1
        subjects: list[str] = []
        common_map = {
            "title": "title",
            "creator": "creator",
            "description": "description",
            "language": "language",
            "identifier": "identifier",
            "publisher": "publisher",
            "rights": "rights",
            "source": "source",
            "coverage": "coverage",
            "relation": "relation",
        }
        for element in list(metadata_element):
            total += 1
            local_name = self._metadata_local_name(element.tag)
            value, clipped = self._metadata_element_text(element, 4_000)
            incomplete = incomplete or clipped
            attributes: list[str] = []
            for raw_name, raw_value in element.attrib.items():
                name, name_clipped = self._metadata_inert_text(self._metadata_local_name(raw_name), 100)
                attr_value, value_clipped = self._metadata_inert_text(raw_value, 500)
                incomplete = incomplete or name_clipped or value_clipped
                if name:
                    attributes.append(f"{name}={attr_value}")
            shown_value = value
            if attributes:
                shown_value = f"{value} [{' ; '.join(attributes)}]".strip()
            shown_value, shown_clipped = self._metadata_inert_text(shown_value, 1_000)
            incomplete = incomplete or shown_clipped
            if len(records) < _MAX_TECHNICAL_METADATA_RECORDS:
                records.append(
                    {
                        "source": "EPUB OPF",
                        "name": local_name,
                        "value_type": "xml",
                        "value": shown_value or "(пустое значение)",
                    }
                )
            else:
                incomplete = True
            mapped = common_map.get(local_name.casefold())
            if mapped and value and mapped not in metadata:
                metadata[mapped] = value
            if local_name.casefold() == "subject" and value:
                subjects.append(value)
            property_name = str(element.attrib.get("property") or "").casefold()
            if property_name == "dcterms:modified" and value:
                metadata.setdefault("modified_date", value)
        if subjects:
            metadata["keywords"] = subjects[:32]
            metadata["keywords_total"] = len(subjects)
            metadata["keywords_shown"] = min(len(subjects), 32)
            incomplete = incomplete or len(subjects) > 32
        if records:
            metadata["stored_properties"] = records
        metadata["stored_properties_total"] = total
        metadata["stored_properties_shown"] = len(records)
        if total > len(records):
            incomplete = True
        own_date = _plausible_document_date(str(metadata.get("modified_date") or ""))
        if own_date:
            metadata["document_date"] = own_date
        if incomplete:
            metadata["metadata_parse_status"] = "partial"
            metadata["technical_metadata_incomplete"] = True
        else:
            metadata["metadata_parse_status"] = "parsed" if total else "absent"
        return self._redact_metadata_value(metadata)

    def _extract_image_metadata(self, content: bytes) -> dict[str, Any]:
        """Read dimensions and a capped EXIF ledger without decoding pixels."""

        from PIL import ExifTags, Image

        metadata: dict[str, Any] = {}
        records: list[dict[str, str]] = []
        total = 0
        incomplete = False

        def exif_value(raw_value: Any) -> tuple[str, bool]:
            if isinstance(raw_value, bytes):
                clipped = len(raw_value) > 4_096
                sample = raw_value[:4_096]
                decoded = sample.decode("utf-8", errors="replace")
                if decoded.count("�") > max(2, len(decoded) // 10):
                    return f"[двоичные данные: {len(raw_value)} байт]", clipped
                value, text_clipped = self._metadata_inert_text(decoded, 1_000)
                return value, clipped or text_clipped
            if isinstance(raw_value, (list, tuple)):
                clipped = len(raw_value) > 32
                value, text_clipped = self._metadata_inert_text(
                    ", ".join(str(item) for item in raw_value[:32]), 1_000
                )
                return value, clipped or text_clipped
            return self._metadata_inert_text(raw_value, 1_000)

        with Image.open(io.BytesIO(content)) as image:
            width, height = image.size
            metadata.update(
                {
                    "image_format": str(image.format or "").upper()[:32],
                    "image_mode": str(image.mode or "")[:32],
                    "width_pixels": max(0, int(width)),
                    "height_pixels": max(0, int(height)),
                    "image_frame_count": max(1, int(getattr(image, "n_frames", 1) or 1)),
                    "image_animated": bool(getattr(image, "is_animated", False)),
                }
            )
            exif = image.getexif()
            ledgers: list[tuple[str, Any]] = [("EXIF", exif)]
            for tag_id, prefix in ((0x8769, "EXIF sub-IFD"), (0x8825, "GPS IFD")):
                with suppress(Exception):
                    nested = exif.get_ifd(tag_id)
                    if nested:
                        ledgers.append((prefix, nested))
            seen: set[tuple[str, int]] = set()
            for source, ledger in ledgers:
                for raw_tag, raw_value in ledger.items():
                    tag = int(raw_tag)
                    identity = (source, tag)
                    if identity in seen:
                        continue
                    seen.add(identity)
                    total += 1
                    name = str(ExifTags.TAGS.get(tag) or ExifTags.GPSTAGS.get(tag) or f"Tag {tag}")
                    value, clipped = exif_value(raw_value)
                    incomplete = incomplete or clipped
                    if len(records) < _MAX_TECHNICAL_METADATA_RECORDS:
                        records.append(
                            {
                                "source": source,
                                "name": name[:200],
                                "value_type": type(raw_value).__name__[:40],
                                "value": value or "(пустое значение)",
                            }
                        )
                    else:
                        incomplete = True
                    if name == "Make" and value:
                        metadata.setdefault("camera_make", value)
                    elif name == "Model" and value:
                        metadata.setdefault("camera_model", value)
                    elif name == "DateTimeOriginal" and value:
                        metadata.setdefault("capture_date", value)
                    elif name == "Orientation" and value:
                        metadata.setdefault("image_orientation", value)
        if records:
            metadata["stored_properties"] = records
        metadata["stored_properties_total"] = total
        metadata["stored_properties_shown"] = len(records)
        if total > len(records):
            incomplete = True
        if incomplete:
            metadata["metadata_parse_status"] = "partial"
            metadata["technical_metadata_incomplete"] = True
        else:
            metadata["metadata_parse_status"] = "parsed"
        return self._redact_metadata_value(metadata)

    @staticmethod
    def _bounded_odf_text(element: Any, limit: int) -> str:
        text = " ".join("".join(element.itertext()).split())
        return text[: max(0, int(limit))]

    @staticmethod
    def _bounded_odf_value(value: Any, limit: int) -> str:
        text = re.sub(r"[\x00-\x1f\x7f]", " ", str(value or ""))
        return " ".join(text.split())[: max(0, int(limit))]

    @staticmethod
    def _odf_iso_datetime(value: str) -> str:
        candidate = str(value or "").strip()
        if not re.fullmatch(
            r"\d{4}-\d{2}-\d{2}(?:T\d{2}:\d{2}:\d{2}(?:\.\d{1,9})?(?:Z|[+-]\d{2}:\d{2})?)?",
            candidate,
        ):
            return ""
        with suppress(ValueError):
            datetime.fromisoformat(candidate.replace("Z", "+00:00"))
            return candidate
        return ""

    def _extract_odf_signature_metadata(
        self,
        archive: zipfile.ZipFile,
        members: Sequence[zipfile.ZipInfo],
    ) -> dict[str, Any]:
        """Describe stored ODF signature XML without validating a signature.

        ODF permits signature files below ``META-INF/`` whose names contain
        ``signatures``.  Presence and inert XML fields are container metadata;
        cryptographic validity requires certificate/path/revocation checks and
        is deliberately outside this parser.
        """

        all_candidates = sorted(
            member.filename
            for member in members
            if not member.is_dir()
            and member.filename.startswith("META-INF/")
            and "signatures" in PurePosixPath(member.filename).name.casefold()
            and member.filename.casefold().endswith(".xml")
        )
        candidates = all_candidates[:_MAX_ODF_SIGNATURE_MEMBERS]
        if not candidates:
            return {}

        signature_ids: list[str] = []
        signature_subjects: list[str] = []
        signature_times: list[str] = []
        signature_ids_total = 0
        signature_subjects_total = 0
        signature_times_total = 0
        signatures = 0
        incomplete = len(all_candidates) > len(candidates)
        for name in candidates:
            info = archive.getinfo(name)
            if info.file_size > _MAX_ODF_SIGNATURE_BYTES:
                incomplete = True
                continue
            try:
                with archive.open(info) as stream:
                    raw, truncated = self._read_stream_preview(stream, _MAX_ODF_SIGNATURE_BYTES)
                if truncated:
                    incomplete = True
                    continue
                root = self._metadata_xml_root(raw)
            except (KeyError, OSError, ValueError, zipfile.BadZipFile):
                incomplete = True
                continue
            if root.tag.rsplit("}", 1)[-1] != "document-signatures":
                incomplete = True
                continue
            for signature in root.iter(f"{{{_XMLDSIG_NS}}}Signature"):
                signatures += 1
                raw_signature_id = self._bounded_odf_value(signature.attrib.get("Id"), 10_000)
                signature_id = raw_signature_id[:200]
                incomplete = incomplete or len(raw_signature_id) > len(signature_id)
                if signature_id:
                    signature_ids_total += 1
                    if signature_id not in signature_ids and len(signature_ids) < 16:
                        signature_ids.append(signature_id)
            for element in root.iter():
                local_name = element.tag.rsplit("}", 1)[-1]
                if local_name == "X509SubjectName":
                    raw_value = self._bounded_odf_text(element, 10_000)
                    value = raw_value[:500]
                    incomplete = incomplete or len(raw_value) > len(value)
                    if value:
                        signature_subjects_total += 1
                        if value not in signature_subjects and len(signature_subjects) < 16:
                            signature_subjects.append(value)
                elif local_name == "SigningTime":
                    raw_value = self._bounded_odf_text(element, 10_000)
                    value = raw_value[:80]
                    incomplete = incomplete or len(raw_value) > len(value)
                    if value:
                        signature_times_total += 1
                        if value not in signature_times and len(signature_times) < 16:
                            signature_times.append(value)

        metadata: dict[str, Any] = {
            "signature_members": candidates,
            "signature_members_total": len(all_candidates),
            "signature_members_shown": len(candidates),
            "signature_count": signatures,
            "signature_validity": "not_checked",
        }
        if signature_ids:
            metadata["signature_ids"] = signature_ids
        metadata["signature_ids_total"] = signature_ids_total
        metadata["signature_ids_shown"] = len(signature_ids)
        if signature_subjects:
            metadata["signature_subjects"] = signature_subjects
        metadata["signature_subjects_total"] = signature_subjects_total
        metadata["signature_subjects_shown"] = len(signature_subjects)
        if signature_times:
            metadata["signature_times"] = signature_times
        metadata["signature_times_total"] = signature_times_total
        metadata["signature_times_shown"] = len(signature_times)
        if (
            signature_ids_total > len(signature_ids)
            or signature_subjects_total > len(signature_subjects)
            or signature_times_total > len(signature_times)
        ):
            incomplete = True
        if incomplete:
            metadata["signature_metadata_incomplete"] = True
        return metadata

    def _redact_metadata_value(self, value: Any) -> Any:
        if isinstance(value, str):
            return _redact_own_secrets(value, self.secret_values)[0]
        if isinstance(value, list):
            return [self._redact_metadata_value(item) for item in value]
        if isinstance(value, dict):
            return {str(key): self._redact_metadata_value(item) for key, item in value.items()}
        return value

    def _extract_opendocument_metadata(self, content: bytes) -> dict[str, Any]:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = self._validate_office_zip(archive)
            names = {member.filename for member in members}
            signature_metadata = self._extract_odf_signature_metadata(archive, members)
            raw = b""
            truncated = False
            metadata_parse_status = "absent"
            if _ODF_META_MEMBER in names:
                info = archive.getinfo(_ODF_META_MEMBER)
                if info.file_size <= _MAX_ODF_METADATA_BYTES:
                    with archive.open(info) as stream:
                        raw, truncated = self._read_stream_preview(stream, _MAX_ODF_METADATA_BYTES)
                    metadata_parse_status = "read"
                else:
                    metadata_parse_status = "too_large"
        if not raw:
            result = {**signature_metadata, "metadata_parse_status": metadata_parse_status}
            if metadata_parse_status != "absent":
                result["technical_metadata_incomplete"] = True
            return result
        if truncated:
            return {
                **signature_metadata,
                "metadata_parse_status": "unsafe_or_truncated",
                "technical_metadata_incomplete": True,
            }
        try:
            root = self._metadata_xml_root(raw)
        except ValueError:
            return {
                **signature_metadata,
                "metadata_parse_status": "unreadable",
                "technical_metadata_incomplete": True,
            }
        office_meta = root.find(f".//{{{_ODF_OFFICE_NS}}}meta")
        if office_meta is None:
            return {
                **signature_metadata,
                "metadata_parse_status": "missing_office_meta",
                "technical_metadata_incomplete": True,
            }

        metadata: dict[str, Any] = {
            **signature_metadata,
            "metadata_parse_status": "parsed",
        }
        technical_incomplete = signature_metadata.get("signature_metadata_incomplete") is True
        keywords: list[str] = []
        keywords_total = 0
        user_defined: list[dict[str, str]] = []
        user_defined_total = 0

        def element_text(element: Any, limit: int) -> str:
            nonlocal technical_incomplete
            value = " ".join("".join(element.itertext()).split())
            if len(value) > limit:
                technical_incomplete = True
            return value[:limit]

        def attribute_text(value: Any, limit: int) -> str:
            nonlocal technical_incomplete
            cleaned = re.sub(r"[\x00-\x1f\x7f]", " ", str(value or ""))
            cleaned = " ".join(cleaned.split())
            if len(cleaned) > limit:
                technical_incomplete = True
            return cleaned[:limit]

        for element in list(office_meta):
            string_spec = _ODF_STRING_TAGS.get(element.tag)
            if string_spec is not None and string_spec[0] not in metadata:
                key, limit = string_spec
                value = element_text(element, limit)
                if (
                    key == "language"
                    and value
                    and not re.fullmatch(r"[A-Za-z]{2,8}(?:-[A-Za-z0-9]{1,8})*", value)
                ):
                    value = ""
                    technical_incomplete = True
                if value:
                    metadata[key] = value
                continue
            if element.tag == f"{{{_ODF_META_NS}}}keyword":
                keywords_total += 1
                keyword = element_text(element, 200)
                if keyword and keyword not in keywords and len(keywords) < 32:
                    keywords.append(keyword)
                elif keyword:
                    technical_incomplete = True
                continue
            if element.tag == f"{{{_ODF_META_NS}}}creation-date":
                raw_value = element_text(element, 64)
                value = self._odf_iso_datetime(raw_value)
                if value:
                    metadata.setdefault("creation_date", value)
                elif raw_value:
                    technical_incomplete = True
                continue
            if element.tag == f"{{{_ODF_DC_NS}}}date":
                raw_value = element_text(element, 64)
                value = self._odf_iso_datetime(raw_value)
                if value:
                    metadata.setdefault("modified_date", value)
                elif raw_value:
                    technical_incomplete = True
                continue
            if element.tag == f"{{{_ODF_META_NS}}}print-date":
                raw_value = element_text(element, 64)
                value = self._odf_iso_datetime(raw_value)
                if value:
                    metadata.setdefault("print_date", value)
                elif raw_value:
                    technical_incomplete = True
                continue
            if element.tag == f"{{{_ODF_META_NS}}}editing-cycles":
                value = element_text(element, 16)
                if value.isdecimal():
                    parsed = int(value)
                    metadata["editing_cycles"] = min(parsed, 2_147_483_647)
                    technical_incomplete = technical_incomplete or parsed > 2_147_483_647
                elif value:
                    technical_incomplete = True
                continue
            if element.tag == f"{{{_ODF_META_NS}}}editing-duration":
                value = element_text(element, 64)
                if re.fullmatch(
                    r"P(?=\d|T)(?:\d+D)?(?:T(?:\d+H)?(?:\d+M)?(?:\d+(?:\.\d+)?S)?)?",
                    value,
                ):
                    metadata["editing_duration"] = value
                elif value:
                    technical_incomplete = True
                continue
            if element.tag == f"{{{_ODF_META_NS}}}document-statistic":
                for attribute, key in _ODF_STATISTIC_ATTRIBUTES.items():
                    value = str(element.attrib.get(attribute) or "")
                    if value.isdecimal():
                        try:
                            parsed = int(value)
                        except ValueError:
                            technical_incomplete = True
                            continue
                        metadata[key] = min(parsed, 2_147_483_647)
                        technical_incomplete = technical_incomplete or parsed > 2_147_483_647
                    elif value:
                        technical_incomplete = True
                continue
            if element.tag == f"{{{_ODF_META_NS}}}template":
                template: dict[str, str] = {}
                title = attribute_text(element.attrib.get(f"{{{_ODF_META_NS}}}title"), 500)
                date_value = self._odf_iso_datetime(
                    attribute_text(element.attrib.get(f"{{{_ODF_META_NS}}}date"), 64)
                )
                href = attribute_text(element.attrib.get(f"{{{_ODF_XLINK_NS}}}href"), 1_000)
                if title:
                    template["title"] = title
                if date_value:
                    template["date"] = date_value
                if href:
                    template["href"] = href
                if template:
                    metadata["template"] = template
                continue
            if element.tag == f"{{{_ODF_META_NS}}}auto-reload":
                auto_reload: dict[str, str] = {}
                href = attribute_text(element.attrib.get(f"{{{_ODF_XLINK_NS}}}href"), 1_000)
                delay = attribute_text(element.attrib.get(f"{{{_ODF_META_NS}}}delay"), 64)
                if href:
                    auto_reload["href"] = href
                if re.fullmatch(
                    r"P(?=\d|T)(?:\d+D)?(?:T(?:\d+H)?(?:\d+M)?(?:\d+(?:\.\d+)?S)?)?",
                    delay,
                ):
                    auto_reload["delay"] = delay
                elif delay:
                    technical_incomplete = True
                if auto_reload:
                    metadata["auto_reload"] = auto_reload
                continue
            if element.tag == f"{{{_ODF_META_NS}}}hyperlink-behaviour":
                behaviour: dict[str, str] = {}
                target = attribute_text(element.attrib.get(f"{{{_ODF_OFFICE_NS}}}target-frame-name"), 200)
                show = attribute_text(element.attrib.get(f"{{{_ODF_XLINK_NS}}}show"), 32)
                if target:
                    behaviour["target_frame_name"] = target
                if show in {"new", "replace", "embed", "other", "none"}:
                    behaviour["show"] = show
                elif show:
                    technical_incomplete = True
                if behaviour:
                    metadata["hyperlink_behaviour"] = behaviour
                continue
            if element.tag == f"{{{_ODF_META_NS}}}user-defined":
                user_defined_total += 1
                name = attribute_text(element.attrib.get(f"{{{_ODF_META_NS}}}name"), 200)
                value_type = attribute_text(
                    element.attrib.get(f"{{{_ODF_META_NS}}}value-type") or "string",
                    32,
                )
                value = element_text(element, 1_000)
                if (
                    len(user_defined) < 32
                    and name
                    and value
                    and re.fullmatch(r"[A-Za-z][A-Za-z0-9._-]{0,31}", value_type)
                ):
                    user_defined.append({"name": name, "value_type": value_type, "value": value})
                else:
                    technical_incomplete = True
                continue
            # ODF extended documents may carry implementation-defined custom
            # metadata.  It is not safe to assign semantics to an unknown XML
            # element, so expose the omission instead of silently dropping it.
            technical_incomplete = True

        if keywords:
            metadata["keywords"] = keywords
        metadata["keywords_total"] = keywords_total
        metadata["keywords_shown"] = len(keywords)
        if user_defined:
            metadata["user_defined"] = user_defined
        metadata["user_defined_total"] = user_defined_total
        metadata["user_defined_shown"] = len(user_defined)
        if keywords_total > len(keywords) or user_defined_total > len(user_defined):
            technical_incomplete = True
        if technical_incomplete:
            metadata["technical_metadata_incomplete"] = True
            metadata["metadata_parse_status"] = "partial"
        metadata = self._redact_metadata_value(metadata)
        own_date = _plausible_document_date(
            str(metadata.get("creation_date") or metadata.get("modified_date") or "")
        )
        if own_date:
            metadata["document_date"] = own_date
        return metadata

    def extract_visual_assets(
        self,
        content: bytes,
        filename: str,
        mime_type: str = "",
        *,
        max_images: int = 4,
        max_pixels: int = 8_000_000,
        max_encoded_bytes: int = 3 * 1024 * 1024,
    ) -> list[VisualAsset]:
        """Extract a few representative images for local vision/OCR.

        The method never renders arbitrary PDF programs or invokes external
        binaries. Direct images are normalized with Pillow; PDF and Office
        documents contribute only already embedded raster images. Every asset
        is downscaled, re-encoded, and bounded before it leaves this module.
        """
        if not isinstance(content, bytes) or len(content) > self.max_input_bytes:
            return []
        max_images = max(1, min(int(max_images), 4))
        safe_name = Path(str(filename or "document")).name
        ext = self._compound_extension(safe_name.casefold())
        detected_mime = (
            (mime_type or mimetypes.guess_type(safe_name)[0] or "").split(";", 1)[0].strip().casefold()
        )
        suffix_known = ext in _KNOWN_DOCUMENT_EXTENSIONS
        candidates: list[tuple[bytes, str]] = []

        if ext in _IMAGE_EXTENSIONS or (
            not suffix_known
            and detected_mime.startswith("image/")
            and detected_mime not in _CONVERTED_OFFICE_MIME_FORMATS
        ):
            candidates.append((content, safe_name))
        elif ext == ".pdf" or (not suffix_known and detected_mime == "application/pdf"):
            candidates.extend(self._pdf_embedded_images(content, max_candidates=max_images * 3))
        elif ext in _OOXML_EXTENSIONS or ext in _OPENDOCUMENT_EXTENSIONS:
            candidates.extend(self._office_embedded_images(content, max_candidates=max_images * 3))

        assets: list[VisualAsset] = []
        seen: set[str] = set()
        for raw, source in candidates:
            if len(raw) > 16 * 1024 * 1024:
                continue
            normalized = self._normalize_visual_asset(
                raw,
                source=source,
                max_pixels=max_pixels,
                max_encoded_bytes=max_encoded_bytes,
            )
            if normalized is None:
                continue
            fingerprint = f"{len(normalized.data)}:{normalized.data[:64]!r}"
            if fingerprint in seen:
                continue
            seen.add(fingerprint)
            assets.append(normalized)
            if len(assets) >= max_images:
                break
        return assets

    @staticmethod
    def local_ocr_available() -> bool:
        """Whether the optional bounded local OCR backend is installed."""

        return local_ocr_available()

    def ocr_visual_assets(
        self,
        assets: Sequence[VisualAsset],
        *,
        deadline: float | None = None,
        executable: str | None = None,
    ) -> LocalOcrResult:
        """Run deterministic local OCR over an already-normalized prefix."""

        return extract_local_ocr(
            assets,
            max_text_chars=self.max_text_chars,
            deadline=deadline,
            executable=executable,
        )

    def render_pdf_pages(
        self,
        content: bytes,
        filename: str,
        mime_type: str = "",
        *,
        max_pages: int = 40,
        max_pixels: int = 8_000_000,
        max_encoded_bytes: int = 1_500_000,
        deadline: float | None = None,
    ) -> VisualPageRender:
        """Render a contiguous, bounded prefix of a PDF entirely in-process.

        ``pypdf`` can recover a text layer and individual embedded images, but
        neither is the visible page: forms, rotations, masks and several image
        tiles may compose one scan page.  PDFium supplies that missing page
        boundary without a shell, temporary files, or an unbounded output
        directory.  Each page is normalized through the same image gate as a
        direct upload before it can reach vision.

        The prefix rule is deliberate.  Skipping a broken middle page and then
        reporting only a count would make ``6 of 7`` look like pages 1--6 even
        when page 3 was the one nobody read.  On the first render failure or
        deadline, stop and return the exact contiguous coverage achieved.
        """

        safe_name = Path(str(filename or "document.pdf")).name
        ext = self._compound_extension(safe_name.casefold())
        detected_mime = (
            (mime_type or mimetypes.guess_type(safe_name)[0] or "").split(";", 1)[0].strip().casefold()
        )
        suffix_known = ext in _KNOWN_DOCUMENT_EXTENSIONS
        if ext != ".pdf" and (suffix_known or detected_mime != "application/pdf"):
            return VisualPageRender((), 0, 0, False, False, False, "not_pdf")
        if not isinstance(content, bytes) or len(content) > self.max_input_bytes:
            return VisualPageRender((), 0, 0, False, False, False, "invalid_pdf_input")

        max_pages = max(1, min(int(max_pages), 128))
        max_pixels = max(100_000, min(int(max_pixels), 12_000_000))
        max_encoded_bytes = max(64_000, min(int(max_encoded_bytes), 3 * 1024 * 1024))
        try:
            import pypdfium2 as pdfium  # type: ignore[import-untyped]
        except ImportError:
            return VisualPageRender((), 0, 0, False, False, False, "pdf_renderer_unavailable")

        assets: list[VisualAsset] = []
        pages_total = 0
        deadline_reached = False
        error = ""
        try:
            document = pdfium.PdfDocument(content)
        except Exception as exc:
            LOGGER.debug("PDF page renderer could not open document (%s)", type(exc).__name__)
            return VisualPageRender((), 0, 0, False, False, False, "pdf_open_failed")
        try:
            pages_total = max(0, int(len(document)))
            for page_index in range(min(pages_total, max_pages)):
                if deadline is not None and time.monotonic() >= deadline:
                    deadline_reached = True
                    break
                page = None
                bitmap = None
                try:
                    page = document[page_index]
                    width, height = page.get_size()
                    page_area = float(width) * float(height)
                    if (
                        not math.isfinite(page_area)
                        or not math.isfinite(float(width))
                        or not math.isfinite(float(height))
                        or width <= 0
                        or height <= 0
                    ):
                        error = "pdf_page_has_invalid_dimensions"
                        break
                    # 2.5x is sharp enough for ordinary 72-DPI PDF coordinates;
                    # oversized pages are scaled down before PDFium allocates a
                    # bitmap, so the configured pixel ceiling is real, not a
                    # post-render resize.
                    max_axis = min(_MAX_PDF_RENDER_AXIS_PIXELS, max_pixels)
                    scale = min(
                        2.5,
                        math.sqrt(max_pixels / page_area),
                        max_axis / float(width),
                        max_axis / float(height),
                    )
                    if not math.isfinite(scale) or scale <= 0:
                        error = "pdf_page_has_invalid_scale"
                        break
                    # PDFium rounds both axes to whole pixels before allocating
                    # the bitmap.  An area-only floating-point check is not
                    # sufficient for an extreme MediaBox such as 80M x 0.1:
                    # ceil(width*scale) can otherwise request a huge row even
                    # though width*height looks harmless.  Find the largest
                    # uniform scale whose actual rounded allocation obeys both
                    # the area and per-axis ceilings.
                    render_width = max(1, math.ceil(float(width) * scale))
                    render_height = max(1, math.ceil(float(height) * scale))
                    if (
                        render_width > max_axis
                        or render_height > max_axis
                        or render_width * render_height > max_pixels
                    ):
                        low = 0.0
                        high = scale
                        for _ in range(48):
                            candidate = (low + high) / 2.0
                            candidate_width = max(1, math.ceil(float(width) * candidate))
                            candidate_height = max(1, math.ceil(float(height) * candidate))
                            if (
                                candidate_width <= max_axis
                                and candidate_height <= max_axis
                                and candidate_width * candidate_height <= max_pixels
                            ):
                                low = candidate
                            else:
                                high = candidate
                        scale = low
                        render_width = max(1, math.ceil(float(width) * scale))
                        render_height = max(1, math.ceil(float(height) * scale))
                    if (
                        not math.isfinite(scale)
                        or scale <= 0
                        or render_width > max_axis
                        or render_height > max_axis
                        or render_width * render_height > max_pixels
                    ):
                        error = "pdf_page_render_budget_exceeded"
                        break
                    bitmap = page.render(scale=scale)
                    image = bitmap.to_pil()
                    encoded = io.BytesIO()
                    image.save(encoded, format="JPEG", quality=90, optimize=False, progressive=False)
                    normalized = self._normalize_visual_asset(
                        encoded.getvalue(),
                        source=f"pdf-page-{page_index + 1}-render",
                        max_pixels=max_pixels,
                        max_encoded_bytes=max_encoded_bytes,
                    )
                    if normalized is None:
                        error = "pdf_page_normalization_failed"
                        break
                    assets.append(normalized)
                except Exception as exc:
                    LOGGER.debug(
                        "PDF page %s rendering failed (%s)",
                        page_index + 1,
                        type(exc).__name__,
                    )
                    error = "pdf_page_render_failed"
                    break
                finally:
                    if bitmap is not None:
                        with suppress(Exception):
                            bitmap.close()
                    if page is not None:
                        with suppress(Exception):
                            page.close()
        finally:
            with suppress(Exception):
                document.close()

        pages_rendered = len(assets)
        page_cap_reached = pages_total > max_pages and pages_rendered >= max_pages
        pages_truncated = pages_rendered < pages_total
        return VisualPageRender(
            tuple(assets),
            pages_total,
            pages_rendered,
            pages_truncated,
            deadline_reached,
            page_cap_reached,
            error,
        )

    def visual_source_pages(self, content: bytes, filename: str, mime_type: str = "") -> int:
        """Сколько страниц (или вложенных картинок) есть у документа ВСЕГО.

        `extract_visual_assets` берёт из них лишь несколько — распознавание идёт
        через модель, и каждая страница стоит места в запросе. Само по себе это
        честная цена; молчание о ней — нет: скан на сорок страниц читался по
        четырём картинкам, и человек получал ответ, уверенный, что прочитано всё.

        Число берётся отдельным дешёвым проходом (у PDF это оглавление, у офисных
        — список членов архива), а не выводится из числа взятых картинок: у одной
        страницы их может быть несколько, и «взято 4» не значит «страниц 4».
        """
        if not isinstance(content, bytes) or len(content) > self.max_input_bytes:
            return 0
        safe_name = Path(str(filename or "document")).name
        ext = self._compound_extension(safe_name.casefold())
        detected_mime = (
            (mime_type or mimetypes.guess_type(safe_name)[0] or "").split(";", 1)[0].strip().casefold()
        )
        suffix_known = ext in _KNOWN_DOCUMENT_EXTENSIONS
        if ext in _IMAGE_EXTENSIONS or (
            not suffix_known
            and detected_mime.startswith("image/")
            and detected_mime not in _CONVERTED_OFFICE_MIME_FORMATS
        ):
            return 1
        if ext == ".pdf" or (not suffix_known and detected_mime == "application/pdf"):
            with suppress(Exception):
                from pypdf import PdfReader

                reader = PdfReader(io.BytesIO(content), strict=False)
                if reader.is_encrypted:
                    with suppress(Exception):
                        if reader.decrypt("") == 0:
                            return 0
                return len(reader.pages)
            return 0
        if ext in _OOXML_EXTENSIONS or ext in _OPENDOCUMENT_EXTENSIONS:
            with suppress(Exception), zipfile.ZipFile(io.BytesIO(content)) as archive:
                return sum(
                    1
                    for info in archive.infolist()
                    if not info.is_dir()
                    and info.filename.casefold().startswith(
                        ("word/media/", "ppt/media/", "xl/media/", "pictures/")
                    )
                    and Path(info.filename).suffix.casefold() in _IMAGE_EXTENSIONS
                )
            return 0
        return 0

    @staticmethod
    def _normalize_visual_asset(
        content: bytes,
        *,
        source: str,
        max_pixels: int,
        max_encoded_bytes: int,
    ) -> VisualAsset | None:
        try:
            from PIL import Image, ImageOps
        except ImportError:
            return None
        try:
            with Image.open(io.BytesIO(content)) as opened_image:
                width, height = opened_image.size
                if width <= 0 or height <= 0 or width * height > 40_000_000:
                    return None
                opened_image.load()
                image = ImageOps.exif_transpose(opened_image)
                if image.width * image.height > max_pixels:
                    scale = (max_pixels / float(image.width * image.height)) ** 0.5
                    target = (max(1, int(image.width * scale)), max(1, int(image.height * scale)))
                    image.thumbnail(target)
                if image.mode not in {"RGB", "L"}:
                    background = Image.new("RGB", image.size, "white")
                    if "A" in image.getbands():
                        background.paste(image, mask=image.getchannel("A"))
                    else:
                        background.paste(image.convert("RGB"))
                    image = background
                elif image.mode == "L":
                    image = image.convert("RGB")

                for quality in (86, 76, 66):
                    output = io.BytesIO()
                    image.save(output, format="JPEG", quality=quality, optimize=False, progressive=False)
                    encoded = output.getvalue()
                    if len(encoded) <= max_encoded_bytes:
                        return VisualAsset(
                            data=encoded,
                            mime_type="image/jpeg",
                            source=source[:200],
                            width=image.width,
                            height=image.height,
                        )
                return None
        except Exception as exc:
            LOGGER.debug("Visual asset normalization failed (%s)", type(exc).__name__)
            return None

    @staticmethod
    def _pdf_embedded_images(content: bytes, *, max_candidates: int) -> list[tuple[bytes, str]]:
        try:
            from pypdf import PdfReader
        except ImportError:
            return []
        output: list[tuple[bytes, str]] = []
        try:
            reader = PdfReader(io.BytesIO(content), strict=False)
            if reader.is_encrypted:
                try:
                    if reader.decrypt("") == 0:
                        return []
                except Exception:
                    return []
            page_count = min(len(reader.pages), 250)
            if page_count <= 0:
                return []
            # Sample across the document instead of over-representing page one.
            indices = sorted(
                {
                    min(page_count - 1, round(position * (page_count - 1) / max(1, max_candidates - 1)))
                    for position in range(max_candidates)
                }
            )
            for page_index in indices:
                try:
                    images = reader.pages[page_index].images
                except Exception:
                    continue
                for image_index, image in enumerate(images):
                    data = bytes(image.data)
                    if 0 < len(data) <= 16 * 1024 * 1024:
                        output.append((data, f"pdf-page-{page_index + 1}-image-{image_index + 1}"))
                    if len(output) >= max_candidates:
                        return output
        except Exception as exc:
            LOGGER.debug("PDF embedded-image extraction failed (%s)", type(exc).__name__)
        return output

    @staticmethod
    def _office_embedded_images(content: bytes, *, max_candidates: int) -> list[tuple[bytes, str]]:
        output: list[tuple[bytes, str]] = []
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as archive:
                members = [
                    info
                    for info in archive.infolist()
                    if not info.is_dir()
                    and info.filename.casefold().startswith(
                        ("word/media/", "ppt/media/", "xl/media/", "pictures/")
                    )
                    and Path(info.filename).suffix.casefold() in _IMAGE_EXTENSIONS
                ]
                members.sort(key=lambda info: info.file_size, reverse=True)
                for info in members[: max_candidates * 2]:
                    if info.file_size <= 0 or info.file_size > 16 * 1024 * 1024:
                        continue
                    if info.compress_size and info.file_size / max(1, info.compress_size) > _MAX_ZIP_RATIO:
                        continue
                    try:
                        data = archive.read(info, pwd=None)
                    except Exception as exc:
                        # A single corrupted media stream (bit flip, interrupted save,
                        # adversarial upload) must not abort the whole document's
                        # extraction — zipfile.read() can raise things well outside
                        # the outer (OSError, BadZipFile, RuntimeError) tuple, e.g.
                        # zlib.error on a broken deflate stream. Skip this image and
                        # keep the others; found by adversarial review.
                        LOGGER.debug("Office embedded image is corrupted (%s)", type(exc).__name__)
                        continue
                    if len(data) != info.file_size:
                        continue
                    output.append((data, info.filename))
                    if len(output) >= max_candidates:
                        break
        except (OSError, zipfile.BadZipFile, RuntimeError) as exc:
            LOGGER.debug("Office embedded-image extraction failed (%s)", type(exc).__name__)
        return output

    @staticmethod
    def _compound_extension(filename: str) -> str:
        for suffix in (".tar.gz", ".tar.bz2", ".tar.xz", ".tar.zst"):
            if filename.endswith(suffix):
                return suffix
        return Path(filename).suffix.casefold()

    @staticmethod
    def _decode(content: bytes) -> str:
        if content.startswith((b"\xff\xfe", b"\xfe\xff")):
            return content.decode("utf-16", errors="replace")
        if content.startswith(b"\xef\xbb\xbf"):
            return content.decode("utf-8-sig", errors="replace")
        for encoding in ("utf-8", "cp1251", "latin-1"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="replace")

    def _bounded_text_source(self, content: bytes, *, structured: bool = False) -> tuple[str, bool]:
        hard_limit = _MAX_STRUCTURED_PARSE_BYTES if structured else _MAX_TEXT_PARSE_BYTES
        # Enough bytes for the configured character budget without letting a
        # 50 MiB log/HTML/JSON file amplify through Python parser structures.
        derived_limit = max(64 * 1024, self.max_text_chars * (2 if structured else 4) + 4)
        parse_limit = min(self.max_input_bytes, hard_limit, derived_limit)
        prefix = content[:parse_limit]
        return self._decode(prefix).replace("\x00", ""), len(prefix) < len(content)

    def _append_bounded(
        self,
        parts: list[str],
        value: str,
        used: int,
        *,
        separator: str = "\n",
    ) -> tuple[int, bool]:
        """Append text without letting the assembled result exceed its budget."""

        if not value:
            return used, False
        separator_size = len(separator) if parts else 0
        available = self.max_text_chars - used - separator_size
        if available <= 0:
            return used, True
        clipped = value[:available]
        parts.append(clipped)
        return used + separator_size + len(clipped), len(clipped) != len(value)

    def _extract_text(self, content: bytes, ext: str) -> DocumentResult:
        structured = ext in {".json", ".csv", ".tsv", ".xml"}
        text, source_truncated = self._bounded_text_source(content, structured=structured)
        metadata: dict[str, Any] = {"format": ext.lstrip(".") or "text"}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        if ext == ".json":
            if source_truncated or len(content) > _MAX_STRUCTURED_PARSE_BYTES:
                metadata["json_valid"] = None
                metadata["json_pretty_skipped"] = True
            else:
                try:
                    text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
                    metadata["json_valid"] = True
                except json.JSONDecodeError:
                    metadata["json_valid"] = False
        elif ext in {".csv", ".tsv"}:
            delimiter = self._table_delimiter(text, default="\t" if ext == ".tsv" else ",")
            metadata["delimiter"] = delimiter
            try:
                rendered_rows: list[str] = []
                parsed_rows: list[list[str]] = []
                rendered_chars = 0
                rows_read = 0
                for row in csv.reader(io.StringIO(text), delimiter=delimiter):
                    if rows_read >= _MAX_TABULAR_ROWS:
                        metadata["rows_truncated"] = True
                        break
                    rows_read += 1
                    parsed_rows.append([cell.strip() for cell in row])
                    rendered = " | ".join(cell.strip() for cell in row)
                    remaining = self.max_text_chars - rendered_chars
                    if remaining <= 0:
                        metadata["rows_truncated"] = True
                        break
                    rendered_rows.append(rendered[:remaining])
                    rendered_chars += min(len(rendered), remaining) + 1
                    if len(rendered) > remaining or rendered_chars >= self.max_text_chars:
                        metadata["rows_truncated"] = True
                        break
                text = "\n".join(rendered_rows)
                metadata["rows_read"] = rows_read
                # Структура таблицы — тем же построителем, что у `.xlsx`. Только
                # если разбор дошёл до конца: индекс, построенный по обрезанной
                # таблице, утверждал бы о ней то, чего в ней нет.
                if not metadata.get("rows_truncated"):
                    text, extra, structure = self._csv_structure(parsed_rows, text)
                    metadata.update(extra)
                    if structure is not None:
                        return DocumentResult(text, metadata, office_structure_index=structure)
            except csv.Error:
                pass
        elif ext == ".xml":
            text = self._strip_xml_tags(text)
        return DocumentResult(text, metadata)

    @staticmethod
    def _table_delimiter(text: str, *, default: str) -> str:
        """Каким знаком разделены колонки — по самому тексту, а не по расширению.

        Русский Excel сохраняет CSV с ТОЧКОЙ С ЗАПЯТОЙ: это его локальное
        умолчание, и такой файл — обычный рабочий документ, а не экзотика. Пока
        разделитель был зашит запятой, каждая строка такого файла становилась
        ОДНОЙ ячейкой: таблица уходила модели плоским текстом, а вид
        «ячейка | ячейка», на который рассчитан и точный путь по таблицам, не
        появлялся вовсе.

        Выбор — не `csv.Sniffer`: тот эвристичен, бросает исключение на коротких
        файлах и на однострочных заголовках, и объяснить его выбор человеку
        нечем. Здесь правило простое и проверяемое: побеждает знак, который даёт
        одинаковое число колонок (больше одной) в наибольшем числе первых строк.
        Ничья — за умолчанием расширения.
        """
        lines = [line for line in text.splitlines()[:20] if line.strip()]
        if not lines:
            return default
        best = default
        best_score = 0
        for candidate in (",", ";", "\t", "|"):
            counts = [line.count(candidate) for line in lines]
            positive = [value for value in counts if value > 0]
            if not positive:
                continue
            # Сколько колонок — решает САМОЕ ЧАСТОЕ значение, а не первая строка.
            # По первой строке считать нельзя: над шапкой часто стоит название
            # отчёта или пустая строка, и один такой заголовок уводил выбор.
            columns = max(set(positive), key=positive.count)
            agreeing = sum(1 for value in counts if value == columns)
            # Две согласные строки — минимум: одна строка согласна сама с собой
            # всегда, и это не признак таблицы, а признак одной запятой в прозе.
            if agreeing < 2:
                continue
            score = agreeing * (columns + 1)
            if score > best_score:
                best_score, best = score, candidate
        return best

    def _extract_html(self, content: bytes) -> DocumentResult:
        source, source_truncated = self._bounded_text_source(content, structured=True)
        soup = BeautifulSoup(source, "lxml")
        title = soup.title.get_text(" ", strip=True) if soup.title else ""
        for tag in soup(["script", "style", "svg", "noscript", "nav", "footer", "header"]):
            tag.decompose()
        root = soup.find("main") or soup.find("article") or soup.body or soup
        lines = [" ".join(line.split()) for line in root.get_text("\n").splitlines()]
        metadata: dict[str, Any] = {"format": "html", "title": html.unescape(title)}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        return DocumentResult(
            "\n".join(line for line in lines if line),
            metadata,
        )

    @staticmethod
    def _html_to_text(source: str) -> str:
        """Видимый текст из размеченного куска — тем же способом, что у страницы.

        Отдельная функция, а не повторное `_extract_html`: у письма и у главы книги
        нет ни заголовка страницы, ни `<main>`, и разметка приходит уже строкой.
        """
        soup = BeautifulSoup(source, "lxml")
        for tag in soup(["script", "style", "svg", "noscript"]):
            tag.decompose()
        lines = [" ".join(line.split()) for line in soup.get_text("\n").splitlines()]
        return "\n".join(line for line in lines if line)

    def _validate_zip(
        self,
        archive: zipfile.ZipFile,
        *,
        total_limit: int | None = None,
        member_limit: int | None = None,
        allow_encrypted: bool = False,
    ) -> list[zipfile.ZipInfo]:
        members = archive.infolist()
        if len(members) > self.max_archive_entries:
            raise ArchiveLimitError(
                f"Archive entry count is {len(members)}; limit is {self.max_archive_entries}"
            )
        expanded_limit = self.max_archive_uncompressed_bytes if total_limit is None else total_limit
        total = 0
        for member in members:
            _safe_archive_member_name(member.filename)
            unix_mode = (int(member.external_attr) >> 16) & 0xFFFF
            if stat.S_ISLNK(unix_mode):
                raise ArchiveLimitError("Archive links are not supported")
            if member.flag_bits & 0x1 and not allow_encrypted:
                raise ArchiveLimitError("Encrypted ZIP entries are not supported")
            if member.is_dir():
                continue
            member_size = max(0, int(member.file_size))
            if member_limit is not None and member_size > member_limit:
                raise ArchiveLimitError("Archive member size exceeds configured limit")
            total += member_size
            if total > expanded_limit:
                raise ArchiveLimitError("Archive uncompressed size exceeds configured limit")
            compressed = max(1, int(member.compress_size))
            if member_size > 1024 * 1024 and member_size / compressed > _MAX_ZIP_RATIO:
                raise ArchiveLimitError("Suspicious ZIP compression ratio")
        return members

    def _validate_office_zip(self, archive: zipfile.ZipFile) -> list[zipfile.ZipInfo]:
        return self._validate_zip(
            archive,
            total_limit=min(self.max_archive_uncompressed_bytes, _MAX_OFFICE_EXPANDED_BYTES),
            member_limit=min(self.max_archive_uncompressed_bytes, _MAX_OFFICE_MEMBER_BYTES),
        )

    def _normalized_ooxml_main_type(
        self,
        content: bytes,
        *,
        main_part: str,
        canonical_type: str,
        alias_types: frozenset[str],
    ) -> tuple[bytes, bool]:
        """Normalize one allowlisted OOXML main type for strict Python readers.

        DOCX/PPTX libraries reject valid template, slideshow and macro-enabled
        containers before reading their otherwise identical main XML.  Rewrite
        only the exact main-part override, after validating the whole ZIP.  No
        other relationship or content-type declaration is inferred or changed.
        """

        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = self._validate_office_zip(archive)
            type_members = [member for member in members if member.filename == _CONTENT_TYPES_MEMBER]
            if len(type_members) != 1:
                raise ValueError("OOXML content types member must be unique")
            type_member = type_members[0]
            if type_member.file_size > _MAX_OOXML_CONTENT_TYPES_BYTES:
                raise ArchiveLimitError("OOXML content types member exceeds limit")
            with archive.open(type_member) as source:
                raw_types, truncated = self._read_stream_preview(
                    source,
                    _MAX_OOXML_CONTENT_TYPES_BYTES,
                )
            if truncated:
                raise ArchiveLimitError("OOXML content types member exceeds limit")
            root = self._metadata_xml_root(raw_types)
            overrides = [
                node
                for node in root.iter()
                if self._metadata_local_name(getattr(node, "tag", "")) == "Override"
                and str(node.get("PartName") or "") == main_part
            ]
            if len(overrides) != 1:
                raise ValueError("OOXML main content type must be unique")
            override = overrides[0]
            content_type = str(override.get("ContentType") or "")
            if content_type == canonical_type:
                return content, False
            if content_type not in alias_types:
                return content, False

            from lxml import etree  # type: ignore[import-untyped]

            override.set("ContentType", canonical_type)
            normalized_types = etree.tostring(
                root,
                encoding="UTF-8",
                xml_declaration=True,
            )
            if len(normalized_types) > _MAX_OOXML_CONTENT_TYPES_BYTES:
                raise ArchiveLimitError("normalized OOXML content types member exceeds limit")
            output = io.BytesIO()
            with zipfile.ZipFile(output, "w") as normalized:
                normalized.comment = archive.comment
                for member in members:
                    payload = (
                        normalized_types if member.filename == _CONTENT_TYPES_MEMBER else archive.read(member)
                    )
                    normalized.writestr(member, payload)
        normalized_content = output.getvalue()
        with zipfile.ZipFile(io.BytesIO(normalized_content)) as archive:
            self._validate_office_zip(archive)
        return normalized_content, True

    def _archive_member_limit(self, budget: _ArchiveBudget) -> int:
        """Bytes one nested member may consume at this point in the upload.

        The old fixed 128 KiB preview ceiling made ordinary nested PDF and DOCX
        files unreadable.  The upload already has stronger shared ceilings: an
        individual document cannot exceed ``max_input_bytes`` and every nesting
        level spends from the same expansion budget.  Use those live limits
        instead of a second, lossy fixed cap.
        """

        return max(0, min(self.max_input_bytes, budget.expanded_bytes))

    @staticmethod
    def _read_stream_limited(
        stream: _BinaryReadable,
        limit: int,
        *,
        deadline: float | None = None,
    ) -> bytes:
        chunks: list[bytes] = []
        total = 0
        while True:
            if deadline is not None and time.monotonic() >= deadline:
                raise ArchiveExtractionError("Archive member read exceeded its deadline")
            chunk = stream.read(min(64 * 1024, limit + 1 - total))
            if not chunk:
                break
            chunks.append(chunk)
            total += len(chunk)
            if total > limit:
                raise ArchiveLimitError("Decompressed member exceeds configured limit")
        return b"".join(chunks)

    @staticmethod
    def _read_stream_preview(stream: _BinaryReadable, limit: int) -> tuple[bytes, bool]:
        chunks: list[bytes] = []
        total = 0
        while total < limit:
            chunk = stream.read(min(64 * 1024, limit - total))
            if not chunk:
                return b"".join(chunks), False
            chunks.append(chunk)
            total += len(chunk)
        return b"".join(chunks), bool(stream.read(1))

    def _extract_docx(self, content: bytes) -> DocumentResult:
        normalized_content, main_type_normalized = self._normalized_ooxml_main_type(
            content,
            main_part=_WORD_MAIN_PART,
            canonical_type=_WORD_CANONICAL_MAIN_TYPE,
            alias_types=_WORD_ALIAS_MAIN_TYPES,
        )
        try:
            from docx import Document

            document = Document(io.BytesIO(normalized_content))
            text, office_structure_index, extraction_truncated = build_docx_text_and_structure(
                document,
                max_text_chars=self.max_text_chars,
                content=normalized_content,
            )
            metadata: dict[str, Any] = {
                "format": "docx",
                "paragraphs": len(document.paragraphs),
                "tables": len(document.tables),
            }
            if main_type_normalized:
                metadata["main_content_type_normalized"] = True
            if extraction_truncated:
                metadata["extraction_truncated"] = True
                metadata["text_truncated"] = True
            return DocumentResult(
                text,
                metadata,
                office_structure_index=office_structure_index,
            )
        except ImportError:
            return self._extract_xml_zip_text(
                normalized_content,
                "word/document.xml",
                "docx",
            )

    def _extract_doc(self, content: bytes, *, deadline: float | None = None) -> DocumentResult:
        """Legacy Word 97-2003. Parsed from bytes; see `friday.documents._ole`.

        Measured on a real 3.19 GB working folder: 206 files in this format, 130 MB,
        and every one of them arrived as `[File: NAME.doc; …]` — an Inbox item with
        nothing in it to review. 197 of the 206 now read, all recognisably Russian
        text, no replacement characters and no control characters left behind.

        Several deployed ``.doc`` files are actually RTF or OOXML exported under
        Word's legacy suffix. Their container magic is unambiguous, so route RTF
        directly and allow a real ZIP container through the bounded converter;
        arbitrary non-OLE bytes still fail before LibreOffice can reinterpret
        renamed plain text as a successful Word document.
        """
        from friday.documents._ole import OLE_SIGNATURE, OleError, extract_doc_text

        if content.lstrip().startswith(b"{\\rtf"):
            result = self._extract_rtf(content)
            return DocumentResult(
                result.text,
                {**result.metadata, "declared_format": "doc", "format": "rtf"},
                result.success,
                result.error,
            )
        zip_container = content.startswith(b"PK") and zipfile.is_zipfile(io.BytesIO(content))
        if not content.startswith(OLE_SIGNATURE) and not zip_container:
            # LibreOffice deliberately accepts plain text under a misleading
            # ``.doc`` suffix. That is not format attestation: arbitrary bytes
            # would otherwise be promoted as a successfully read Word file.
            # Legacy Word 97-2003 is OLE2; the only admitted non-OLE variant is
            # the explicit RTF magic above.
            return DocumentResult(
                "",
                {"format": "doc"},
                False,
                "unsupported_legacy_doc",
            )
        try:
            text, metadata = extract_doc_text(content)
        except OleError:
            converted = self._extract_converted_office(content, "doc", deadline=deadline)
            if converted.success:
                return converted
            return DocumentResult(
                "",
                {"format": "doc", "conversion_error": converted.error},
                False,
                "unsupported_legacy_doc",
            )
        return DocumentResult(text, metadata)

    def _extract_converted_office(
        self,
        content: bytes,
        source_format: str,
        *,
        deadline: float | None,
    ) -> DocumentResult:
        """Convert one closed Office import family, then parse its safe target."""

        from friday.documents._office_convert import convert_legacy_office

        converted = convert_legacy_office(
            content,
            source_format,
            deadline=deadline,
            max_output_bytes=_MAX_OFFICE_EXPANDED_BYTES,
        )
        conversion_metadata: dict[str, Any] = {
            "format": source_format,
            "converted_format": converted.target_format,
            "parser": "libreoffice",
        }
        if not converted.success:
            return DocumentResult("", conversion_metadata, False, converted.error)
        if converted.target_format == "odg":
            parsed = self._extract_xml_zip_text(
                converted.content,
                "content.xml",
                "odg",
            )
        elif converted.target_format == "pdf":
            # `.wps` is shared by Writer and Calc import filters.  PDF is the
            # one bounded target both document services can export without
            # guessing a source family from attacker-controlled MIME data.
            parsed = self._extract_pdf(
                converted.content,
                deadline=self._pdf_parse_deadline(deadline),
            )
        else:
            parser = {
                "docx": self._extract_docx,
                "xlsx": self._extract_xlsx,
                "pptx": self._extract_pptx,
            }.get(converted.target_format)
            if parser is None:  # closed mapping invariant
                return DocumentResult(
                    "",
                    conversion_metadata,
                    False,
                    "legacy_office_conversion_unsupported",
                )
            parsed = parser(converted.content)
        return DocumentResult(
            parsed.text,
            {**parsed.metadata, **conversion_metadata},
            parsed.success,
            parsed.error,
            parsed.office_structure_index,
        )

    def _extract_xlsx(self, content: bytes) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            self._validate_office_zip(archive)
        try:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
        except ImportError:
            return DocumentResult("", {"format": "xlsx"}, False, "openpyxl is not installed")
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        # Formula expressions are not source text and never enter the index.  A
        # second read-only view exists solely to distinguish a genuinely empty
        # cell from a formula whose cached result is absent.  Failure of this
        # auxiliary view makes coverage incomplete; it must not change legacy
        # extraction success or text.
        try:
            formula_workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=False)
        except Exception:  # noqa: BLE001 - fail closed in coverage, preserve old parser result
            formula_workbook = None
        sheet_count = len(workbook.sheetnames)
        try:
            text, office_structure_index, extraction_truncated, row_count = build_xlsx_text_and_structure(
                workbook,
                formula_workbook,
                content=content,
                max_text_chars=self.max_text_chars,
                max_rows=_MAX_TABULAR_ROWS,
            )
        finally:
            workbook.close()
            if formula_workbook is not None:
                formula_workbook.close()
        metadata: dict[str, Any] = {
            "format": "xlsx",
            "sheets": sheet_count,
            "rows_read": row_count,
        }
        if extraction_truncated:
            metadata["extraction_truncated"] = True
            metadata["text_truncated"] = True
        return DocumentResult(
            text,
            metadata,
            office_structure_index=office_structure_index,
        )

    #: Сколько строк CSV получают ПОЛНУЮ структуру таблицы.
    #:
    #: Индекс строится тем же построителем, что у `.xlsx`, а тот работает с
    #: настоящей книгой — значит книгу надо собрать и сохранить в память. Это
    #: стоит времени и памяти линейно по строкам, и на выгрузке в сто тысяч строк
    #: цена перестаёт окупаться. Выше потолка CSV остаётся тем, чем был, — ровным
    #: текстом с колонками, — и об этом говорится вслух.
    _CSV_STRUCTURE_MAX_ROWS = 5_000

    def _csv_structure(
        self, rows: list[list[str]], text: str
    ) -> tuple[str, dict[str, Any], dict[str, Any] | None]:
        """Дать CSV ту же структуру таблицы, что есть у `.xlsx`.

        До сих пор структура доезжала до модели только у `.docx` и `.xlsx`: CSV
        оставался текстом, и точный путь по таблице — «сколько всего позиций»,
        «перечисли всех» — для него не работал вовсе. При этом CSV это ровно
        таблица, и терять её структуру только из-за формата файла нечестно.

        Собирается настоящая книга и прогоняется через ТОТ ЖЕ построитель, что и
        `.xlsx`. Не свой облегчённый разбор: второй построитель индекса разошёлся
        бы с первым на первой же правке, а индекс проверяется валидатором, который
        отбрасывает несогласованное МОЛЧА — вместе с точным путём.
        """
        # Строки приходят РАЗОБРАННЫМИ, а не разбираются здесь заново из текста:
        # к этому месту текст уже склеен видом «ячейка | ячейка», и повторный
        # разбор давал ОДНУ колонку на строку. Индекс при этом получался
        # валидным и полным — и совершенно бесполезным: ни одной записи, ни
        # одного кандидата. Поймано сравнением с тем же файлом в `.xlsx`.
        metadata: dict[str, Any] = {}
        if not rows:
            return text, metadata, None
        if len(rows) > self._CSV_STRUCTURE_MAX_ROWS:
            # Молчаливого отказа быть не должно: человек, спросивший «сколько
            # всего», получит обычный ответ модели вместо точного, и знать об
            # этом он обязан заранее.
            metadata["office_structure_skipped"] = "too_many_rows"
            return text, metadata, None
        try:
            from openpyxl import Workbook, load_workbook  # type: ignore[import-untyped]
        except ImportError:
            metadata["office_structure_skipped"] = "openpyxl_missing"
            return text, metadata, None
        try:
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "CSV"
            for row in rows:
                sheet.append([cell.strip() for cell in row])
            buffer = io.BytesIO()
            workbook.save(buffer)
            workbook.close()
            synthetic = buffer.getvalue()
            values = load_workbook(io.BytesIO(synthetic), read_only=True, data_only=True)
            formulas = load_workbook(io.BytesIO(synthetic), read_only=True, data_only=False)
            try:
                built, index, truncated, _ = build_xlsx_text_and_structure(
                    values,
                    formulas,
                    content=synthetic,
                    max_text_chars=self.max_text_chars,
                    max_rows=_MAX_TABULAR_ROWS,
                )
            finally:
                values.close()
                formulas.close()
        except Exception as exc:  # noqa: BLE001 — структура необязательна, текст важнее
            LOGGER.info("CSV structure build failed (%s)", type(exc).__name__)
            metadata["office_structure_skipped"] = f"build_failed:{type(exc).__name__}"
            return text, metadata, None
        if truncated:
            metadata["text_truncated"] = True
        return built, metadata, index

    def _extract_pptx(self, content: bytes) -> DocumentResult:
        normalized_content, main_type_normalized = self._normalized_ooxml_main_type(
            content,
            main_part=_PRESENTATION_MAIN_PART,
            canonical_type=_PRESENTATION_CANONICAL_MAIN_TYPE,
            alias_types=_PRESENTATION_ALIAS_MAIN_TYPES,
        )
        try:
            from pptx import Presentation
        except ImportError:
            return DocumentResult("", {"format": "pptx"}, False, "python-pptx is not installed")
        presentation = Presentation(io.BytesIO(normalized_content))
        slides: list[str] = []
        used = 0
        extraction_truncated = False
        for number, slide in enumerate(presentation.slides, 1):
            text = [
                shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
            ]
            if text:
                used, clipped = self._append_bounded(
                    slides,
                    f"--- Slide {number} ---\n" + "\n".join(text),
                    used,
                    separator="\n\n",
                )
                if clipped:
                    extraction_truncated = True
                    break
        metadata: dict[str, Any] = {"format": "pptx", "slides": len(presentation.slides)}
        if main_type_normalized:
            metadata["main_content_type_normalized"] = True
        if extraction_truncated:
            metadata["extraction_truncated"] = True
        return DocumentResult("\n\n".join(slides), metadata)

    def _extract_rtf(self, content: bytes) -> DocumentResult:
        source, source_truncated = self._bounded_text_source(content)
        metadata: dict[str, Any] = {"format": "rtf"}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        try:
            from striprtf.striprtf import rtf_to_text  # type: ignore[import-not-found]
        except ImportError:
            # A conservative fallback removes control words but keeps visible text.
            source = re.sub(r"\\'[0-9a-fA-F]{2}", " ", source)
            source = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", source)
            source = source.replace("{", " ").replace("}", " ")
            metadata["parser"] = "fallback"
            return DocumentResult(" ".join(source.split()), metadata)
        metadata["parser"] = "striprtf"
        return DocumentResult(rtf_to_text(source), metadata)

    @staticmethod
    def _extract_msg(content: bytes) -> DocumentResult:
        from friday.documents._ole import OleError, extract_msg_text

        try:
            text, metadata = extract_msg_text(content)
        except OleError:
            return DocumentResult("", {"format": "msg"}, False, "unsupported_outlook_msg")
        return DocumentResult(text, metadata)

    def _extract_email(self, content: bytes) -> DocumentResult:
        """Письмо целиком: заголовки, которые человек читает, и тело.

        Почта разбиралась только органом-импортёром почтового ящика. Файл `.eml`,
        присланный в чат или лежащий в папке, отвергался как «формат не
        поддерживается» — при том что это обычный текст с заголовками, и стандартная
        библиотека разбирает его без единой зависимости.

        Берётся `text/plain`-часть, а при её отсутствии — `text/html`, очищенный тем
        же способом, что и обычная веб-страница. Вложения НЕ разворачиваются: у них
        свой путь через архивы, и разворачивать их здесь значило бы обойти
        собственные потолки этого пути.
        """
        from email import policy
        from email.parser import BytesParser

        source, source_truncated = content[:_MAX_TEXT_PARSE_BYTES], len(content) > _MAX_TEXT_PARSE_BYTES
        message = BytesParser(policy=policy.default).parsebytes(source)
        metadata: dict[str, Any] = {"format": "eml"}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        lines: list[str] = []
        for header, label in (
            ("From", "От"),
            ("To", "Кому"),
            ("Cc", "Копия"),
            ("Date", "Дата"),
            ("Subject", "Тема"),
        ):
            value = str(message.get(header) or "").strip()
            if value:
                lines.append(f"{label}: {value}")
        own_date = _plausible_document_date(_email_iso_date(str(message.get("Date") or "")))
        if own_date:
            # Дата письма — его собственная дата, а не день, когда файл попал в
            # архив: то же правило, что у docx и pdf.
            metadata["document_date"] = own_date
        attachments = [
            str(part.get_filename() or "").strip()
            for part in message.walk()
            if part.get_content_disposition() == "attachment"
        ]
        attachments = [name for name in attachments if name]
        if attachments:
            # Названы вслух: их содержимое сюда не разворачивается, и человек
            # должен знать, что в письме было что-то ещё.
            metadata["attachment_names"] = attachments[:20]
            lines.append("Вложения: " + ", ".join(attachments[:20]))
        body = message.get_body(preferencelist=("plain", "html"))
        text = ""
        if body is not None:
            payload = body.get_content()
            text = str(payload)
            if body.get_content_type() == "text/html":
                text = self._html_to_text(text)
                metadata["parser"] = "html-body"
        if lines:
            text = "\n".join(lines) + ("\n\n" + text if text.strip() else "")
        return DocumentResult(text.strip(), metadata)

    def _extract_epub(self, content: bytes) -> DocumentResult:
        """Книга: главы лежат отдельными XHTML внутри обычного zip.

        Порядок глав берётся из имён членов архива, а не из `content.opf`: чтение
        манифеста добавило бы разбор ещё одного XML ради порядка, который у
        подавляющего большинства книг и так задан именами. Расхождение возможно и
        названо здесь, а не спрятано.
        """
        chapters: list[str] = []
        skipped = 0
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = self._validate_office_zip(archive)
            candidates = sorted(
                (
                    info
                    for info in members
                    if Path(info.filename).suffix.casefold() in {".xhtml", ".html", ".htm"}
                ),
                key=lambda info: info.filename,
            )
            budget = min(self.max_archive_uncompressed_bytes, _MAX_OFFICE_MEMBER_BYTES)
            for info in candidates:
                if info.file_size > budget:
                    skipped += 1
                    continue
                budget -= info.file_size
                with archive.open(info) as stream:
                    data, _ = self._read_stream_preview(stream, _MAX_STRUCTURED_PARSE_BYTES)
                chapter = self._html_to_text(self._decode(data))
                if chapter.strip():
                    chapters.append(chapter)
        metadata: dict[str, Any] = {"format": "epub", "chapters_read": len(chapters)}
        if skipped:
            # Молчаливая потеря главы читалась бы как «в книге этого нет».
            metadata["chapters_skipped"] = skipped
        return DocumentResult("\n\n".join(chapters), metadata)

    def _extract_xml_zip_text(self, content: bytes, member_name: str, format_name: str) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = self._validate_office_zip(archive)
            names = {member.filename for member in members}
            if member_name not in names:
                return DocumentResult("", {"format": format_name}, False, f"Missing member: {member_name}")
            info = archive.getinfo(member_name)
            member_limit = min(self.max_archive_uncompressed_bytes, _MAX_OFFICE_MEMBER_BYTES)
            if info.file_size > member_limit:
                raise ArchiveLimitError("Office document XML exceeds configured limit")
            parse_limit = min(member_limit, _MAX_STRUCTURED_PARSE_BYTES)
            with archive.open(info) as stream:
                data, source_truncated = self._read_stream_preview(stream, parse_limit)
        metadata: dict[str, Any] = {"format": format_name}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        return DocumentResult(self._strip_xml_tags(self._decode(data)), metadata)

    def _extract_pdf(self, content: bytes, *, deadline: float | None = None) -> DocumentResult:
        """Постранично, с потолком страниц, потолком знаков и — если задан —
        СРОКОМ.

        Срок нужен потому, что отмена снаружи здесь не работает: `asyncio.timeout`
        вокруг `to_thread` возвращает управление вызывающему, но сам поток
        продолжает молотить страницу столько, сколько ей нужно. Один
        патологический content stream — и поток пула занят навсегда, а пул общий
        со всем остальным, что уходит с event loop. Проверять срок можно только
        МЕЖДУ страницами: прервать `pypdf` внутри одной страницы нечем.
        """
        try:
            from pypdf import PdfReader
        except ImportError:
            return DocumentResult("", {"format": "pdf"}, False, "pypdf is not installed")
        reader = PdfReader(io.BytesIO(content), strict=False)
        if reader.is_encrypted:
            try:
                if reader.decrypt("") == 0:
                    return DocumentResult("", {"format": "pdf"}, False, "Encrypted PDF is not supported")
            except Exception:
                return DocumentResult("", {"format": "pdf"}, False, "Encrypted PDF is not supported")
        pages: list[str] = []
        pages_read = 0
        text_chars = 0
        extraction_truncated = False
        deadline_hit = False
        # Сколько страниц В ДОКУМЕНТЕ, а не сколько мы согласились прочитать.
        # Потолок в 250 страниц срабатывал молча: 251-я и дальше не попадали ни в
        # текст, ни в признаки, и человек, приславший том на 400 страниц, узнавал
        # об этом, только не найдя в нём того, что там есть.
        total_pages = len(reader.pages)
        for page in itertools.islice(reader.pages, 250):
            if deadline is not None and time.monotonic() >= deadline:
                deadline_hit = True
                extraction_truncated = True
                break
            pages_read += 1
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_chars, clipped = self._append_bounded(
                    pages,
                    page_text.strip(),
                    text_chars,
                    separator="\n\n",
                )
                if clipped:
                    extraction_truncated = True
                    break
            if text_chars >= self.max_text_chars:
                extraction_truncated = True
                break
        metadata: dict[str, Any] = {
            "format": "pdf",
            "pages_read": pages_read,
            "page_limit": 250,
            "total_pages": total_pages,
        }
        if total_pages > 250:
            # Отдельный признак, а не общий `extraction_truncated`: причина обрезки
            # человеку важна. «Не уместилось по объёму» лечится вопросом о начале
            # документа, а «страниц больше, чем читаем» означает, что конца тома
            # система не видела вовсе.
            metadata["pages_truncated"] = True
            metadata["extraction_truncated"] = True
        if extraction_truncated:
            metadata["extraction_truncated"] = True
        if deadline_hit:
            # Отдаём то, что успели прочитать, и честно называем причину: пустой
            # ответ был бы неотличим от «в этом PDF нет текста».
            metadata["parse_deadline_reached"] = True
        return DocumentResult(
            "\n\n".join(pages),
            metadata,
        )

    def _extract_archive(
        self,
        content: bytes,
        filename: str,
        ext: str,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None = None,
        password: str | None = None,
    ) -> DocumentResult:
        if ext == ".zip":
            return self._extract_zip(content, depth, budget, deadline, password)
        if ext in {".tar", ".tar.gz", ".tar.bz2", ".tar.xz"}:
            return self._extract_tar(content, ext, depth, budget, deadline, password)
        if ext in {".gz", ".bz2", ".xz", ".zst"}:
            decompressed = self._decompress_single(content, ext, budget, deadline)
            inner_name = filename[: -len(ext)] or "decompressed.txt"
            return self.extract(
                decompressed,
                inner_name,
                archive_password=password,
                _depth=depth + 1,
                _budget=budget,
                _deadline=deadline,
            )
        if ext == ".rar":
            return self._extract_rar(content, depth, budget, deadline, password)
        if ext == ".7z":
            return self._extract_7z(content, depth, budget, deadline, password)
        if ext == ".tar.zst":
            decompressed = self._decompress_single(content, ".zst", budget, deadline)
            return self._extract_tar(decompressed, ".tar", depth, budget, deadline, password)
        return DocumentResult("", {"format": ext.lstrip(".")}, False, "Unsupported archive format")

    def _decompress_single(
        self,
        content: bytes,
        ext: str,
        budget: _ArchiveBudget,
        deadline: float | None = None,
    ) -> bytes:
        limit = min(self.max_archive_uncompressed_bytes, self.max_input_bytes, budget.expanded_bytes)
        if ext == ".gz":
            stream: _BinaryReadable = gzip.GzipFile(fileobj=io.BytesIO(content))
        elif ext == ".bz2":
            stream = bz2.BZ2File(io.BytesIO(content))
        elif ext == ".xz":
            stream = lzma.LZMAFile(io.BytesIO(content))  # noqa: SIM115 - closed via closing() below
        elif ext == ".zst":
            try:
                import zstandard  # type: ignore[import-not-found]
            except ImportError as exc:
                raise ValueError("zstandard is not installed") from exc
            stream = zstandard.ZstdDecompressor().stream_reader(io.BytesIO(content))
        else:
            raise ValueError(f"Unsupported compressor: {ext}")
        with closing(stream):
            data = self._read_stream_limited(stream, limit, deadline=deadline)
        budget.spend_bytes(len(data))
        return data

    def _member_preview(
        self,
        name: str,
        data: bytes,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None = None,
        password: str | None = None,
    ) -> tuple[str, bool]:
        result = self.extract(
            data,
            name,
            archive_password=password,
            _depth=depth + 1,
            _budget=budget,
            _deadline=deadline,
        )
        # A nested encrypted archive is still an encrypted part of this upload.
        # Do not flatten that into the generic "partial archive" bit or the
        # caller would never know that supplying a password can complete it.
        if result.error == "archive_password_required":
            raise ArchivePasswordRequired
        if result.error == "archive_password_invalid":
            raise ArchivePasswordInvalid
        if not result.success:
            return "", False
        metadata = result.metadata or {}
        complete = not any(
            metadata.get(flag)
            for flag in (
                "text_truncated",
                "extraction_truncated",
                "rows_truncated",
                "source_truncated_for_parse",
                "parse_deadline_reached",
                "pages_truncated",
                "archive_budget_exhausted",
            )
        )
        if not result.text:
            return "", complete
        # A member is already bounded by the live per-upload expansion budget
        # and ``max_input_bytes`` before this recursive parse, and the outer
        # extractor applies its global text ceiling.  A second text slice here
        # would buy no safety: it would only hide a readable member's tail.
        return f"\n--- {name} ---\n{result.text}", complete

    def _extract_zip(
        self,
        content: bytes,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None = None,
        password: str | None = None,
    ) -> DocumentResult:
        try:
            return self._extract_zip_members(content, depth, budget, deadline, password)
        except (
            ArchiveBackendUnavailable,
            ArchiveExtractionError,
            ArchiveLimitError,
            ArchivePasswordInvalid,
            ArchivePasswordRequired,
        ):
            raise
        except Exception as exc:
            # ZIP parser/decompressor diagnostics are attacker-controlled input
            # details.  Keep them out of the public result and model context.
            raise ArchiveExtractionError from exc

    def _extract_zip_members(
        self,
        content: bytes,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None,
        password: str | None,
    ) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as catalog, ExitStack() as stack:
            members = self._validate_zip(catalog, allow_encrypted=True)
            encrypted = any(member.flag_bits & 0x1 for member in members if not member.is_dir())
            if encrypted and not password:
                raise ArchivePasswordRequired

            reader: Any = catalog
            read_members: list[Any] = members
            password_bytes: bytes | None = None
            if encrypted:
                try:
                    import pyzipper  # type: ignore[import-untyped]
                except ImportError as exc:  # pragma: no cover - required runtime dependency
                    raise ArchiveBackendUnavailable from exc
                password_bytes = str(password).encode("utf-8")
                reader = stack.enter_context(pyzipper.AESZipFile(io.BytesIO(content), mode="r"))
                reader.setpassword(password_bytes)
                read_members = list(reader.infolist())
                if len(read_members) != len(members) or any(
                    str(left.filename) != str(right.filename)
                    for left, right in zip(members, read_members, strict=True)
                ):
                    raise ArchiveExtractionError

            indexed_files = [
                (index, member, read_members[index])
                for index, member in enumerate(members)
                if not member.is_dir()
            ]
            parts = [
                f"ZIP archive: {len(indexed_files)} files",
                *(member.filename for _index, member, _reader_member in indexed_files[:100]),
            ]

            # Password validity cannot depend on preview eligibility.  Put the
            # smallest encrypted file first and read it through authenticated
            # EOF/CRC.  If it is larger than the live bound, a correct password
            # ends in archive_limit_exceeded while a wrong one normally fails at
            # the encryption verifier; neither can be reported as success.
            validation_index: int | None = None
            ordered_files = indexed_files
            if encrypted:
                encrypted_files = [item for item in indexed_files if item[1].flag_bits & 0x1]
                if not encrypted_files:
                    raise ArchiveExtractionError
                validation = min(encrypted_files, key=lambda item: (max(0, int(item[1].file_size)), item[0]))
                validation_index = validation[0]
                ordered_files = [validation, *(item for item in indexed_files if item[0] != validation_index)]

            previewed = 0
            exhausted = False
            for index, member, reader_member in ordered_files:
                mandatory_validation = validation_index is not None and index == validation_index
                if deadline is not None and time.monotonic() >= deadline:
                    if mandatory_validation:
                        raise ArchiveExtractionError
                    exhausted = True
                    break
                member_limit = self._archive_member_limit(budget)
                if member.file_size > member_limit and not mandatory_validation:
                    exhausted = True
                    continue
                # Count the decompression, not the success. Incrementing only when a
                # preview came back meant a member that yields no text — an inner
                # archive of binary blobs, say — never advanced the cap.
                if not budget.take_preview():
                    if mandatory_validation:
                        raise ArchiveLimitError("Archive preview budget cannot validate password")
                    exhausted = True
                    break
                previewed += 1
                try:
                    with reader.open(reader_member, pwd=password_bytes) as stream:
                        data = self._read_stream_limited(stream, member_limit, deadline=deadline)
                except ArchiveLimitError:
                    raise
                except (RuntimeError, zipfile.BadZipFile) as exc:
                    if member.flag_bits & 0x1:
                        raise ArchivePasswordInvalid from exc
                    raise ArchiveExtractionError from exc
                except ArchiveExtractionError:
                    raise
                except Exception as exc:
                    raise ArchiveExtractionError from exc
                budget.spend_bytes(len(data))
                preview, member_complete = self._member_preview(
                    member.filename,
                    data,
                    depth,
                    budget,
                    deadline,
                    password,
                )
                if not member_complete:
                    exhausted = True
                if preview:
                    parts.append(preview)
        metadata: dict[str, Any] = {
            "format": "zip",
            "files": len(indexed_files),
            "previewed_files": previewed,
            "encrypted": encrypted,
        }
        if exhausted:
            # ТАR это говорил, ZIP молчал — при том что ZIP на входе встречается
            # чаще всех. Человек получал список из тридцати имён и содержимое
            # двадцати четырёх, и ничто не отличало «прочитано всё» от «прочитана
            # часть».
            metadata["archive_budget_exhausted"] = True
        return DocumentResult("\n".join(parts), metadata)

    def _extract_tar(
        self,
        content: bytes,
        ext: str,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None = None,
        password: str | None = None,
    ) -> DocumentResult:
        # Streaming mode avoids materialising an attacker-controlled member
        # list before the configured entry limit can be enforced.
        with tarfile.open(fileobj=io.BytesIO(content), mode="r|*") as archive:
            entry_count = 0
            file_count = 0
            total = 0
            names: list[str] = []
            previews: list[str] = []
            previewed = 0
            exhausted = False
            for member in archive:
                if deadline is not None and time.monotonic() >= deadline:
                    exhausted = True
                    break
                entry_count += 1
                if entry_count > self.max_archive_entries:
                    raise ArchiveLimitError(f"Archive entry count exceeds limit {self.max_archive_entries}")
                _safe_archive_member_name(member.name)
                if member.issym() or member.islnk() or member.isdev():
                    raise ArchiveLimitError("Archive links and device entries are not supported")
                if not member.isfile():
                    continue
                file_count += 1
                member_size = max(0, int(member.size))
                total += member_size
                if total > self.max_archive_uncompressed_bytes:
                    raise ArchiveLimitError("Archive uncompressed size exceeds configured limit")
                member_limit = self._archive_member_limit(budget)
                if member_size > budget.expanded_bytes:
                    # Streaming past it would itself spend more than the shared
                    # expansion budget, so stop instead of silently walking it.
                    exhausted = True
                    break
                budget.spend_bytes(member_size)
                if len(names) < 100:
                    names.append(member.name)
                # Walking the members is itself the cost here: streaming mode has to
                # read past every member it skips. When the upload's allowance is
                # gone, stop walking rather than skip in a loop — that is where the
                # nested bomb's 107 seconds were actually spent.
                if member_size > member_limit:
                    exhausted = True
                    continue
                stream = archive.extractfile(member)
                if stream is None:
                    exhausted = True
                    continue
                if not budget.take_preview():
                    exhausted = True
                    break
                previewed += 1  # decompressions, not successes — see _extract_zip
                with stream:
                    data = self._read_stream_limited(stream, member_limit, deadline=deadline)
                preview, member_complete = self._member_preview(
                    member.name,
                    data,
                    depth,
                    budget,
                    deadline,
                    password,
                )
                if not member_complete:
                    exhausted = True
                if preview:
                    previews.append(preview)
            parts = [f"TAR archive: {file_count} files", *names, *previews]
        metadata: dict[str, Any] = {
            "format": ext.lstrip("."),
            "files": file_count,
            "previewed_files": previewed,
        }
        if exhausted:
            # Said out loud: the listing is partial, and a caller reading `files`
            # as "what the archive holds" would otherwise be quietly wrong.
            metadata["archive_budget_exhausted"] = True
        return DocumentResult("\n".join(parts), metadata)

    @staticmethod
    def _read_rar_member_with_tool(
        source_path: str,
        member_name: str,
        *,
        tool: str,
        password: str | None,
        limit: int,
        deadline: float | None,
    ) -> bytes:
        """Stream one member through official UnRAR without exposing a password.

        Bare ``-p`` is the documented prompt mode and accepts a redirected pipe;
        the secret is therefore stdin-only, never argv or environment. ``p``
        writes member bytes to stdout, so no extracted member touches disk. The
        received RAR itself is supplied as one private 0600 temporary input per
        archive because UnRAR has no memory-buffer CLI.
        """

        effective_deadline = min(
            deadline if deadline is not None else math.inf,
            time.monotonic() + _RAR_MEMBER_TIMEOUT_SEC,
        )
        safe_member = _safe_archive_member_name(member_name)
        password_switch = "-p" if password is not None else "-p-"
        command = [
            tool,
            "p",
            "-inul",
            "-cfg-",
            _RAR_DICTIONARY_SWITCH,
            password_switch,
            "--",
            source_path,
            safe_member,
        ]
        process = subprocess.Popen(  # noqa: S603 - fixed executable and argv contract
            command,
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
            close_fds=True,
            start_new_session=True,
            # Do not expose service credentials, HOME-based configuration, or
            # attacker-influenced locale state to the external decoder.
            env={"LANG": "C.UTF-8", "LC_ALL": "C.UTF-8"},
        )
        try:
            if process.stdin is None or process.stdout is None:  # pragma: no cover - Popen invariant
                raise ArchiveExtractionError
            if password is not None:
                # RAR itself truncates at 127 Unicode characters. Sending
                # anything beyond that can only become input to a second
                # prompt, so apply the same ceiling before the pipe.
                process.stdin.write((str(password)[:127] + "\n").encode("utf-8"))
            process.stdin.close()
            descriptor = process.stdout.fileno()
            os.set_blocking(descriptor, False)
            output = bytearray()
            eof = False
            with selectors.DefaultSelector() as selector:
                selector.register(descriptor, selectors.EVENT_READ)
                while not eof:
                    remaining = effective_deadline - time.monotonic()
                    if remaining <= 0:
                        raise ArchiveExtractionError
                    events = selector.select(min(0.25, remaining))
                    if not events and process.poll() is None:
                        continue
                    while True:
                        try:
                            chunk = os.read(descriptor, min(64 * 1024, limit + 1 - len(output)))
                        except BlockingIOError:
                            break
                        if not chunk:
                            eof = True
                            break
                        output.extend(chunk)
                        if len(output) > limit:
                            raise ArchiveLimitError("Decompressed RAR member exceeds configured limit")
                    if process.poll() is not None and not events and not eof:
                        # One more pass observes either buffered bytes or EOF.
                        continue
            return_code = process.wait(timeout=max(0.1, effective_deadline - time.monotonic()))
            if return_code == 11:
                raise ArchivePasswordInvalid
            if return_code != 0:
                raise ArchiveExtractionError
            return bytes(output)
        except BaseException:
            if process.poll() is None:
                process.kill()
            with suppress(Exception):
                process.wait(timeout=1)
            raise
        finally:
            if process.stdout is not None:
                process.stdout.close()
            if process.stdin is not None:
                process.stdin.close()

    def _extract_rar(
        self,
        content: bytes,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None = None,
        password: str | None = None,
    ) -> DocumentResult:
        try:
            import rarfile  # type: ignore[import-untyped]
        except ImportError as exc:  # pragma: no cover - required runtime dependency
            raise ArchiveBackendUnavailable from exc

        try:
            return self._extract_rar_members(
                rarfile,
                content,
                depth,
                budget,
                deadline,
                password,
            )
        except rarfile.RarWrongPassword as exc:
            raise ArchivePasswordInvalid from exc
        except (
            ArchiveBackendUnavailable,
            ArchiveExtractionError,
            ArchiveLimitError,
            ArchivePasswordInvalid,
            ArchivePasswordRequired,
        ):
            raise
        except Exception as exc:
            raise ArchiveExtractionError from exc

    def _extract_rar_members(
        self,
        rarfile: Any,
        content: bytes,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None,
        password: str | None,
    ) -> DocumentResult:
        with rarfile.RarFile(io.BytesIO(content)) as archive, ExitStack() as stack:
            encrypted = bool(archive.needs_password())
            if encrypted and not password:
                raise ArchivePasswordRequired
            if password:
                archive.setpassword(str(password))
            members = archive.infolist()
            encrypted = encrypted or bool(archive.needs_password())
            if encrypted and not password:
                raise ArchivePasswordRequired
            if len(members) > self.max_archive_entries:
                raise ArchiveLimitError("RAR entry count exceeds configured limit")
            for member in members:
                _safe_archive_member_name(member.filename)
                if member.is_symlink():
                    raise ArchiveLimitError("Archive links are not supported")
            indexed_files = [(index, member) for index, member in enumerate(members) if not member.isdir()]
            total = sum(max(0, int(member.file_size)) for _index, member in indexed_files)
            if total > self.max_archive_uncompressed_bytes:
                raise ArchiveLimitError("RAR uncompressed size exceeds configured limit")
            parts = [
                f"RAR archive: {len(indexed_files)} files",
                *(member.filename for _index, member in indexed_files[:100]),
            ]

            validation_index: int | None = None
            ordered_files = indexed_files
            if encrypted:
                encrypted_files = [item for item in indexed_files if item[1].needs_password()]
                if encrypted_files:
                    validation = min(
                        encrypted_files,
                        key=lambda item: (max(0, int(item[1].file_size)), item[0]),
                    )
                    validation_index = validation[0]
                    ordered_files = [
                        validation,
                        *(item for item in indexed_files if item[0] != validation_index),
                    ]

            source: Any | None = None
            tool = ""

            def rar_source_path() -> str:
                nonlocal source, tool
                if not tool:
                    tool = str(shutil.which("unrar") or "")
                    if not tool:
                        raise ArchiveBackendUnavailable
                if source is None:
                    source = stack.enter_context(
                        tempfile.NamedTemporaryFile(prefix="friday-rar-", suffix=".rar")
                    )
                    os.chmod(source.name, 0o600)
                    source.write(content)
                    source.flush()
                return str(source.name)

            previewed = 0
            exhausted = False
            for index, member in ordered_files:
                mandatory_validation = validation_index is not None and index == validation_index
                if deadline is not None and time.monotonic() >= deadline:
                    if mandatory_validation:
                        raise ArchiveExtractionError
                    exhausted = True
                    break
                member_limit = self._archive_member_limit(budget)
                if member.file_size > member_limit and not mandatory_validation:
                    exhausted = True
                    continue
                if not budget.take_preview():
                    if mandatory_validation:
                        raise ArchiveLimitError("Archive preview budget cannot validate password")
                    exhausted = True
                    break
                previewed += 1  # decompressions, not successes — see _extract_zip
                data = self._read_rar_member_with_tool(
                    rar_source_path(),
                    member.filename,
                    tool=tool,
                    password=password if encrypted else None,
                    limit=member_limit,
                    deadline=deadline,
                )
                budget.spend_bytes(len(data))
                preview, member_complete = self._member_preview(
                    member.filename,
                    data,
                    depth,
                    budget,
                    deadline,
                    password,
                )
                if not member_complete:
                    exhausted = True
                if preview:
                    parts.append(preview)
        rar_metadata: dict[str, Any] = {
            "format": "rar",
            "files": len(indexed_files),
            "previewed_files": previewed,
            "encrypted": encrypted,
        }
        if exhausted:
            rar_metadata["archive_budget_exhausted"] = True
        return DocumentResult("\n".join(parts), rar_metadata)

    def _validate_7z_coder_folders(self, folders: Sequence[Any]) -> None:
        """Reject 7z coder parameters that can demand unbounded RAM or CPU."""

        if len(folders) > self.max_archive_entries:
            raise ArchiveLimitError("7z folder count exceeds configured limit")
        for folder in folders:
            coders = list(getattr(folder, "coders", ()) or ())
            if not coders or len(coders) > 4:
                raise ArchiveExtractionError
            for coder in coders:
                method = bytes(coder.get("method") or b"")
                raw_properties = coder.get("properties")
                if raw_properties is not None and not isinstance(raw_properties, bytes):
                    raise ArchiveExtractionError
                properties = raw_properties or b""
                dictionary_bytes = 0

                if method == b"\x21":  # LZMA2
                    if len(properties) != 1 or properties[0] > 40:
                        raise ArchiveExtractionError
                    value = int(properties[0])
                    dictionary_bytes = 0xFFFFFFFF if value == 40 else (2 | (value & 1)) << (value // 2 + 11)
                elif method == b"\x03\x01\x01":  # LZMA1
                    if len(properties) != 5:
                        raise ArchiveExtractionError
                    dictionary_bytes = int.from_bytes(properties[1:5], "little")
                elif method == b"\x03\x04\x01":  # PPMd
                    if len(properties) not in {5, 7}:
                        raise ArchiveExtractionError
                    dictionary_bytes = int.from_bytes(properties[1:5], "little")
                elif method == b"\x06\xf1\x07\x01":  # 7zAES
                    if len(properties) < 2 or properties[0] & 0xC0 == 0:
                        raise ArchiveExtractionError
                    cycles_power = properties[0] & 0x3F
                    salt_size = ((properties[0] >> 7) & 1) + (properties[1] >> 4)
                    iv_size = ((properties[0] >> 6) & 1) + (properties[1] & 0x0F)
                    if len(properties) != 2 + salt_size + iv_size:
                        raise ArchiveExtractionError
                    if cycles_power != 0x3F and cycles_power > _MAX_7Z_AES_CYCLES_POWER:
                        raise ArchiveLimitError("7z AES work factor exceeds configured limit")

                if dictionary_bytes > _MAX_ARCHIVE_DICTIONARY_BYTES:
                    raise ArchiveLimitError("7z dictionary exceeds configured limit")

    def _preflight_7z_header(self, content: bytes) -> None:
        """Inspect an encoded 7z header before py7zr constructs its decoders.

        Py7zr otherwise instantiates the encoded-header LZMA dictionary while
        opening the archive, too early for the ordinary post-open coder check.
        Its metadata parser is safe to use here: it only describes folders and
        does not construct a decompressor.
        """

        from py7zr.archiveinfo import HeaderStreamsInfo, SignatureHeader
        from py7zr.helpers import calculate_crc32
        from py7zr.properties import MAGIC_7Z, PROPERTY

        source = io.BytesIO(content)
        if source.read(len(MAGIC_7Z)) != MAGIC_7Z:
            raise ArchiveExtractionError
        source.seek(0)
        signature = SignatureHeader.retrieve(source)
        after_header = source.tell()
        next_size = int(signature.nextheadersize)
        next_offset = int(signature.nextheaderofs)
        if next_size < 0 or next_size > _MAX_7Z_HEADER_BYTES or next_offset < 0:
            raise ArchiveLimitError("7z header exceeds configured limit")
        next_position = after_header + next_offset
        if next_position < after_header or next_position + next_size > len(content):
            raise ArchiveExtractionError
        source.seek(next_position)
        header_bytes = source.read(next_size)
        if len(header_bytes) != next_size or calculate_crc32(header_bytes) != signature.nextheadercrc:
            raise ArchiveExtractionError
        if not header_bytes:
            return
        header = io.BytesIO(header_bytes)
        marker = header.read(1)
        if marker == PROPERTY.HEADER:
            return
        if marker != PROPERTY.ENCODED_HEADER:
            raise ArchiveExtractionError
        streams = HeaderStreamsInfo.retrieve(header)
        folders = list(streams.unpackinfo.folders)
        self._validate_7z_coder_folders(folders)
        expanded_header = sum(
            max(0, int(folder.unpacksizes[-1])) for folder in folders if getattr(folder, "unpacksizes", None)
        )
        if expanded_header > _MAX_7Z_HEADER_BYTES:
            raise ArchiveLimitError("Expanded 7z header exceeds configured limit")

    def _validate_open_7z_coders(self, archive: Any) -> None:
        streams = getattr(getattr(archive, "header", None), "main_streams", None)
        unpackinfo = getattr(streams, "unpackinfo", None)
        folders = list(getattr(unpackinfo, "folders", ()) or ())
        if folders:
            self._validate_7z_coder_folders(folders)

    @staticmethod
    def _encrypted_7z_member_names(archive: Any) -> set[str]:
        """Map data-encrypted 7z folders back to their member names.

        Encryption is a folder/packed-stream property in the 7z format, not an
        archive-wide promise.  In particular, an archive may contain an older
        plain folder followed by a newly-added encrypted one.  Py7zr attaches
        the parsed folder object to every ``ArchiveFile`` in
        ``_real_get_contents``; use that mapping so password validation cannot
        accidentally stop at the ordinary 24-preview cap.
        """

        encrypted_names: set[str] = set()
        for member in getattr(archive, "files", ()) or ():
            folder = getattr(member, "folder", None)
            coders = list(getattr(folder, "coders", ()) or ())
            if not any(bytes(coder.get("method") or b"") == b"\x06\xf1\x07\x01" for coder in coders):
                continue
            if bool(getattr(member, "is_directory", False)):
                continue
            encrypted_names.add(_safe_archive_member_name(getattr(member, "filename", "")))
        return encrypted_names

    def _extract_7z(
        self,
        content: bytes,
        depth: int,
        budget: _ArchiveBudget,
        deadline: float | None = None,
        password: str | None = None,
    ) -> DocumentResult:
        try:
            import py7zr  # type: ignore[import-not-found]
        except ImportError as exc:  # pragma: no cover - required runtime dependency
            raise ArchiveBackendUnavailable from exc

        try:
            self._preflight_7z_header(content)
        except (ArchiveExtractionError, ArchiveLimitError):
            raise
        except Exception as exc:
            raise ArchiveExtractionError from exc

        # Probe without a password first.  That distinguishes an encrypted
        # header from a generally corrupt file before a wrong password can make
        # py7zr fail with a codec-level TypeError/LZMAError.
        content_encrypted = False
        try:
            with py7zr.SevenZipFile(io.BytesIO(content), mode="r") as probe:
                self._validate_open_7z_coders(probe)
                encrypted = bool(probe.needs_password())
                content_encrypted = encrypted
                entries = probe.list()
        except py7zr.PasswordRequired:
            encrypted = True
            entries = []
        except (ArchiveExtractionError, ArchiveLimitError):
            raise
        except Exception as exc:
            raise ArchiveExtractionError from exc
        if encrypted and not password:
            raise ArchivePasswordRequired

        try:
            with py7zr.SevenZipFile(
                io.BytesIO(content),
                mode="r",
                password=str(password) if password is not None else None,
                max_extract_size=min(self.max_archive_uncompressed_bytes, budget.expanded_bytes),
            ) as archive:
                self._validate_open_7z_coders(archive)
                entries = archive.list()
                encrypted_member_names = self._encrypted_7z_member_names(archive)
                if content_encrypted and not encrypted_member_names:
                    # The password-free probe observed data encryption, so an
                    # absent member mapping is parser ambiguity, not proof that
                    # the supplied password is valid.  Fail closed.
                    raise ArchiveExtractionError
                if len(entries) > self.max_archive_entries:
                    raise ArchiveLimitError("7z entry count exceeds configured limit")
                names: list[str] = []
                files: list[tuple[Any, str, int]] = []
                seen_names: set[str] = set()
                total = 0
                for entry in entries:
                    name = _safe_archive_member_name(getattr(entry, "filename", ""))
                    if name in seen_names:
                        raise ArchiveLimitError("Duplicate 7z member name")
                    seen_names.add(name)
                    if bool(getattr(entry, "is_symlink", False)):
                        raise ArchiveLimitError("Archive links are not supported")
                    if bool(getattr(entry, "is_directory", False)):
                        continue
                    member_size = max(0, int(getattr(entry, "uncompressed", 0) or 0))
                    total += member_size
                    if total > self.max_archive_uncompressed_bytes:
                        raise ArchiveLimitError("7z uncompressed size exceeds configured limit")
                    names.append(name)
                    files.append((entry, name, member_size))

                validation_item: tuple[Any, str, int] | None = None
                if encrypted_member_names:
                    validation_candidates = [item for item in files if item[1] in encrypted_member_names]
                    if len(validation_candidates) != len(encrypted_member_names):
                        raise ArchiveExtractionError
                    validation_item = min(validation_candidates, key=lambda item: (item[2], item[1]))
                ordered_files = files
                if validation_item is not None:
                    ordered_files = [
                        validation_item,
                        *(item for item in files if item[1] != validation_item[1]),
                    ]

                selected: list[tuple[Any, str, int]] = []
                exhausted = False
                reserved_bytes = 0
                for item in ordered_files:
                    mandatory_validation = validation_item is not None and item[1] == validation_item[1]
                    if deadline is not None and time.monotonic() >= deadline:
                        if mandatory_validation:
                            raise ArchiveExtractionError
                        exhausted = True
                        break
                    member_limit = max(
                        0,
                        min(
                            self.max_input_bytes,
                            budget.expanded_bytes - reserved_bytes,
                        ),
                    )
                    if item[2] > member_limit and not mandatory_validation:
                        exhausted = True
                        continue
                    if not budget.take_preview():
                        if mandatory_validation:
                            raise ArchiveLimitError("Archive preview budget cannot validate password")
                        exhausted = True
                        break
                    selected.append(item)
                    reserved_bytes += item[2]

                # Content encryption is not authenticated by listing metadata.
                # Header encryption was already authenticated while opening the
                # archive.  An archive with only an encrypted header has no
                # encrypted data member to force; an unusual empty encrypted
                # archive still needs a bounded decoder attempt if possible.
                if encrypted and files and not selected:
                    if deadline is not None and time.monotonic() >= deadline:
                        raise ArchiveExtractionError
                    if not budget.take_preview():
                        raise ArchiveLimitError("Archive preview budget cannot validate password")
                    selected = [min(files, key=lambda item: (item[2], item[1]))]
                    exhausted = True

                previews: list[str] = []
                if selected:
                    extraction_limit = max(0, budget.expanded_bytes)
                    factory = _Bounded7zFactory(
                        member_limit=min(self.max_input_bytes, extraction_limit),
                        total_limit=extraction_limit,
                        deadline=deadline,
                    )
                    archive.extract(
                        targets=[name for _entry, name, _size in selected],
                        factory=factory,
                    )
                    for _entry, name, _size in selected:
                        if deadline is not None and time.monotonic() >= deadline:
                            exhausted = True
                            break
                        product = factory.get(name)
                        product.seek(0)
                        data = product.read()
                        budget.spend_bytes(len(data))
                        preview, member_complete = self._member_preview(
                            name,
                            data,
                            depth,
                            budget,
                            deadline,
                            password,
                        )
                        if not member_complete:
                            exhausted = True
                        if preview:
                            previews.append(preview)
        except py7zr.DecompressionBombError as exc:
            raise ArchiveLimitError("7z extraction exceeds configured limit") from exc
        except ArchiveLimitError:
            raise
        except (ArchiveExtractionError, ArchivePasswordRequired, ArchivePasswordInvalid):
            raise
        except Exception as exc:
            if encrypted and password is not None:
                raise ArchivePasswordInvalid from exc
            raise ArchiveExtractionError from exc

        previewed = len(selected)
        metadata: dict[str, Any] = {
            "format": "7z",
            "files": len(files),
            "previewed_files": previewed,
            "encrypted": encrypted,
        }
        if exhausted:
            metadata["archive_budget_exhausted"] = True
        return DocumentResult(
            "\n".join([f"7z archive: {len(files)} files", *names[:100], *previews]),
            metadata,
        )

    @staticmethod
    def _strip_xml_tags(xml_text: str) -> str:
        text = re.sub(r"<[^>]+>", " ", xml_text)
        text = html.unescape(text)
        return " ".join(text.split())
