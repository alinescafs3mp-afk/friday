"""Safe, in-memory document text extraction.

The extractor never writes archive members to disk.  Archive entry counts,
uncompressed byte totals, individual reads, and recursion depth are bounded so
uploaded files cannot expand without limit.  Optional Office/RTF/RAR support is
activated when the corresponding dependency is installed.
"""

from __future__ import annotations

import bz2
import csv
import gzip
import html
import io
import itertools
import json
import logging
import lzma
import mimetypes
import re
import tarfile
import zipfile
from contextlib import closing, suppress
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Protocol

from bs4 import BeautifulSoup

LOGGER = logging.getLogger(__name__)
_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".xml",
    ".log",
    ".py",
    ".ps1",
    ".sh",
    ".ini",
    ".cfg",
    ".conf",
    ".toml",
    ".sql",
    ".js",
    ".ts",
    ".css",
    ".rst",
}
_HTML_EXTENSIONS = {".html", ".htm"}
_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".odt", ".rtf"}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
_ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".bz2", ".xz", ".7z", ".rar", ".zst"}
_MAX_NESTING_DEPTH = 2
_MAX_ARCHIVE_PREVIEW_FILES = 24
_MAX_MEMBER_PREVIEW_BYTES = 128 * 1024
_MAX_STRUCTURED_PARSE_BYTES = 2 * 1024 * 1024
_MAX_TEXT_PARSE_BYTES = 8 * 1024 * 1024
_MAX_OFFICE_EXPANDED_BYTES = 128 * 1024 * 1024
_MAX_OFFICE_MEMBER_BYTES = 64 * 1024 * 1024
_MAX_ZIP_RATIO = 500.0
_MAX_TABULAR_ROWS = 100_000

# Дата САМОГО документа, взятая из провенанса файла, а не угаданная из текста.
#
# Существующий фильтр по датам намеренно означает «документ УПОМИНАЕТ дату»:
# документ называет несколько дат, и какая из них его собственная, текст не
# говорит. Но у файла есть собственная дата, которую никто не выдумывал, — её
# записал редактор при сохранении (docProps/core.xml у docx/xlsx, /CreationDate
# у pdf). На архиве владельца это единственный способ отличить документ 2019
# года от документа 2025-го: у всех 1537 объектов дата загрузки одна и та же —
# день импорта.
#
# Берётся ТОЛЬКО из формата. Файловое mtime сюда не годится: копирование на
# флешку переписывает его у всех файлов разом, и «дата документа» стала бы
# датой копирования — то самое угадывание, от которого фильтр и уходил.
_OFFICE_CORE_PROPERTIES = "docProps/core.xml"
_CORE_DATE_RE = re.compile(
    r"<(?:dcterms:)?(created|modified)[^>]*>([0-9]{4}-[0-9]{2}-[0-9]{2})", re.IGNORECASE
)
# 1900 — ниже этого у офисных файлов лежат только служебные нули и мусор
# конвертеров; 2100 — потолок против «31.12.9999», который ставят генераторы.
_DOCUMENT_DATE_MIN = "1900-01-01"
_DOCUMENT_DATE_MAX = "2100-01-01"


def _plausible_document_date(value: str) -> str | None:
    """ISO-дата, если она вообще похожа на настоящую дату документа."""
    candidate = (value or "").strip()[:10]
    if len(candidate) != 10 or candidate[4] != "-" or candidate[7] != "-":
        return None
    if not (_DOCUMENT_DATE_MIN <= candidate < _DOCUMENT_DATE_MAX):
        return None
    try:
        datetime.strptime(candidate, "%Y-%m-%d")  # noqa: DTZ007 - календарная дата без времени
    except ValueError:
        return None
    return candidate


def _office_document_date(content: bytes) -> str | None:
    """`dcterms:created` из docProps/core.xml; при отсутствии — `modified`."""
    with suppress(Exception):
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            if _OFFICE_CORE_PROPERTIES not in archive.namelist():
                return None
            with archive.open(_OFFICE_CORE_PROPERTIES) as handle:
                # Ограничение чтения: core.xml — это килобайты, и раздутый член
                # архива не должен превращаться в чтение сотен мегабайт.
                raw = handle.read(64 * 1024).decode("utf-8", errors="replace")
        found: dict[str, str] = {}
        for match in _CORE_DATE_RE.finditer(raw):
            found.setdefault(match.group(1).casefold(), match.group(2))
        for key in ("created", "modified"):
            candidate = _plausible_document_date(found.get(key, ""))
            if candidate:
                return candidate
    return None


# Форматы, у которых внутри zip лежит docProps/core.xml.
_OFFICE_DATE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


def _pdf_document_date_from_bytes(content: bytes) -> str | None:
    """Дата PDF по байтам — работает и для скана без текстового слоя."""
    with suppress(Exception):
        from pypdf import PdfReader

        return _pdf_document_date(PdfReader(io.BytesIO(content), strict=False))
    return None


def _pdf_document_date(reader: Any) -> str | None:
    """`/CreationDate` вида `D:20230412...` из метаданных PDF."""
    with suppress(Exception):
        info = reader.metadata or {}
        for key in ("/CreationDate", "/ModDate"):
            raw = str(info.get(key) or "")
            digits = raw[2:10] if raw.startswith("D:") else raw[:8]
            if len(digits) == 8 and digits.isdigit():
                candidate = _plausible_document_date(f"{digits[:4]}-{digits[4:6]}-{digits[6:8]}")
                if candidate:
                    return candidate
    return None


@dataclass(frozen=True)
class DocumentResult:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    success: bool = True
    error: str = ""

    def to_dict(self, *, preview_chars: int = 2_000) -> dict[str, Any]:
        return {
            "text": self.text[:preview_chars],
            "text_length": len(self.text),
            "metadata": self.metadata,
            "success": self.success,
            "error": self.error,
            "truncated": len(self.text) > preview_chars,
        }


@dataclass(frozen=True)
class VisualAsset:
    """Bounded, normalized image supplied to the local vision model."""

    data: bytes
    mime_type: str
    source: str
    width: int
    height: int

    def to_dict(self) -> dict[str, Any]:
        return {
            "mime_type": self.mime_type,
            "source": self.source,
            "width": self.width,
            "height": self.height,
            "size_bytes": len(self.data),
        }


class ArchiveLimitError(ValueError):
    """An archive exceeded configured expansion or entry limits."""


class _BinaryReadable(Protocol):
    """Minimal binary stream contract shared by archive/decompressor readers."""

    def read(self, size: int = -1) -> bytes: ...

    def close(self) -> None: ...


class _ArchiveBudget:
    """What ONE upload may spend on unpacking, across every nesting level.

    The limits used to be per archive: each nested member started again with a
    full allowance of 24 previews and `max_archive_uncompressed_bytes`, so the
    ceilings multiplied instead of dividing. Measured: a 3.0 MB zip of 24
    `.tar.gz` members (440 x 250 KiB inside each) expanded to ~2.7 GB and held
    the event loop for 107 seconds, returning `success=True` and raising nothing
    — the operator's "250 MB per upload" was in fact 250 MB per level.
    """

    __slots__ = ("expanded_bytes", "previews")

    def __init__(self, *, previews: int, expanded_bytes: int) -> None:
        self.previews = max(0, int(previews))
        self.expanded_bytes = max(0, int(expanded_bytes))

    def take_preview(self) -> bool:
        if self.previews <= 0:
            return False
        self.previews -= 1
        return True

    def spend_bytes(self, count: int) -> None:
        self.expanded_bytes = max(0, self.expanded_bytes - max(0, int(count)))

    @property
    def spent_out(self) -> bool:
        return self.previews <= 0 or self.expanded_bytes <= 0


class DocumentExtractor:
    def __init__(
        self,
        *,
        max_archive_entries: int = 500,
        max_archive_uncompressed_bytes: int = 250 * 1024 * 1024,
        max_text_chars: int = 2_000_000,
        max_input_bytes: int = 50 * 1024 * 1024,
    ) -> None:
        self.max_archive_entries = max(1, int(max_archive_entries))
        self.max_archive_uncompressed_bytes = max(1024, int(max_archive_uncompressed_bytes))
        self.max_text_chars = max(10_000, int(max_text_chars))
        self.max_input_bytes = max(1, int(max_input_bytes))

    def extract(
        self,
        content: bytes,
        filename: str,
        mime_type: str = "",
        *,
        _depth: int = 0,
        _budget: _ArchiveBudget | None = None,
    ) -> DocumentResult:
        if not isinstance(content, bytes):
            return DocumentResult("", success=False, error="Document content must be bytes")
        if _depth > _MAX_NESTING_DEPTH:
            return DocumentResult("", success=False, error="Archive nesting limit exceeded")
        # One budget per upload, created at the top and carried down the nesting.
        budget = _budget or _ArchiveBudget(
            previews=_MAX_ARCHIVE_PREVIEW_FILES,
            expanded_bytes=self.max_archive_uncompressed_bytes,
        )
        if len(content) > self.max_input_bytes:
            return DocumentResult(
                "",
                {"input_bytes": len(content), "max_input_bytes": self.max_input_bytes},
                False,
                "Document input exceeds configured size limit",
            )

        safe_name = Path(str(filename or "document")).name
        lowered = safe_name.casefold()
        ext = self._compound_extension(lowered)
        detected_mime = mime_type or mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
        try:
            if ext in _TEXT_EXTENSIONS:
                result = self._extract_text(content, ext)
            elif ext in _HTML_EXTENSIONS or detected_mime in {"text/html", "application/xhtml+xml"}:
                result = self._extract_html(content)
            elif ext == ".pdf" or detected_mime == "application/pdf":
                result = self._extract_pdf(content)
            elif ext == ".docx":
                result = self._extract_docx(content)
            elif ext == ".doc" or detected_mime == "application/msword":
                result = self._extract_doc(content)
            elif ext == ".xlsx":
                result = self._extract_xlsx(content)
            elif ext == ".pptx":
                result = self._extract_pptx(content)
            elif ext == ".odt":
                result = self._extract_xml_zip_text(content, "content.xml", "odt")
            elif ext == ".rtf":
                result = self._extract_rtf(content)
            elif ext in _ARCHIVE_EXTENSIONS or ext.startswith(".tar."):
                result = self._extract_archive(content, safe_name, ext, _depth, budget)
            elif detected_mime.startswith("text/"):
                result = self._extract_text(content, ext or ".txt")
            else:
                result = DocumentResult(
                    "",
                    {"filename": safe_name, "mime_type": detected_mime, "size": len(content)},
                    False,
                    f"Unsupported document format: {ext or detected_mime}",
                )
        except ArchiveLimitError as exc:
            result = DocumentResult("", {"filename": safe_name}, False, str(exc))
        except Exception as exc:  # defensive boundary for optional parsers
            LOGGER.info("Document extraction failed for %s: %s", safe_name, exc)
            result = DocumentResult("", {"filename": safe_name}, False, str(exc))

        text = result.text[: self.max_text_chars]
        metadata = {
            "filename": safe_name,
            "mime_type": detected_mime,
            "input_bytes": len(content),
            **result.metadata,
        }
        if len(result.text) > self.max_text_chars:
            metadata["text_truncated"] = True
            metadata["original_text_chars"] = len(result.text)
        # Собственная дата документа снимается ЗДЕСЬ, а не внутри разборщиков, и
        # потому переживает неудачу разбора: скан без текстового слоя, битый docx
        # и файл незнакомого генератора всё равно несут дату, которую записал
        # редактор. На корпусе владельца 35 файлов не читаются вовсе — их место в
        # хронологии от этого не исчезает.
        if "document_date" not in metadata:
            own_date = (
                _pdf_document_date_from_bytes(content)
                if ext == ".pdf" or detected_mime == "application/pdf"
                else (_office_document_date(content) if ext in _OFFICE_DATE_EXTENSIONS else None)
            )
            if own_date:
                metadata["document_date"] = own_date
        return DocumentResult(text, metadata, result.success, result.error)

    def extract_visual_assets(
        self,
        content: bytes,
        filename: str,
        mime_type: str = "",
        *,
        max_images: int = 4,
        max_pixels: int = 8_000_000,
        max_encoded_bytes: int = 3 * 1024 * 1024,
    ) -> list[VisualAsset]:
        """Extract a few representative images for local vision/OCR.

        The method never renders arbitrary PDF programs or invokes external
        binaries. Direct images are normalized with Pillow; PDF and Office
        documents contribute only already embedded raster images. Every asset
        is downscaled, re-encoded, and bounded before it leaves this module.
        """
        if not isinstance(content, bytes) or len(content) > self.max_input_bytes:
            return []
        max_images = max(1, min(int(max_images), 4))
        safe_name = Path(str(filename or "document")).name
        ext = self._compound_extension(safe_name.casefold())
        detected_mime = (mime_type or mimetypes.guess_type(safe_name)[0] or "").split(";", 1)[0]
        candidates: list[tuple[bytes, str]] = []

        if ext in _IMAGE_EXTENSIONS or detected_mime.startswith("image/"):
            candidates.append((content, safe_name))
        elif ext == ".pdf" or detected_mime == "application/pdf":
            candidates.extend(self._pdf_embedded_images(content, max_candidates=max_images * 3))
        elif ext in {".docx", ".pptx", ".xlsx", ".odt"}:
            candidates.extend(self._office_embedded_images(content, max_candidates=max_images * 3))

        assets: list[VisualAsset] = []
        seen: set[str] = set()
        for raw, source in candidates:
            if len(raw) > 16 * 1024 * 1024:
                continue
            normalized = self._normalize_visual_asset(
                raw,
                source=source,
                max_pixels=max_pixels,
                max_encoded_bytes=max_encoded_bytes,
            )
            if normalized is None:
                continue
            fingerprint = f"{len(normalized.data)}:{normalized.data[:64]!r}"
            if fingerprint in seen:
                continue
            seen.add(fingerprint)
            assets.append(normalized)
            if len(assets) >= max_images:
                break
        return assets

    @staticmethod
    def _normalize_visual_asset(
        content: bytes,
        *,
        source: str,
        max_pixels: int,
        max_encoded_bytes: int,
    ) -> VisualAsset | None:
        try:
            from PIL import Image, ImageOps
        except ImportError:
            return None
        try:
            with Image.open(io.BytesIO(content)) as opened_image:
                width, height = opened_image.size
                if width <= 0 or height <= 0 or width * height > 40_000_000:
                    return None
                opened_image.load()
                image = ImageOps.exif_transpose(opened_image)
                if image.width * image.height > max_pixels:
                    scale = (max_pixels / float(image.width * image.height)) ** 0.5
                    target = (max(1, int(image.width * scale)), max(1, int(image.height * scale)))
                    image.thumbnail(target)
                if image.mode not in {"RGB", "L"}:
                    background = Image.new("RGB", image.size, "white")
                    if "A" in image.getbands():
                        background.paste(image, mask=image.getchannel("A"))
                    else:
                        background.paste(image.convert("RGB"))
                    image = background
                elif image.mode == "L":
                    image = image.convert("RGB")

                for quality in (86, 76, 66):
                    output = io.BytesIO()
                    image.save(output, format="JPEG", quality=quality, optimize=False, progressive=False)
                    encoded = output.getvalue()
                    if len(encoded) <= max_encoded_bytes:
                        return VisualAsset(
                            data=encoded,
                            mime_type="image/jpeg",
                            source=source[:200],
                            width=image.width,
                            height=image.height,
                        )
                return None
        except Exception:
            LOGGER.debug("Visual asset normalization failed for %s", source, exc_info=True)
            return None

    @staticmethod
    def _pdf_embedded_images(content: bytes, *, max_candidates: int) -> list[tuple[bytes, str]]:
        try:
            from pypdf import PdfReader
        except ImportError:
            return []
        output: list[tuple[bytes, str]] = []
        try:
            reader = PdfReader(io.BytesIO(content), strict=False)
            if reader.is_encrypted:
                try:
                    if reader.decrypt("") == 0:
                        return []
                except Exception:
                    return []
            page_count = min(len(reader.pages), 250)
            if page_count <= 0:
                return []
            # Sample across the document instead of over-representing page one.
            indices = sorted(
                {
                    min(page_count - 1, round(position * (page_count - 1) / max(1, max_candidates - 1)))
                    for position in range(max_candidates)
                }
            )
            for page_index in indices:
                try:
                    images = reader.pages[page_index].images
                except Exception:
                    continue
                for image_index, image in enumerate(images):
                    data = bytes(image.data)
                    if 0 < len(data) <= 16 * 1024 * 1024:
                        output.append((data, f"pdf-page-{page_index + 1}-image-{image_index + 1}"))
                    if len(output) >= max_candidates:
                        return output
        except Exception:
            LOGGER.debug("PDF embedded-image extraction failed", exc_info=True)
        return output

    @staticmethod
    def _office_embedded_images(content: bytes, *, max_candidates: int) -> list[tuple[bytes, str]]:
        output: list[tuple[bytes, str]] = []
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as archive:
                members = [
                    info
                    for info in archive.infolist()
                    if not info.is_dir()
                    and info.filename.casefold().startswith(
                        ("word/media/", "ppt/media/", "xl/media/", "pictures/")
                    )
                    and Path(info.filename).suffix.casefold() in _IMAGE_EXTENSIONS
                ]
                members.sort(key=lambda info: info.file_size, reverse=True)
                for info in members[: max_candidates * 2]:
                    if info.file_size <= 0 or info.file_size > 16 * 1024 * 1024:
                        continue
                    if info.compress_size and info.file_size / max(1, info.compress_size) > _MAX_ZIP_RATIO:
                        continue
                    try:
                        data = archive.read(info, pwd=None)
                    except Exception:
                        # A single corrupted media stream (bit flip, interrupted save,
                        # adversarial upload) must not abort the whole document's
                        # extraction — zipfile.read() can raise things well outside
                        # the outer (OSError, BadZipFile, RuntimeError) tuple, e.g.
                        # zlib.error on a broken deflate stream. Skip this image and
                        # keep the others; found by adversarial review.
                        LOGGER.debug("Office embedded image %r is corrupted", info.filename, exc_info=True)
                        continue
                    if len(data) != info.file_size:
                        continue
                    output.append((data, info.filename))
                    if len(output) >= max_candidates:
                        break
        except (OSError, zipfile.BadZipFile, RuntimeError):
            LOGGER.debug("Office embedded-image extraction failed", exc_info=True)
        return output

    @staticmethod
    def _compound_extension(filename: str) -> str:
        for suffix in (".tar.gz", ".tar.bz2", ".tar.xz", ".tar.zst"):
            if filename.endswith(suffix):
                return suffix
        return Path(filename).suffix.casefold()

    @staticmethod
    def _decode(content: bytes) -> str:
        if content.startswith((b"\xff\xfe", b"\xfe\xff")):
            return content.decode("utf-16", errors="replace")
        if content.startswith(b"\xef\xbb\xbf"):
            return content.decode("utf-8-sig", errors="replace")
        for encoding in ("utf-8", "cp1251", "latin-1"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="replace")

    def _bounded_text_source(self, content: bytes, *, structured: bool = False) -> tuple[str, bool]:
        hard_limit = _MAX_STRUCTURED_PARSE_BYTES if structured else _MAX_TEXT_PARSE_BYTES
        # Enough bytes for the configured character budget without letting a
        # 50 MiB log/HTML/JSON file amplify through Python parser structures.
        derived_limit = max(64 * 1024, self.max_text_chars * (2 if structured else 4) + 4)
        parse_limit = min(self.max_input_bytes, hard_limit, derived_limit)
        prefix = content[:parse_limit]
        return self._decode(prefix).replace("\x00", ""), len(prefix) < len(content)

    def _append_bounded(
        self,
        parts: list[str],
        value: str,
        used: int,
        *,
        separator: str = "\n",
    ) -> tuple[int, bool]:
        """Append text without letting the assembled result exceed its budget."""

        if not value:
            return used, False
        separator_size = len(separator) if parts else 0
        available = self.max_text_chars - used - separator_size
        if available <= 0:
            return used, True
        clipped = value[:available]
        parts.append(clipped)
        return used + separator_size + len(clipped), len(clipped) != len(value)

    def _extract_text(self, content: bytes, ext: str) -> DocumentResult:
        structured = ext in {".json", ".csv", ".tsv", ".xml"}
        text, source_truncated = self._bounded_text_source(content, structured=structured)
        metadata: dict[str, Any] = {"format": ext.lstrip(".") or "text"}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        if ext == ".json":
            if source_truncated or len(content) > _MAX_STRUCTURED_PARSE_BYTES:
                metadata["json_valid"] = None
                metadata["json_pretty_skipped"] = True
            else:
                try:
                    text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
                    metadata["json_valid"] = True
                except json.JSONDecodeError:
                    metadata["json_valid"] = False
        elif ext in {".csv", ".tsv"}:
            delimiter = "\t" if ext == ".tsv" else ","
            try:
                rendered_rows: list[str] = []
                rendered_chars = 0
                rows_read = 0
                for row in csv.reader(io.StringIO(text), delimiter=delimiter):
                    if rows_read >= _MAX_TABULAR_ROWS:
                        metadata["rows_truncated"] = True
                        break
                    rows_read += 1
                    rendered = " | ".join(cell.strip() for cell in row)
                    remaining = self.max_text_chars - rendered_chars
                    if remaining <= 0:
                        metadata["rows_truncated"] = True
                        break
                    rendered_rows.append(rendered[:remaining])
                    rendered_chars += min(len(rendered), remaining) + 1
                    if len(rendered) > remaining or rendered_chars >= self.max_text_chars:
                        metadata["rows_truncated"] = True
                        break
                text = "\n".join(rendered_rows)
                metadata["rows_read"] = rows_read
            except csv.Error:
                pass
        elif ext == ".xml":
            text = self._strip_xml_tags(text)
        return DocumentResult(text, metadata)

    def _extract_html(self, content: bytes) -> DocumentResult:
        source, source_truncated = self._bounded_text_source(content, structured=True)
        soup = BeautifulSoup(source, "lxml")
        title = soup.title.get_text(" ", strip=True) if soup.title else ""
        for tag in soup(["script", "style", "svg", "noscript", "nav", "footer", "header"]):
            tag.decompose()
        root = soup.find("main") or soup.find("article") or soup.body or soup
        lines = [" ".join(line.split()) for line in root.get_text("\n").splitlines()]
        metadata: dict[str, Any] = {"format": "html", "title": html.unescape(title)}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        return DocumentResult(
            "\n".join(line for line in lines if line),
            metadata,
        )

    def _validate_zip(
        self,
        archive: zipfile.ZipFile,
        *,
        total_limit: int | None = None,
        member_limit: int | None = None,
    ) -> list[zipfile.ZipInfo]:
        members = archive.infolist()
        if len(members) > self.max_archive_entries:
            raise ArchiveLimitError(
                f"Archive entry count is {len(members)}; limit is {self.max_archive_entries}"
            )
        expanded_limit = self.max_archive_uncompressed_bytes if total_limit is None else total_limit
        total = 0
        for member in members:
            if member.flag_bits & 0x1:
                raise ArchiveLimitError("Encrypted ZIP entries are not supported")
            if member.is_dir():
                continue
            member_size = max(0, int(member.file_size))
            if member_limit is not None and member_size > member_limit:
                raise ArchiveLimitError("Archive member size exceeds configured limit")
            total += member_size
            if total > expanded_limit:
                raise ArchiveLimitError("Archive uncompressed size exceeds configured limit")
            compressed = max(1, int(member.compress_size))
            if member_size > 1024 * 1024 and member_size / compressed > _MAX_ZIP_RATIO:
                raise ArchiveLimitError("Suspicious ZIP compression ratio")
        return members

    def _validate_office_zip(self, archive: zipfile.ZipFile) -> list[zipfile.ZipInfo]:
        return self._validate_zip(
            archive,
            total_limit=min(self.max_archive_uncompressed_bytes, _MAX_OFFICE_EXPANDED_BYTES),
            member_limit=min(self.max_archive_uncompressed_bytes, _MAX_OFFICE_MEMBER_BYTES),
        )

    @staticmethod
    def _read_stream_limited(stream: _BinaryReadable, limit: int) -> bytes:
        chunks: list[bytes] = []
        total = 0
        while True:
            chunk = stream.read(min(64 * 1024, limit + 1 - total))
            if not chunk:
                break
            chunks.append(chunk)
            total += len(chunk)
            if total > limit:
                raise ArchiveLimitError("Decompressed member exceeds configured limit")
        return b"".join(chunks)

    @staticmethod
    def _read_stream_preview(stream: _BinaryReadable, limit: int) -> tuple[bytes, bool]:
        chunks: list[bytes] = []
        total = 0
        while total < limit:
            chunk = stream.read(min(64 * 1024, limit - total))
            if not chunk:
                return b"".join(chunks), False
            chunks.append(chunk)
            total += len(chunk)
        return b"".join(chunks), bool(stream.read(1))

    def _extract_docx(self, content: bytes) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            self._validate_office_zip(archive)
        try:
            from docx import Document

            document = Document(io.BytesIO(content))
            parts: list[str] = []
            used = 0
            extraction_truncated = False
            for paragraph in document.paragraphs:
                used, clipped = self._append_bounded(parts, paragraph.text.strip(), used)
                if clipped:
                    extraction_truncated = True
                    break
            for table in document.tables:
                if extraction_truncated:
                    break
                for row in table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        used, clipped = self._append_bounded(parts, " | ".join(values), used)
                        if clipped:
                            extraction_truncated = True
                            break
            metadata: dict[str, Any] = {
                "format": "docx",
                "paragraphs": len(document.paragraphs),
                "tables": len(document.tables),
            }
            if extraction_truncated:
                metadata["extraction_truncated"] = True
            return DocumentResult(
                "\n".join(parts),
                metadata,
            )
        except ImportError:
            return self._extract_xml_zip_text(content, "word/document.xml", "docx")

    def _extract_doc(self, content: bytes) -> DocumentResult:
        """Legacy Word 97-2003. Parsed from bytes; see `jericho.documents._ole`.

        Measured on a real 3.19 GB working folder: 206 files in this format, 130 MB,
        and every one of them arrived as `[File: NAME.doc; …]` — an Inbox item with
        nothing in it to review. 197 of the 206 now read, all recognisably Russian
        text, no replacement characters and no control characters left behind.

        A file that is not a compound document at all (five of those 206) still
        reports unsupported: the reviewer needs "Jericho cannot read this format"
        and "the file is damaged" to stay different sentences.
        """
        from jericho.documents._ole import OleError, extract_doc_text

        try:
            text, metadata = extract_doc_text(content)
        except OleError as exc:
            return DocumentResult("", {"format": "doc"}, False, f"Unsupported document format: {exc}")
        return DocumentResult(text, metadata)

    def _extract_xlsx(self, content: bytes) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            self._validate_office_zip(archive)
        try:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
        except ImportError:
            return DocumentResult("", {"format": "xlsx"}, False, "openpyxl is not installed")
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts: list[str] = []
        used = 0
        row_count = 0
        extraction_truncated = False
        sheet_count = len(workbook.sheetnames)
        try:
            for sheet in workbook.worksheets:
                used, clipped = self._append_bounded(parts, f"--- Sheet: {sheet.title} ---", used)
                if clipped:
                    extraction_truncated = True
                    break
                for row in sheet.iter_rows(values_only=True):
                    if row_count >= _MAX_TABULAR_ROWS:
                        extraction_truncated = True
                        break
                    row_count += 1
                    values = [str(value) if value is not None else "" for value in row]
                    if any(values):
                        used, clipped = self._append_bounded(parts, " | ".join(values), used)
                        if clipped:
                            extraction_truncated = True
                            break
                if extraction_truncated:
                    break
        finally:
            workbook.close()
        metadata: dict[str, Any] = {
            "format": "xlsx",
            "sheets": sheet_count,
            "rows_read": row_count,
        }
        if extraction_truncated:
            metadata["extraction_truncated"] = True
        return DocumentResult(
            "\n".join(parts),
            metadata,
        )

    def _extract_pptx(self, content: bytes) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            self._validate_office_zip(archive)
        try:
            from pptx import Presentation
        except ImportError:
            return DocumentResult("", {"format": "pptx"}, False, "python-pptx is not installed")
        presentation = Presentation(io.BytesIO(content))
        slides: list[str] = []
        used = 0
        extraction_truncated = False
        for number, slide in enumerate(presentation.slides, 1):
            text = [
                shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
            ]
            if text:
                used, clipped = self._append_bounded(
                    slides,
                    f"--- Slide {number} ---\n" + "\n".join(text),
                    used,
                    separator="\n\n",
                )
                if clipped:
                    extraction_truncated = True
                    break
        metadata: dict[str, Any] = {"format": "pptx", "slides": len(presentation.slides)}
        if extraction_truncated:
            metadata["extraction_truncated"] = True
        return DocumentResult("\n\n".join(slides), metadata)

    def _extract_rtf(self, content: bytes) -> DocumentResult:
        source, source_truncated = self._bounded_text_source(content)
        metadata: dict[str, Any] = {"format": "rtf"}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        try:
            from striprtf.striprtf import rtf_to_text  # type: ignore[import-not-found]
        except ImportError:
            # A conservative fallback removes control words but keeps visible text.
            source = re.sub(r"\\'[0-9a-fA-F]{2}", " ", source)
            source = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", source)
            source = source.replace("{", " ").replace("}", " ")
            metadata["parser"] = "fallback"
            return DocumentResult(" ".join(source.split()), metadata)
        metadata["parser"] = "striprtf"
        return DocumentResult(rtf_to_text(source), metadata)

    def _extract_xml_zip_text(self, content: bytes, member_name: str, format_name: str) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = self._validate_office_zip(archive)
            names = {member.filename for member in members}
            if member_name not in names:
                return DocumentResult("", {"format": format_name}, False, f"Missing member: {member_name}")
            info = archive.getinfo(member_name)
            member_limit = min(self.max_archive_uncompressed_bytes, _MAX_OFFICE_MEMBER_BYTES)
            if info.file_size > member_limit:
                raise ArchiveLimitError("Office document XML exceeds configured limit")
            parse_limit = min(member_limit, _MAX_STRUCTURED_PARSE_BYTES)
            with archive.open(info) as stream:
                data, source_truncated = self._read_stream_preview(stream, parse_limit)
        metadata: dict[str, Any] = {"format": format_name}
        if source_truncated:
            metadata["source_truncated_for_parse"] = True
        return DocumentResult(self._strip_xml_tags(self._decode(data)), metadata)

    def _extract_pdf(self, content: bytes) -> DocumentResult:
        try:
            from pypdf import PdfReader
        except ImportError:
            return DocumentResult("", {"format": "pdf"}, False, "pypdf is not installed")
        reader = PdfReader(io.BytesIO(content), strict=False)
        if reader.is_encrypted:
            try:
                if reader.decrypt("") == 0:
                    return DocumentResult("", {"format": "pdf"}, False, "Encrypted PDF is not supported")
            except Exception:
                return DocumentResult("", {"format": "pdf"}, False, "Encrypted PDF is not supported")
        pages: list[str] = []
        pages_read = 0
        text_chars = 0
        extraction_truncated = False
        for page in itertools.islice(reader.pages, 250):
            pages_read += 1
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_chars, clipped = self._append_bounded(
                    pages,
                    page_text.strip(),
                    text_chars,
                    separator="\n\n",
                )
                if clipped:
                    extraction_truncated = True
                    break
            if text_chars >= self.max_text_chars:
                extraction_truncated = True
                break
        metadata: dict[str, Any] = {
            "format": "pdf",
            "pages_read": pages_read,
            "page_limit": 250,
        }
        if extraction_truncated:
            metadata["extraction_truncated"] = True
        return DocumentResult(
            "\n\n".join(pages),
            metadata,
        )

    def _extract_archive(
        self, content: bytes, filename: str, ext: str, depth: int, budget: _ArchiveBudget
    ) -> DocumentResult:
        if ext == ".zip":
            return self._extract_zip(content, depth, budget)
        if ext in {".tar", ".tar.gz", ".tar.bz2", ".tar.xz"}:
            return self._extract_tar(content, ext, depth, budget)
        if ext in {".gz", ".bz2", ".xz", ".zst"}:
            decompressed = self._decompress_single(content, ext, budget)
            inner_name = filename[: -len(ext)] or "decompressed.txt"
            return self.extract(decompressed, inner_name, _depth=depth + 1, _budget=budget)
        if ext == ".rar":
            return self._extract_rar(content, depth, budget)
        if ext == ".7z":
            return self._extract_7z(content)
        if ext == ".tar.zst":
            decompressed = self._decompress_single(content, ".zst", budget)
            return self._extract_tar(decompressed, ".tar", depth, budget)
        return DocumentResult("", {"format": ext.lstrip(".")}, False, "Unsupported archive format")

    def _decompress_single(self, content: bytes, ext: str, budget: _ArchiveBudget) -> bytes:
        limit = min(self.max_archive_uncompressed_bytes, self.max_input_bytes, budget.expanded_bytes)
        if ext == ".gz":
            stream: _BinaryReadable = gzip.GzipFile(fileobj=io.BytesIO(content))
        elif ext == ".bz2":
            stream = bz2.BZ2File(io.BytesIO(content))
        elif ext == ".xz":
            stream = lzma.LZMAFile(io.BytesIO(content))  # noqa: SIM115 - closed via closing() below
        elif ext == ".zst":
            try:
                import zstandard  # type: ignore[import-not-found]
            except ImportError as exc:
                raise ValueError("zstandard is not installed") from exc
            stream = zstandard.ZstdDecompressor().stream_reader(io.BytesIO(content))
        else:
            raise ValueError(f"Unsupported compressor: {ext}")
        with closing(stream):
            data = self._read_stream_limited(stream, limit)
        budget.spend_bytes(len(data))
        return data

    def _member_preview(self, name: str, data: bytes, depth: int, budget: _ArchiveBudget) -> str:
        result = self.extract(data, name, _depth=depth + 1, _budget=budget)
        if not result.success or not result.text:
            return ""
        return f"\n--- {name} ---\n{result.text[:20_000]}"

    def _extract_zip(self, content: bytes, depth: int, budget: _ArchiveBudget) -> DocumentResult:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = self._validate_zip(archive)
            files = [member for member in members if not member.is_dir()]
            parts = [f"ZIP archive: {len(files)} files", *(member.filename for member in files[:100])]
            previewed = 0
            for member in files:
                if member.file_size > _MAX_MEMBER_PREVIEW_BYTES:
                    continue
                # Count the decompression, not the success. Incrementing only when a
                # preview came back meant a member that yields no text — an inner
                # archive of binary blobs, say — never advanced the cap, so the loop
                # decompressed EVERY member. Nested, that is a decompression bomb the
                # cap was supposed to bound: a 99 KB upload (24 x 24 x 500 zero-filled
                # members) held the event loop for 31 seconds.
                #
                # The allowance belongs to the UPLOAD, not to this archive: a nested
                # member used to start over with a full one.
                if not budget.take_preview():
                    break
                previewed += 1
                with archive.open(member) as stream:
                    data = self._read_stream_limited(stream, _MAX_MEMBER_PREVIEW_BYTES)
                budget.spend_bytes(len(data))
                preview = self._member_preview(member.filename, data, depth, budget)
                if preview:
                    parts.append(preview)
        return DocumentResult(
            "\n".join(parts),
            {"format": "zip", "files": len(files), "previewed_files": previewed},
        )

    def _extract_tar(self, content: bytes, ext: str, depth: int, budget: _ArchiveBudget) -> DocumentResult:
        # Streaming mode avoids materialising an attacker-controlled member
        # list before the configured entry limit can be enforced.
        with tarfile.open(fileobj=io.BytesIO(content), mode="r|*") as archive:
            entry_count = 0
            file_count = 0
            total = 0
            names: list[str] = []
            previews: list[str] = []
            previewed = 0
            exhausted = False
            for member in archive:
                entry_count += 1
                if entry_count > self.max_archive_entries:
                    raise ArchiveLimitError(f"Archive entry count exceeds limit {self.max_archive_entries}")
                if member.issym() or member.islnk() or member.isdev():
                    raise ArchiveLimitError("Archive links and device entries are not supported")
                if not member.isfile():
                    continue
                file_count += 1
                member_size = max(0, int(member.size))
                total += member_size
                if total > self.max_archive_uncompressed_bytes:
                    raise ArchiveLimitError("Archive uncompressed size exceeds configured limit")
                budget.spend_bytes(member_size)
                if len(names) < 100:
                    names.append(member.name)
                # Walking the members is itself the cost here: streaming mode has to
                # read past every member it skips. When the upload's allowance is
                # gone, stop walking rather than skip in a loop — that is where the
                # nested bomb's 107 seconds were actually spent.
                if budget.expanded_bytes <= 0:
                    exhausted = True
                    break
                if member_size > _MAX_MEMBER_PREVIEW_BYTES:
                    continue
                stream = archive.extractfile(member)
                if stream is None:
                    continue
                if not budget.take_preview():
                    exhausted = True
                    break
                previewed += 1  # decompressions, not successes — see _extract_zip
                with stream:
                    data = self._read_stream_limited(stream, _MAX_MEMBER_PREVIEW_BYTES)
                preview = self._member_preview(member.name, data, depth, budget)
                if preview:
                    previews.append(preview)
            parts = [f"TAR archive: {file_count} files", *names, *previews]
        metadata: dict[str, Any] = {
            "format": ext.lstrip("."),
            "files": file_count,
            "previewed_files": previewed,
        }
        if exhausted:
            # Said out loud: the listing is partial, and a caller reading `files`
            # as "what the archive holds" would otherwise be quietly wrong.
            metadata["archive_budget_exhausted"] = True
        return DocumentResult("\n".join(parts), metadata)

    def _extract_rar(self, content: bytes, depth: int, budget: _ArchiveBudget) -> DocumentResult:
        try:
            import rarfile  # type: ignore[import-untyped]
        except ImportError:
            return DocumentResult("", {"format": "rar"}, False, "rarfile is not installed")
        with rarfile.RarFile(io.BytesIO(content)) as archive:
            if archive.needs_password():
                raise ArchiveLimitError("Encrypted RAR archives are not supported")
            members = archive.infolist()
            if len(members) > self.max_archive_entries:
                raise ArchiveLimitError("RAR entry count exceeds configured limit")
            files = [member for member in members if not member.isdir()]
            total = sum(max(0, int(member.file_size)) for member in files)
            if total > self.max_archive_uncompressed_bytes:
                raise ArchiveLimitError("RAR uncompressed size exceeds configured limit")
            parts = [f"RAR archive: {len(files)} files", *(member.filename for member in files[:100])]
            previewed = 0
            for member in files:
                if member.file_size > _MAX_MEMBER_PREVIEW_BYTES:
                    continue
                if not budget.take_preview():
                    break
                previewed += 1  # decompressions, not successes — see _extract_zip
                with archive.open(member) as stream:
                    data = self._read_stream_limited(stream, _MAX_MEMBER_PREVIEW_BYTES)
                budget.spend_bytes(len(data))
                preview = self._member_preview(member.filename, data, depth, budget)
                if preview:
                    parts.append(preview)
        return DocumentResult(
            "\n".join(parts), {"format": "rar", "files": len(files), "previewed_files": previewed}
        )

    def _extract_7z(self, content: bytes) -> DocumentResult:
        try:
            import py7zr  # type: ignore[import-not-found]
        except ImportError:
            return DocumentResult("", {"format": "7z"}, False, "py7zr is not installed")
        with py7zr.SevenZipFile(io.BytesIO(content), mode="r") as archive:
            entries = archive.list()
            if len(entries) > self.max_archive_entries:
                raise ArchiveLimitError("7z entry count exceeds configured limit")
            total = sum(max(0, int(getattr(entry, "uncompressed", 0) or 0)) for entry in entries)
            if total > self.max_archive_uncompressed_bytes:
                raise ArchiveLimitError("7z uncompressed size exceeds configured limit")
            names = [str(getattr(entry, "filename", "")) for entry in entries]
        # Listing is intentional: py7zr extraction APIs write to a target path in
        # several versions.  Jericho does not permit archive members on disk.
        return DocumentResult(
            "\n".join([f"7z archive: {len(names)} files", *names[:100]]),
            {"format": "7z", "files": len(names), "previewed_files": 0},
        )

    @staticmethod
    def _strip_xml_tags(xml_text: str) -> str:
        text = re.sub(r"<[^>]+>", " ", xml_text)
        text = html.unescape(text)
        return " ".join(text.split())
